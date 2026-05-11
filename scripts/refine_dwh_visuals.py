from __future__ import annotations

import math
import textwrap
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Rectangle
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches


ROOT = Path("/Users/orperetz/Documents/shenkar/הנדסת נתונים/קורס")
SLIDES_DIR = ROOT / "slides"
WEEK04_DIR = ROOT / "diagrams" / "week04"
WEEK05_DIR = ROOT / "diagrams" / "week05"


PALETTE = {
    "source": {"fill": "#B3E5FC", "edge": "#81D4FA", "text": "#01579B"},
    "process": {"fill": "#C8E6C9", "edge": "#A5D6A7", "text": "#1B5E20"},
    "decision": {"fill": "#FFE0B2", "edge": "#FFCC80", "text": "#E65100"},
    "failure": {"fill": "#FFCDD2", "edge": "#EF9A9A", "text": "#B71C1C"},
    "success": {"fill": "#C5E1A5", "edge": "#AED581", "text": "#33691E"},
    "store": {"fill": "#BBDEFB", "edge": "#90CAF9", "text": "#0D47A1"},
    "neutral": {"fill": "#FFFFFF", "edge": "#B0BEC5", "text": "#37474F"},
    "muted": {"fill": "#ECEFF1", "edge": "#CFD8DC", "text": "#455A64"},
}

ARROW_COLOR = "#5C6B73"
PANEL_EDGE = "#CFD8DC"
PANEL_FILL = "#FAFAFA"
FONT = "DejaVu Sans"


def setup_figure(width=14, height=7):
    fig, ax = plt.subplots(figsize=(width, height), dpi=220)
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    return fig, ax


def save(fig, path: Path):
    path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.03)
    plt.close(fig)


def panel(ax, x, y, w, h, title=None):
    patch = FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle="round,pad=0.012,rounding_size=0.02",
        linewidth=1.5,
        edgecolor=PANEL_EDGE,
        facecolor=PANEL_FILL,
    )
    ax.add_patch(patch)
    if title:
        ax.text(
            x + w / 2,
            y + h - 0.035,
            title,
            ha="center",
            va="center",
            fontsize=16,
            fontweight="bold",
            color="#37474F",
            family=FONT,
        )


def box(ax, x, y, w, h, text, kind="neutral", tag=None, fontsize=15, align="center"):
    style = PALETTE[kind]
    patch = FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle="round,pad=0.01,rounding_size=0.02",
        linewidth=1.6,
        edgecolor=style["edge"],
        facecolor=style["fill"],
    )
    ax.add_patch(patch)
    ax.text(
        x + (w / 2 if align == "center" else 0.02),
        y + h / 2,
        text,
        ha=align,
        va="center",
        fontsize=fontsize,
        fontweight="semibold",
        color=style["text"],
        family=FONT,
        wrap=True,
    )


def arrow(ax, start, end, label=None, label_pos=0.5, dy=0.018, fontsize=12):
    patch = FancyArrowPatch(
        start,
        end,
        arrowstyle="-|>",
        mutation_scale=16,
        linewidth=2,
        color=ARROW_COLOR,
    )
    ax.add_patch(patch)
    if label:
        x = start[0] + (end[0] - start[0]) * label_pos
        y = start[1] + (end[1] - start[1]) * label_pos + dy
        ax.text(
            x,
            y,
            label,
            ha="center",
            va="center",
            fontsize=fontsize,
            color="#455A64",
            family=FONT,
        )


def dashed_arrow(ax, start, end, label=None, label_pos=0.5, dy=0.02, fontsize=11):
    patch = FancyArrowPatch(
        start,
        end,
        arrowstyle="-|>",
        mutation_scale=14,
        linewidth=1.6,
        color=ARROW_COLOR,
        linestyle="--",
    )
    ax.add_patch(patch)
    if label:
        x = start[0] + (end[0] - start[0]) * label_pos
        y = start[1] + (end[1] - start[1]) * label_pos + dy
        ax.text(
            x,
            y,
            label,
            ha="center",
            va="center",
            fontsize=fontsize,
            color="#546E7A",
            family=FONT,
        )


def annotate(ax, x, y, text, fontsize=13, color="#455A64", ha="center"):
    ax.text(x, y, text, ha=ha, va="center", fontsize=fontsize, color=color, family=FONT)


def wrapped(text: str, width: int) -> str:
    return "\n".join(textwrap.wrap(text, width=width))


def make_week04_hybrid():
    fig, ax = setup_figure(14, 7.4)
    panel(ax, 0.02, 0.08, 0.96, 0.84, "ShopNow Hybrid Analytics Architecture")

    sources = [
        ("Orders\nPostgres", 0.06, 0.67),
        ("Payments\nStripe API", 0.06, 0.51),
        ("Customers\nCRM", 0.06, 0.35),
        ("Clickstream\nJSON files", 0.06, 0.19),
    ]
    for label, x, y in sources:
        box(ax, x, y, 0.13, 0.11, label, "source", tag="source", fontsize=13)

    box(ax, 0.27, 0.30, 0.17, 0.28, "Bronze\nraw landing", "store", tag="storage", fontsize=17)
    box(ax, 0.50, 0.39, 0.16, 0.11, "Standardize keys\nDedup + DQ checks", "process", tag="process", fontsize=14)
    box(ax, 0.73, 0.30, 0.16, 0.28, "Silver\ntrusted entities", "store", tag="storage", fontsize=17)
    box(ax, 0.73, 0.67, 0.16, 0.11, "Gold facts + dims\nCertified metrics", "success", tag="success", fontsize=14)
    box(ax, 0.91, 0.45, 0.06, 0.18, "BI\nDashboards", "neutral", tag="use", fontsize=13)

    for _, _, y in sources:
        arrow(ax, (0.19, y + 0.055), (0.27, 0.44), label=None)
    arrow(ax, (0.44, 0.44), (0.50, 0.445))
    arrow(ax, (0.66, 0.445), (0.73, 0.44))
    arrow(ax, (0.81, 0.58), (0.81, 0.67), label="model marts", label_pos=0.55, dy=0.03)
    arrow(ax, (0.89, 0.72), (0.91, 0.54), label="query", label_pos=0.4, dy=0.03)

    annotate(ax, 0.355, 0.24, "Raw, replayable, schema-on-read", fontsize=12)
    annotate(ax, 0.81, 0.24, "Cleaned joins, conformed keys", fontsize=12)
    annotate(ax, 0.81, 0.82, "Finance / retention / LTV", fontsize=12)

    save(fig, WEEK04_DIR / "week4_hybrid_architecture_rebuild.png")


def make_week04_extraction():
    fig, ax = setup_figure(14, 7.2)
    panel(ax, 0.02, 0.08, 0.96, 0.84, "ShopNow Extraction Contracts")

    xs = [0.04, 0.28, 0.52, 0.76]
    source_labels = ["Orders DB", "Stripe events", "CRM master", "Clickstream files"]
    extract_labels = [
        "updated_at > watermark",
        "cursor > last_event_id",
        "daily snapshot + diff",
        "new files by date prefix",
    ]
    for x, s_label, e_label in zip(xs, source_labels, extract_labels):
        box(ax, x, 0.66, 0.18, 0.12, s_label, "source", tag="source", fontsize=14)
        box(ax, x, 0.42, 0.18, 0.12, wrapped(e_label, 18), "process", tag="extract", fontsize=13)
        arrow(ax, (x + 0.09, 0.66), (x + 0.09, 0.54))
    box(ax, 0.30, 0.15, 0.40, 0.14, "Staging / landing zone\nraw payload + ingestion metadata", "store", tag="storage", fontsize=16)
    for x in xs:
        arrow(ax, (x + 0.09, 0.42), (0.50, 0.29))
    annotate(ax, 0.50, 0.10, "Goal: pull only new data while keeping source systems responsive", fontsize=13)
    save(fig, WEEK04_DIR / "week4_extraction_contracts_rebuild.png")


def make_week04_safe_loading():
    fig, ax = setup_figure(14, 7.2)
    panel(ax, 0.02, 0.08, 0.96, 0.84, "Safe Loading for ShopNow")

    box(ax, 0.07, 0.37, 0.18, 0.16, "stg_orders\n+ stg_customers", "store", tag="staging", fontsize=15)
    box(ax, 0.32, 0.37, 0.17, 0.16, "Dedup by\nbusiness key", "process", tag="process", fontsize=15)
    box(ax, 0.56, 0.55, 0.16, 0.14, "MERGE\nfact_order_line", "success", tag="load", fontsize=15)
    box(ax, 0.56, 0.25, 0.16, 0.14, "UPSERT\ndim_customer", "success", tag="load", fontsize=15)
    box(ax, 0.80, 0.41, 0.15, 0.16, "Gold target\nsame revenue once", "store", tag="target", fontsize=15)
    box(ax, 0.34, 0.72, 0.18, 0.12, "Control table\nlast_successful_batch", "store", tag="state", fontsize=13)

    arrow(ax, (0.25, 0.45), (0.32, 0.45))
    arrow(ax, (0.49, 0.45), (0.56, 0.62))
    arrow(ax, (0.49, 0.45), (0.56, 0.32))
    arrow(ax, (0.72, 0.62), (0.80, 0.49))
    arrow(ax, (0.72, 0.32), (0.80, 0.49))
    dashed_arrow(ax, (0.87, 0.41), (0.43, 0.72), label="update checkpoint", label_pos=0.45, dy=0.02)

    annotate(ax, 0.17, 0.24, "Input can be replayed", fontsize=12)
    annotate(ax, 0.87, 0.24, "Rerun does not duplicate facts", fontsize=12)
    save(fig, WEEK04_DIR / "week4_safe_loading_rebuild.png")


def make_week04_etl_vs_elt():
    fig, ax = setup_figure(14, 7.2)
    panel(ax, 0.03, 0.10, 0.44, 0.80, "ETL")
    panel(ax, 0.53, 0.10, 0.44, 0.80, "ELT")

    # ETL
    box(ax, 0.08, 0.66, 0.12, 0.12, "Sources", "source", tag="source", fontsize=15)
    box(ax, 0.24, 0.66, 0.12, 0.12, "Extract", "process", tag="step", fontsize=15)
    box(ax, 0.24, 0.43, 0.12, 0.12, "Transform\nbefore load", "process", tag="step", fontsize=14)
    box(ax, 0.24, 0.20, 0.12, 0.12, "Load", "success", tag="step", fontsize=15)
    box(ax, 0.40, 0.43, 0.08, 0.20, "Warehouse", "store", tag="target", fontsize=15)
    arrow(ax, (0.20, 0.72), (0.24, 0.72))
    arrow(ax, (0.30, 0.66), (0.30, 0.55))
    arrow(ax, (0.30, 0.43), (0.30, 0.32))
    arrow(ax, (0.36, 0.26), (0.40, 0.43))
    annotate(ax, 0.25, 0.14, "Quality gates happen before analysts see the data", fontsize=12)

    # ELT
    box(ax, 0.58, 0.66, 0.12, 0.12, "Sources", "source", tag="source", fontsize=15)
    box(ax, 0.74, 0.66, 0.12, 0.12, "Extract", "process", tag="step", fontsize=15)
    box(ax, 0.74, 0.43, 0.12, 0.12, "Load raw", "success", tag="step", fontsize=15)
    box(ax, 0.90, 0.43, 0.06, 0.20, "Lakehouse /\nWarehouse", "store", tag="target", fontsize=13)
    box(ax, 0.74, 0.20, 0.12, 0.12, "Transform\ninside platform", "process", tag="step", fontsize=13)
    arrow(ax, (0.70, 0.72), (0.74, 0.72))
    arrow(ax, (0.80, 0.66), (0.80, 0.55))
    arrow(ax, (0.86, 0.49), (0.90, 0.49))
    arrow(ax, (0.90, 0.43), (0.80, 0.32))
    annotate(ax, 0.76, 0.14, "Cloud compute makes late transformation practical", fontsize=12)

    save(fig, WEEK04_DIR / "week4_etl_vs_elt_rebuild.png")


def make_week04_cdc_options():
    fig, ax = setup_figure(14, 7.0)
    panel(ax, 0.02, 0.08, 0.96, 0.84, "CDC Options Compared")
    cols = [
        ("Snapshot", "Small lookup tables", "slow", "manual diff", "No"),
        ("Timestamp", "Orders with updated_at", "light", "weak deletes", "Partial"),
        ("Log-based", "Payments / order status", "light", "best fidelity", "Yes"),
        ("Trigger", "Only if logs unavailable", "medium", "app impact", "Yes"),
    ]
    x_positions = [0.04, 0.28, 0.52, 0.76]
    for x, (title, good_for, load, note, deletes) in zip(x_positions, cols):
        box(ax, x, 0.66, 0.18, 0.12, title, "process", tag="mode", fontsize=15)
        box(ax, x, 0.46, 0.18, 0.10, wrapped(f"Best for\n{good_for}", 16), "source", tag="fit", fontsize=12)
        box(ax, x, 0.30, 0.18, 0.10, f"Source load:\n{load}", "decision", tag="cost", fontsize=12)
        box(ax, x, 0.14, 0.18, 0.10, f"Deletes:\n{deletes}\n{note}", "muted", tag="tradeoff", fontsize=11)
    save(fig, WEEK04_DIR / "week4_cdc_options_rebuild.png")


def make_week04_cdc_rerun():
    fig, ax = setup_figure(14, 7.1)
    panel(ax, 0.02, 0.08, 0.96, 0.84, "CDC + Watermark + Rerun Logic")

    box(ax, 0.06, 0.43, 0.16, 0.14, "orders_log\nLSN 1041..1080", "source", tag="source", fontsize=15)
    box(ax, 0.30, 0.43, 0.16, 0.14, "Read only\n> last_lsn", "process", tag="extract", fontsize=15)
    box(ax, 0.54, 0.43, 0.16, 0.14, "Dedup + MERGE\nfacts / dims", "success", tag="load", fontsize=15)
    box(ax, 0.78, 0.43, 0.16, 0.14, "Target state\ncommitted once", "store", tag="target", fontsize=15)
    box(ax, 0.54, 0.72, 0.18, 0.12, "Control table\nlast_committed_lsn=1080", "store", tag="state", fontsize=13)
    box(ax, 0.30, 0.17, 0.24, 0.11, "Rerun after failure:\nrestart from saved LSN", "failure", tag="rerun", fontsize=13)

    arrow(ax, (0.22, 0.50), (0.30, 0.50))
    arrow(ax, (0.46, 0.50), (0.54, 0.50))
    arrow(ax, (0.70, 0.50), (0.78, 0.50))
    dashed_arrow(ax, (0.63, 0.57), (0.63, 0.72), label="commit watermark", label_pos=0.55, dy=0.03)
    dashed_arrow(ax, (0.54, 0.17), (0.39, 0.43), label="skip old changes", label_pos=0.55, dy=0.02)

    save(fig, WEEK04_DIR / "week4_cdc_rerun_rebuild.png")


def make_week04_sttm():
    fig, ax = setup_figure(14, 7.0)
    panel(ax, 0.02, 0.08, 0.96, 0.84, "ShopNow STTM Mapping Flow")

    box(ax, 0.06, 0.63, 0.20, 0.12, "crm_customer_id\norder_status\ngross_amount", "source", tag="source fields", fontsize=14)
    box(ax, 0.06, 0.28, 0.20, 0.12, "currency\ntimezone\nrefund_amount", "source", tag="source fields", fontsize=14)
    box(ax, 0.39, 0.52, 0.22, 0.16, "Normalize codes\nMap natural -> surrogate keys\nConvert currency/time", "process", tag="rules", fontsize=13)
    box(ax, 0.72, 0.63, 0.20, 0.12, "customer_key\nstatus\nnet_revenue", "success", tag="target columns", fontsize=14)
    box(ax, 0.72, 0.28, 0.20, 0.12, "order_ts_utc\nrefund_flag\nunknown member", "success", tag="target columns", fontsize=14)

    arrow(ax, (0.26, 0.69), (0.39, 0.60))
    arrow(ax, (0.26, 0.34), (0.39, 0.58))
    arrow(ax, (0.61, 0.60), (0.72, 0.69))
    arrow(ax, (0.61, 0.56), (0.72, 0.34))

    annotate(ax, 0.50, 0.14, "The mapping spec is the contract between messy operational fields and stable analytical columns", fontsize=13)
    save(fig, WEEK04_DIR / "week4_sttm_rebuild.png")


def make_week04_idempotency():
    fig, ax = setup_figure(14, 6.8)
    panel(ax, 0.02, 0.08, 0.96, 0.84, "Idempotency: same input -> same final state")

    box(ax, 0.07, 0.43, 0.16, 0.16, "Batch D\n(order lines 1..N)", "source", tag="input", fontsize=15)
    box(ax, 0.33, 0.43, 0.18, 0.16, "Load f(D)\nMERGE / UPSERT", "process", tag="operation", fontsize=15)
    box(ax, 0.61, 0.43, 0.16, 0.16, "Run again\nsame batch D", "source", tag="rerun", fontsize=15)
    box(ax, 0.80, 0.43, 0.13, 0.16, "Target\nunchanged", "success", tag="result", fontsize=15)
    arrow(ax, (0.23, 0.51), (0.33, 0.51))
    arrow(ax, (0.51, 0.51), (0.80, 0.51))
    dashed_arrow(ax, (0.69, 0.43), (0.42, 0.43), label="same transformation rules", label_pos=0.5, dy=-0.05)
    annotate(ax, 0.50, 0.20, "Reprocessing cannot duplicate revenue if the write path is deterministic", fontsize=13)
    save(fig, WEEK04_DIR / "week4_idempotency_rebuild.png")


def make_week04_merge_vs_overwrite():
    fig, ax = setup_figure(14, 7.0)
    panel(ax, 0.03, 0.10, 0.44, 0.80, "MERGE")
    panel(ax, 0.53, 0.10, 0.44, 0.80, "Partition Overwrite")

    box(ax, 0.10, 0.60, 0.14, 0.12, "Mixed inserts\n+ updates", "source", tag="input", fontsize=14)
    box(ax, 0.29, 0.60, 0.14, 0.12, "Match on\nbusiness key", "process", tag="logic", fontsize=14)
    box(ax, 0.18, 0.30, 0.18, 0.14, "Best when\nchanges hit only some rows", "success", tag="use", fontsize=13)
    arrow(ax, (0.24, 0.66), (0.29, 0.66))
    arrow(ax, (0.36, 0.60), (0.27, 0.44))

    box(ax, 0.60, 0.60, 0.14, 0.12, "Whole day /\nmonth slice", "source", tag="input", fontsize=14)
    box(ax, 0.79, 0.60, 0.14, 0.12, "Replace one\npartition", "process", tag="logic", fontsize=14)
    box(ax, 0.68, 0.30, 0.18, 0.14, "Best when\nbackfilling controlled windows", "success", tag="use", fontsize=13)
    arrow(ax, (0.74, 0.66), (0.79, 0.66))
    arrow(ax, (0.86, 0.60), (0.77, 0.44))

    annotate(ax, 0.50, 0.05, "Choose MERGE for mixed row-level change streams; choose overwrite when you intentionally rebuild a full slice", fontsize=12)
    save(fig, WEEK04_DIR / "week4_merge_vs_overwrite_rebuild.png")


def make_week05_star_vs_snowflake():
    fig, ax = setup_figure(14, 7.2)
    panel(ax, 0.03, 0.10, 0.44, 0.80, "Star Schema")
    panel(ax, 0.53, 0.10, 0.44, 0.80, "Snowflake Schema")

    # star
    box(ax, 0.18, 0.44, 0.14, 0.12, "sales_fact", "store", tag="fact", fontsize=15)
    box(ax, 0.18, 0.70, 0.14, 0.10, "dim_date", "store", tag="dimension", fontsize=13)
    box(ax, 0.04, 0.44, 0.12, 0.10, "dim_customer", "store", tag="dimension", fontsize=12)
    box(ax, 0.34, 0.44, 0.12, 0.10, "dim_product", "store", tag="dimension", fontsize=12)
    arrow(ax, (0.25, 0.70), (0.25, 0.56))
    arrow(ax, (0.16, 0.49), (0.18, 0.49))
    arrow(ax, (0.34, 0.49), (0.32, 0.49))
    annotate(ax, 0.25, 0.22, "Few joins\nSimple BI SQL", fontsize=13)

    # snowflake
    box(ax, 0.68, 0.44, 0.14, 0.12, "sales_fact", "store", tag="fact", fontsize=15)
    box(ax, 0.68, 0.70, 0.14, 0.10, "dim_date", "store", tag="dimension", fontsize=13)
    box(ax, 0.55, 0.44, 0.11, 0.10, "dim_customer", "store", tag="dimension", fontsize=11)
    box(ax, 0.84, 0.44, 0.11, 0.10, "dim_product", "store", tag="dimension", fontsize=11)
    box(ax, 0.84, 0.62, 0.11, 0.09, "dim_brand", "store", tag="hierarchy", fontsize=11)
    box(ax, 0.84, 0.78, 0.11, 0.09, "dim_category", "store", tag="hierarchy", fontsize=11)
    arrow(ax, (0.75, 0.70), (0.75, 0.56))
    arrow(ax, (0.66, 0.49), (0.68, 0.49))
    arrow(ax, (0.84, 0.49), (0.82, 0.49))
    arrow(ax, (0.895, 0.54), (0.895, 0.62))
    arrow(ax, (0.895, 0.71), (0.895, 0.78))
    annotate(ax, 0.75, 0.22, "More joins\nLess hierarchy duplication", fontsize=13)

    save(fig, WEEK05_DIR / "week5_star_vs_snowflake_rebuild.png")


def make_week05_product_hierarchy():
    fig, ax = setup_figure(14, 7.0)
    panel(ax, 0.02, 0.08, 0.96, 0.84, "Snowflaked Product Hierarchy")

    levels = [
        ("category", "Electronics", 0.07),
        ("brand", "Apple", 0.30),
        ("line", "iPhone", 0.53),
        ("product", "iPhone 15 Pro", 0.76),
    ]
    for kind, label, x in levels:
        box(ax, x, 0.52, 0.16, 0.16, label, "store", tag=kind, fontsize=16)
    for idx in range(len(levels) - 1):
        arrow(ax, (levels[idx][2] + 0.16, 0.60), (levels[idx + 1][2], 0.60), label="1-to-many", fontsize=11)
    box(ax, 0.41, 0.20, 0.18, 0.12, "sales_fact joins\nonly on product_key", "success", tag="query path", fontsize=14)
    arrow(ax, (0.84, 0.52), (0.59, 0.32), label="drill up for rollups", label_pos=0.45)
    annotate(ax, 0.50, 0.12, "Normalization keeps the hierarchy reusable, but every extra level adds another join", fontsize=13)
    save(fig, WEEK05_DIR / "week5_product_hierarchy_rebuild.png")


def make_week05_partition_pruning():
    fig, ax = setup_figure(14, 7.0)
    panel(ax, 0.02, 0.08, 0.96, 0.84, "Partition Pruning Cost Intuition")

    # partition strip
    start_x = 0.55
    width = 0.034
    gap = 0.006
    for i in range(12):
        color = PALETTE["store"]["fill"] if i == 7 else "#E3F2FD"
        edge = PALETTE["store"]["edge"] if i == 7 else "#CFD8DC"
        rect = Rectangle((start_x + i * (width + gap), 0.48), width, 0.18, facecolor=color, edgecolor=edge, linewidth=1.2)
        ax.add_patch(rect)
    box(ax, 0.09, 0.56, 0.26, 0.14, "Without filter\nscan 365 / 365 partitions\n~1 TB", "failure", tag="bad query", fontsize=14)
    box(ax, 0.09, 0.28, 0.26, 0.14, "With date_key filter\nscan 1 / 365 partitions\n~2.7 GB", "success", tag="good query", fontsize=14)
    arrow(ax, (0.35, 0.63), (0.55, 0.57), label="full scan", dy=0.03)
    arrow(ax, (0.35, 0.35), (start_x + 7 * (width + gap) + width / 2, 0.48), label="pruned to one day", dy=0.03)
    annotate(ax, 0.77, 0.40, "Highlighted partition = only the day requested by the dashboard", fontsize=12)
    save(fig, WEEK05_DIR / "week5_partition_pruning_rebuild.png")


def make_week05_join_cost():
    fig, ax = setup_figure(14, 7.0)
    panel(ax, 0.03, 0.10, 0.44, 0.80, "Cheap Join")
    panel(ax, 0.53, 0.10, 0.44, 0.80, "Expensive Join")

    # cheap
    box(ax, 0.08, 0.56, 0.14, 0.12, "Pruned fact\n2.7 GB", "store", tag="scan", fontsize=14)
    box(ax, 0.28, 0.56, 0.14, 0.12, "Small dims\n20 MB", "source", tag="broadcast", fontsize=14)
    box(ax, 0.18, 0.30, 0.16, 0.14, "Broadcast join\nfast aggregate", "success", tag="engine", fontsize=14)
    arrow(ax, (0.22, 0.56), (0.24, 0.44))
    arrow(ax, (0.35, 0.56), (0.28, 0.44))

    # expensive
    box(ax, 0.60, 0.56, 0.14, 0.12, "Pruned fact\n2.7 GB", "store", tag="scan", fontsize=14)
    box(ax, 0.80, 0.56, 0.14, 0.12, "Large snowflaked dims\n60 GB", "failure", tag="shuffle", fontsize=13)
    box(ax, 0.70, 0.30, 0.16, 0.14, "Shuffle + spill\nslow join", "failure", tag="engine", fontsize=14)
    arrow(ax, (0.74, 0.56), (0.76, 0.44))
    arrow(ax, (0.87, 0.56), (0.78, 0.44))

    annotate(ax, 0.50, 0.05, "After partition pruning, dimension size decides whether the engine broadcasts or shuffles", fontsize=12)
    save(fig, WEEK05_DIR / "week5_join_cost_rebuild.png")


def make_week05_query_flow():
    fig, ax = setup_figure(14, 6.8)
    panel(ax, 0.02, 0.12, 0.96, 0.76, "BI Query Execution Flow")

    steps = [
        ("Dashboard SQL\nWHERE date_key...", "source"),
        ("Planner checks\npartition predicate", "process"),
        ("Prune fact\npartitions", "success"),
        ("Join customer /\nproduct dims", "process"),
        ("Aggregate\nrevenue by region", "process"),
        ("Small result set\nfor BI", "store"),
    ]
    xs = [0.05, 0.22, 0.39, 0.56, 0.73, 0.88]
    widths = [0.12, 0.13, 0.13, 0.13, 0.13, 0.09]
    for (label, kind), x, w in zip(steps, xs, widths):
        box(ax, x, 0.42, w, 0.18, label, kind, tag="step", fontsize=13)
    for i in range(len(xs) - 1):
        arrow(ax, (xs[i] + widths[i], 0.51), (xs[i + 1], 0.51))
    annotate(ax, 0.50, 0.22, "The expensive work happens early: prune first, then join, then aggregate", fontsize=13)
    save(fig, WEEK05_DIR / "week5_query_flow_rebuild.png")


def generate_diagrams():
    make_week04_hybrid()
    make_week04_extraction()
    make_week04_safe_loading()
    make_week04_etl_vs_elt()
    make_week04_cdc_options()
    make_week04_cdc_rerun()
    make_week04_sttm()
    make_week04_idempotency()
    make_week04_merge_vs_overwrite()
    make_week05_star_vs_snowflake()
    make_week05_product_hierarchy()
    make_week05_partition_pruning()
    make_week05_join_cost()
    make_week05_query_flow()


def remove_shape(shape):
    element = shape._element
    parent = element.getparent()
    parent.remove(element)


def remove_pictures(slide):
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            remove_shape(shape)


def remove_text_containing(slide, needle: str):
    for shape in list(slide.shapes):
        if hasattr(shape, "text") and needle in (shape.text or ""):
            remove_shape(shape)


def set_box(shape, left, top, width, height):
    shape.left = Inches(left)
    shape.top = Inches(top)
    shape.width = Inches(width)
    shape.height = Inches(height)


def add_picture(slide, image_path: Path, left, top, width, height):
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def edit_deck_04():
    path = SLIDES_DIR / "04-DWH-1-refined.pptx"
    prs = Presentation(path)

    # Slide 10: full-width hybrid architecture
    slide = prs.slides[9]
    remove_pictures(slide)
    remove_text_containing(slide, "Bronze keeps raw orders")
    add_picture(slide, WEEK04_DIR / "week4_hybrid_architecture_rebuild.png", 0.7, 1.45, 10.0, 4.55)

    # Slide 12: extraction contracts, keep text left and enlarge visual
    slide = prs.slides[11]
    remove_pictures(slide)
    set_box(slide.shapes[1], 0.75, 1.45, 4.2, 3.65)
    add_picture(slide, WEEK04_DIR / "week4_extraction_contracts_rebuild.png", 5.45, 1.38, 5.55, 4.25)

    # Slide 14: safe loading
    slide = prs.slides[13]
    remove_pictures(slide)
    set_box(slide.shapes[1], 0.75, 1.48, 4.35, 3.55)
    add_picture(slide, WEEK04_DIR / "week4_safe_loading_rebuild.png", 5.35, 1.46, 5.65, 4.15)

    # Slide 15: replace generic ELT image with clear comparison
    slide = prs.slides[14]
    remove_pictures(slide)
    set_box(slide.shapes[1], 1.15, 1.10, 9.2, 1.05)
    add_picture(slide, WEEK04_DIR / "week4_etl_vs_elt_rebuild.png", 0.8, 2.05, 10.4, 3.65)

    # Slide 17: CDC options comparison, diagram-led
    slide = prs.slides[16]
    remove_pictures(slide)
    remove_text_containing(slide, "Full snapshot compare")
    add_picture(slide, WEEK04_DIR / "week4_cdc_options_rebuild.png", 0.8, 1.62, 10.4, 4.25)

    # Slide 18: CDC + watermark + rerun
    slide = prs.slides[17]
    remove_pictures(slide)
    remove_text_containing(slide, "Inserts: each new order")
    add_picture(slide, WEEK04_DIR / "week4_cdc_rerun_rebuild.png", 0.8, 1.62, 10.45, 4.25)

    # Slide 22: STTM flow
    slide = prs.slides[21]
    remove_pictures(slide)
    add_picture(slide, WEEK04_DIR / "week4_sttm_rebuild.png", 0.95, 1.58, 9.95, 4.35)

    # Slide 24: idempotency rule, create diagram-only slide
    slide = prs.slides[23]
    remove_pictures(slide)
    add_picture(slide, WEEK04_DIR / "week4_idempotency_rebuild.png", 0.8, 1.68, 10.4, 4.15)

    # Slide 27: merge vs overwrite, diagram-led
    slide = prs.slides[26]
    remove_pictures(slide)
    remove_text_containing(slide, "MERGE: row-level precision")
    add_picture(slide, WEEK04_DIR / "week4_merge_vs_overwrite_rebuild.png", 0.85, 1.62, 10.35, 4.25)

    prs.save(path)


def edit_deck_05():
    path = SLIDES_DIR / "05-DWH-2-refined.pptx"
    prs = Presentation(path)

    # Slide 19: star vs snowflake
    slide = prs.slides[18]
    remove_pictures(slide)
    add_picture(slide, WEEK05_DIR / "week5_star_vs_snowflake_rebuild.png", 0.82, 1.50, 10.35, 4.45)

    # Slide 20: snowflake hierarchy
    slide = prs.slides[19]
    remove_pictures(slide)
    add_picture(slide, WEEK05_DIR / "week5_product_hierarchy_rebuild.png", 0.95, 1.48, 10.1, 4.45)

    # Slide 24: partition pruning
    slide = prs.slides[23]
    remove_pictures(slide)
    set_box(slide.shapes[1], 0.72, 1.53, 4.35, 2.85)
    add_picture(slide, WEEK05_DIR / "week5_partition_pruning_rebuild.png", 5.38, 1.46, 5.55, 4.20)

    # Slide 25: join cost
    slide = prs.slides[24]
    remove_pictures(slide)
    set_box(slide.shapes[1], 0.72, 1.55, 4.15, 2.85)
    add_picture(slide, WEEK05_DIR / "week5_join_cost_rebuild.png", 5.10, 1.46, 5.80, 4.22)

    # Slide 26: query flow
    slide = prs.slides[25]
    remove_pictures(slide)
    remove_text_containing(slide, "Query flow")
    set_box(slide.shapes[1], 0.75, 1.48, 3.85, 3.15)
    add_picture(slide, WEEK05_DIR / "week5_query_flow_rebuild.png", 4.75, 1.46, 5.95, 4.18)

    prs.save(path)


def main():
    generate_diagrams()
    edit_deck_04()
    edit_deck_05()


if __name__ == "__main__":
    main()
