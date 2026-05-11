from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "build/06-map-reduce.applied.v9-reviewed.pptx"
OUT = ROOT / "build/06-map-reduce.applied.v10-reference-styled.pptx"

RED = RGBColor(196, 111, 104)
BLACK = RGBColor(0, 0, 0)
MUTED = RGBColor(85, 85, 85)


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def style_run(run, size, bold=False, color=BLACK):
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(shape, paragraphs, sizes, bolds=None, aligns=None, color=BLACK):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bolds = bolds or [False] * len(paragraphs)
    aligns = aligns or [PP_ALIGN.LEFT] * len(paragraphs)
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = aligns[i]
        p.space_after = Pt(5 if i else 6)
        for r in p.runs:
            style_run(r, sizes[i], bolds[i], color)


def no_line(shape):
    try:
        shape.line.fill.background()
    except Exception:
        pass


def add_reference_title_frame(slide, title_text, slide_no):
    # Reference template content title frame: thin red outline, no fill, wide top slot.
    title = slide.shapes[0]
    frame = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.67), Inches(0.30), Inches(12.0), Inches(0.70)
    )
    frame.fill.background()
    frame.line.color.rgb = RED
    frame.line.width = Pt(1.0)
    # Keep frame behind title text by moving XML element before the title placeholder.
    sp_tree = slide.shapes._spTree
    sp_tree.remove(frame._element)
    sp_tree.insert(2, frame._element)

    title.left, title.top, title.width, title.height = Inches(0.77), Inches(0.33), Inches(11.80), Inches(0.62)
    no_line(title)
    set_text(title, [title_text], [29 if len(title_text) > 34 else 31], [True], [PP_ALIGN.CENTER])

    # Preserve footer shape content but normalize size/color.
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            txt = shape.text_frame.text.strip()
            if txt == "Or Peretz":
                shape.left, shape.top, shape.width, shape.height = Inches(4.56), Inches(6.95), Inches(4.22), Inches(0.28)
                set_text(shape, ["Or Peretz"], [7.5], [False], [PP_ALIGN.CENTER], MUTED)
                no_line(shape)
            elif txt == str(slide_no):
                shape.left, shape.top, shape.width, shape.height = Inches(9.56), Inches(6.95), Inches(3.11), Inches(0.28)
                set_text(shape, [str(slide_no)], [10], [False], [PP_ALIGN.RIGHT], MUTED)
                no_line(shape)


def get_text(slide, idx):
    shape = slide.shapes[idx]
    return " | ".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())


def normalize_titles(prs):
    for i, slide in enumerate(prs.slides, 1):
        title = get_text(slide, 0).split(" | ")[0].strip() if getattr(slide.shapes[0], "has_text_frame", False) else ""
        if i == 1:
            # Title slide uses the reference center title frame.
            title_shape = slide.shapes[0]
            frame = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.35), Inches(11.33), Inches(0.78)
            )
            frame.fill.background()
            frame.line.color.rgb = RED
            frame.line.width = Pt(1.0)
            title_shape.left, title_shape.top = Inches(1.1), Inches(2.38)
            title_shape.width, title_shape.height = Inches(11.1), Inches(0.72)
            set_text(title_shape, [title or "MapReduce Algorithm"], [38], [True], [PP_ALIGN.CENTER])
            no_line(title_shape)
        else:
            add_reference_title_frame(slide, title, i)


def fix_formal_model(slide):
    # Top title is already applied. Use a clean two-column body below the reference frame.
    left = slide.shapes[0]
    # Find/remove old explanatory text box fragments; title placeholder currently contains only title.
    body = slide.shapes.add_textbox(Inches(0.98), Inches(1.55), Inches(4.10), Inches(3.55))
    set_text(
        body,
        [
            "D is the input multiset.",
            "Map emits intermediate pairs.",
            "Reduce receives one complete key group.",
            "K₂ is the grouping key.",
            "M(X) denotes a multiset over domain X.",
        ],
        [20, 20, 20, 20, 15],
        [False] * 5,
        [PP_ALIGN.LEFT] * 5,
    )
    no_line(body)

    # Move equation images into a compact, reference-template body area.
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    positions = [(5.45, 1.45, 3.95), (5.45, 2.45, 5.10), (5.45, 3.45, 5.60)]
    for pic, (x, y, w) in zip(pics, positions):
        ratio = pic.height / pic.width
        pic.left, pic.top, pic.width, pic.height = Inches(x), Inches(y), Inches(w), Inches(w * ratio)
    # Remove the old note added in v9 because it is now part of the left text.
    for sh in list(slide.shapes):
        if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip().startswith("M(X) denotes"):
            remove_shape(sh)


def fix_formal_algorithm(slide):
    body = slide.shapes[1]
    # The old title placeholder still carries the step text in v9; replace with title-only.
    set_text(body, ["Formal Algorithm"], [31], [True], [PP_ALIGN.CENTER])
    left = slide.shapes.add_textbox(Inches(0.98), Inches(1.55), Inches(4.10), Inches(3.65))
    set_text(
        left,
        [
            "1. Map all records.",
            "2. Partition by key.",
            "3. Group values by key.",
            "4. Reduce each group.",
            "Invariant: same key -> complete group.",
        ],
        [19, 19, 19, 19, 16],
        [False] * 5,
        [PP_ALIGN.LEFT] * 5,
    )
    no_line(left)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    positions = [(5.25, 1.35, 5.75), (5.25, 2.22, 3.10), (5.25, 3.10, 4.65), (5.25, 4.02, 5.90)]
    for pic, (x, y, w) in zip(pics, positions):
        ratio = pic.height / pic.width
        pic.left, pic.top, pic.width, pic.height = Inches(x), Inches(y), Inches(w), Inches(w * ratio)
    for sh in slide.shapes:
        if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip().startswith("merge ="):
            sh.left, sh.top, sh.width, sh.height = Inches(5.30), Inches(5.05), Inches(6.0), Inches(0.30)
            set_text(sh, ["merge = multiset union/concatenation of emitted pairs."], [13.5], [False], [PP_ALIGN.LEFT], MUTED)


def normalize_step_subtitles(prs):
    subtitles = {
        10: "Each split is read independently; map emits (k2, v2) pairs.",
        11: "Local pre-aggregation shrinks repeated keys before shuffle.",
        12: "System partitions by key and groups equal keys together.",
        13: "Reducer receives (k2, [v2]) and emits final records.",
        15: "Goal: count each token across all input lines.",
        16: "Each line can be processed by a separate mapper.",
        17: "For every token w, emit one pair (w, 1).",
        18: "All pairs with the same word are routed to the same reducer.",
        19: "Each reducer sums the list of 1s for its word.",
    }
    for si in [10, 11, 12, 13, 15, 16, 17, 18, 19]:
        slide = prs.slides[si - 1]
        title = slide.shapes[0]
        lines = [p.text for p in title.text_frame.paragraphs if p.text.strip()]
        title_text = lines[0] if lines else ""
        subtitle = subtitles.get(si, lines[1] if len(lines) > 1 else "")
        set_text(title, [title_text], [30], [True], [PP_ALIGN.CENTER])
        if subtitle:
            sub = slide.shapes.add_textbox(Inches(1.15), Inches(1.08), Inches(11.0), Inches(0.34))
            set_text(sub, [subtitle], [15.5], [False], [PP_ALIGN.CENTER], BLACK)
            no_line(sub)


def tune_step_spacing(prs):
    subtitle_text = {
        10: "Each split is read independently; map emits (k2, v2) pairs.",
        11: "Local pre-aggregation shrinks repeated keys before shuffle.",
        12: "System partitions by key and groups equal keys together.",
        13: "Reducer receives (k2, [v2]) and emits final records.",
        15: "Goal: count each token across all input lines.",
        16: "Each line can be processed by a separate mapper.",
        17: "For every token w, emit one pair (w, 1).",
        18: "All pairs with the same word are routed to the same reducer.",
        19: "Each reducer sums the list of 1s for its word.",
    }
    for si in [12, 13]:
        slide = prs.slides[si - 1]
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip() == subtitle_text[si]:
                continue
            if sh.top is not None and Inches(1.0) <= sh.top < Inches(6.7):
                sh.top += Inches(0.35)

    phase_labels = {"Input", "Split", "Map", "Shuffle & Sort", "Reduce"}
    for si in [15, 16, 17, 18, 19]:
        slide = prs.slides[si - 1]
        for sh in slide.shapes:
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip() in phase_labels:
                sh.top = Inches(1.48)


def tune_slide_20(slide):
    # Keep the math and diagram aligned under the reference title slot.
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.left > Inches(6):
            shape.left, shape.top, shape.width, shape.height = Inches(7.15), Inches(1.35), Inches(2.80), Inches(4.45)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and "emitted intermediate" in shape.text_frame.text:
            shape.left, shape.top, shape.width, shape.height = Inches(1.05), Inches(2.70), Inches(5.2), Inches(2.55)
            set_text(
                shape,
                [
                    "E = emitted intermediate pairs",
                    "s = serialized pair size",
                    "Use combiners to reduce E.",
                    "Use compact encoding to reduce s.",
                    "Slowest reducer bounds completion.",
                ],
                [18, 18, 18, 18, 18],
                [False] * 5,
                [PP_ALIGN.LEFT] * 5,
            )
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE and s.left < Inches(6)]
    for pic, x, y, w in zip(pics, [1.05, 1.05], [1.35, 1.88], [3.0, 4.65]):
        ratio = pic.height / pic.width
        pic.left, pic.top, pic.width, pic.height = Inches(x), Inches(y), Inches(w), Inches(w * ratio)


def tune_mixed_slide(slide, image_x=4.65, image_y=1.35, image_w=6.65, image_h=3.85):
    title_shape = slide.shapes[1] if len(slide.shapes) > 1 else None
    title_text = title_shape.text_frame.text.strip() if title_shape is not None and getattr(title_shape, "has_text_frame", False) else ""
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape is not slide.shapes[0]:
            txt = shape.text_frame.text.strip()
            if txt == title_text:
                continue
            if txt and txt not in {"Or Peretz"} and not txt.isdigit():
                shape.left, shape.top, shape.width, shape.height = Inches(0.90), Inches(1.55), Inches(3.55), Inches(3.35)
                size = 18 if len(txt) > 120 else 19
                set_text(shape, [p.text for p in shape.text_frame.paragraphs if p.text.strip()], [size] * len([p for p in shape.text_frame.paragraphs if p.text.strip()]))
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape.left, shape.top, shape.width, shape.height = Inches(image_x), Inches(image_y), Inches(image_w), Inches(image_h)
    if title_shape is not None:
        title_shape.left, title_shape.top, title_shape.width, title_shape.height = Inches(0.77), Inches(0.33), Inches(11.80), Inches(0.62)
        set_text(title_shape, [title_text], [29 if len(title_text) > 34 else 31], [True], [PP_ALIGN.CENTER])


def remove_off_canvas(prs):
    for slide in prs.slides:
        for sh in list(slide.shapes):
            if sh.left is not None and sh.top is not None and (sh.left >= prs.slide_width or sh.top >= prs.slide_height):
                remove_shape(sh)


def main():
    prs = Presentation(SRC)
    normalize_titles(prs)
    fix_formal_model(prs.slides[4])
    fix_formal_algorithm(prs.slides[5])
    normalize_step_subtitles(prs)
    tune_step_spacing(prs)
    tune_slide_20(prs.slides[19])
    tune_mixed_slide(prs.slides[21], 4.70, 1.32, 6.55, 3.82)
    tune_mixed_slide(prs.slides[24], 4.55, 1.35, 6.45, 3.80)
    remove_off_canvas(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
