#!/usr/bin/env python3

from __future__ import annotations

import math
import textwrap
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parent.parent
SLIDES_DIR = ROOT / "slides"
ASSETS_DIR = SLIDES_DIR / "generated_diagrams"

FONT_REGULAR = "/System/Library/Fonts/Supplemental/Arial.ttf"
FONT_BOLD = "/System/Library/Fonts/Supplemental/Arial Bold.ttf"

BG = "#FBFAF8"
INK = "#1F2937"
MUTED = "#5B6472"
LINE = "#D8DDE4"
NAVY = "#1E3A5F"
BLUE = "#3D7EDB"
BLUE_SOFT = "#DCEBFF"
TEAL = "#0F9D8A"
TEAL_SOFT = "#D7F5EF"
GREEN = "#5E9B56"
GREEN_SOFT = "#E4F4DC"
GOLD = "#E5A93D"
GOLD_SOFT = "#FBECCB"
CORAL = "#E76F51"
CORAL_SOFT = "#FCE1DB"
PURPLE = "#7C5CE0"
PURPLE_SOFT = "#E9E2FF"
GRAY_SOFT = "#F2F4F7"
WHITE = "#FFFFFF"


@dataclass
class SlideUpdate:
    slide_no: int
    key: str
    mode: str = "replace_picture"


class Canvas:
    def __init__(self, width: int, height: int):
        self.w = width
        self.h = height
        self.image = Image.new("RGB", (width, height), BG)
        self.draw = ImageDraw.Draw(self.image)

    def x(self, pct: float) -> int:
        return int(self.w * pct / 100.0)

    def y(self, pct: float) -> int:
        return int(self.h * pct / 100.0)


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REGULAR, size)


def box(draw: ImageDraw.ImageDraw, rect, fill, outline=LINE, radius=28, width=3):
    draw.rounded_rectangle(rect, radius=radius, fill=fill, outline=outline, width=width)


def text_size(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.FreeTypeFont) -> tuple[int, int]:
    left, top, right, bottom = draw.textbbox((0, 0), text, font=fnt)
    return right - left, bottom - top


def wrap_lines(draw: ImageDraw.ImageDraw, text: str, fnt: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    lines: list[str] = []
    for para in text.split("\n"):
        words = para.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            trial = f"{current} {word}"
            if text_size(draw, trial, fnt)[0] <= max_width:
                current = trial
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return lines


def draw_multiline(
    draw: ImageDraw.ImageDraw,
    rect,
    text: str,
    size: int,
    color: str = INK,
    bold: bool = False,
    align: str = "center",
    line_gap: int = 10,
    valign: str = "middle",
):
    x1, y1, x2, y2 = rect
    fnt = font(size, bold=bold)
    max_width = max(40, x2 - x1 - 24)
    lines = wrap_lines(draw, text, fnt, max_width)
    metrics = [text_size(draw, line or " ", fnt) for line in lines]
    total_h = sum(h for _, h in metrics) + line_gap * max(0, len(lines) - 1)
    if valign == "top":
        y = y1
    elif valign == "bottom":
        y = y2 - total_h
    else:
        y = y1 + ((y2 - y1) - total_h) / 2
    for line, (w, h) in zip(lines, metrics):
        if align == "left":
            x = x1
        elif align == "right":
            x = x2 - w
        else:
            x = x1 + ((x2 - x1) - w) / 2
        draw.text((x, y), line, fill=color, font=fnt)
        y += h + line_gap


def draw_list(draw: ImageDraw.ImageDraw, x: int, y: int, width: int, items: list[str], size: int = 24, color: str = INK):
    fnt = font(size)
    bullet = font(size + 4, bold=True)
    cursor = y
    for item in items:
        draw.text((x, cursor), "•", fill=color, font=bullet)
        lines = wrap_lines(draw, item, fnt, width - 34)
        line_h = text_size(draw, "Ag", fnt)[1]
        for idx, line in enumerate(lines):
            draw.text((x + 32, cursor + idx * (line_h + 6)), line, fill=color, font=fnt)
        cursor += len(lines) * (line_h + 6) + 14


def arrow(draw: ImageDraw.ImageDraw, start, end, color=BLUE, width=8, head=18):
    x1, y1 = start
    x2, y2 = end
    draw.line((x1, y1, x2, y2), fill=color, width=width)
    angle = math.atan2(y2 - y1, x2 - x1)
    left = (
        x2 - head * math.cos(angle) + head * 0.6 * math.sin(angle),
        y2 - head * math.sin(angle) - head * 0.6 * math.cos(angle),
    )
    right = (
        x2 - head * math.cos(angle) - head * 0.6 * math.sin(angle),
        y2 - head * math.sin(angle) + head * 0.6 * math.cos(angle),
    )
    draw.polygon([end, left, right], fill=color)


def pill(draw: ImageDraw.ImageDraw, rect, text, fill, text_fill=INK):
    box(draw, rect, fill=fill, outline=fill, radius=(rect[3] - rect[1]) // 2, width=1)
    draw_multiline(draw, rect, text, size=22, color=text_fill, bold=True)


def card(draw: ImageDraw.ImageDraw, rect, title: str, body: str | None = None, fill=WHITE, accent=NAVY, title_size=30, body_size=22):
    x1, y1, x2, y2 = rect
    box(draw, rect, fill=fill, outline=LINE, radius=30, width=3)
    draw.rounded_rectangle((x1, y1, x2, y1 + 18), radius=30, fill=accent)
    draw_multiline(draw, (x1 + 24, y1 + 28, x2 - 24, y1 + 92), title, size=title_size, color=INK, bold=True)
    if body:
        draw_multiline(draw, (x1 + 24, y1 + 104, x2 - 24, y2 - 24), body, size=body_size, color=MUTED, bold=False, align="left", valign="top")


def draw_grid_background(canvas: Canvas):
    for x in range(0, canvas.w, 80):
        canvas.draw.line((x, 0, x, canvas.h), fill="#F3F4F6", width=1)
    for y in range(0, canvas.h, 80):
        canvas.draw.line((0, y, canvas.w, y), fill="#F3F4F6", width=1)


def diagram_dw_vs_lake(canvas: Canvas):
    draw = canvas.draw
    draw_grid_background(canvas)
    left = (canvas.x(5), canvas.y(16), canvas.x(45), canvas.y(83))
    right = (canvas.x(55), canvas.y(16), canvas.x(95), canvas.y(83))
    card(draw, left, "Warehouse", fill=WHITE, accent=BLUE)
    card(draw, right, "Lake", fill=WHITE, accent=TEAL)
    draw_list(draw, left[0] + 28, left[1] + 102, left[2] - left[0] - 56, [
        "Curated tables, strict contracts",
        "Fast BI and certified metrics",
        "Schema-on-write, governance first",
    ], size=24)
    draw_list(draw, right[0] + 28, right[1] + 102, right[2] - right[0] - 56, [
        "Raw + processed zones",
        "Cheap storage for high-volume data",
        "Flexible schema-on-read exploration",
    ], size=24)
    band = (canvas.x(24), canvas.y(6), canvas.x(76), canvas.y(13))
    pill(draw, band, "Most teams keep both: lake for raw, warehouse for trusted BI", GOLD_SOFT, NAVY)
    arrow(draw, (canvas.x(45), canvas.y(49)), (canvas.x(55), canvas.y(49)), color=GOLD, width=10, head=22)


def diagram_dw_choice(canvas: Canvas):
    draw = canvas.draw
    headers = ["Decision lens", "Warehouse", "Lake"]
    cols = [canvas.x(6), canvas.x(36), canvas.x(63), canvas.x(94)]
    top = canvas.y(14)
    row_h = canvas.y(14)
    labels = [
        ("Primary use", "Certified KPIs", "Raw capture + discovery"),
        ("Latency shape", "Fast BI dashboards", "Flexible downstream transforms"),
        ("Best when", "Business needs one truth", "Volume and variability dominate"),
        ("Risk", "More modeling upfront", "Governance can drift"),
    ]
    for i in range(3):
        box(draw, (cols[i], top, cols[i + 1], top + row_h), fill=[GRAY_SOFT, BLUE_SOFT, TEAL_SOFT][i], outline=LINE, radius=24)
        draw_multiline(draw, (cols[i] + 8, top + 10, cols[i + 1] - 8, top + row_h - 10), headers[i], size=28, color=INK, bold=True)
    y = top + row_h + 18
    for label, wh, lake in labels:
        values = [label, wh, lake]
        fills = [WHITE, "#F7FBFF", "#F5FCFA"]
        for i in range(3):
            box(draw, (cols[i], y, cols[i + 1], y + row_h), fill=fills[i], outline=LINE, radius=24)
            draw_multiline(draw, (cols[i] + 16, y + 14, cols[i + 1] - 16, y + row_h - 14), values[i], size=24, color=INK if i == 0 else MUTED, bold=i == 0, align="left")
        y += row_h + 12
    pill(draw, (canvas.x(28), canvas.y(86), canvas.x(72), canvas.y(93)), "Choose by contract strength, not by trend", GOLD_SOFT, NAVY)


def diagram_why_dw(canvas: Canvas):
    draw = canvas.draw
    sources = [(10, 18, "OLTP orders"), (10, 41, "CRM profiles"), (10, 64, "Payments / logs")]
    warehouse_rect = (canvas.x(39), canvas.y(24), canvas.x(63), canvas.y(76))
    bi_rect = (canvas.x(75), canvas.y(28), canvas.x(93), canvas.y(72))
    for x, y, label in sources:
        rect = (canvas.x(x), canvas.y(y), canvas.x(x + 18), canvas.y(y + 14))
        card(draw, rect, label, fill=WHITE, accent=CORAL, title_size=24)
        arrow(draw, (rect[2], (rect[1] + rect[3]) // 2), (warehouse_rect[0] - 12, canvas.y(y + 7)), color=CORAL, width=7, head=16)
    card(draw, warehouse_rect, "Curated warehouse", "Historical snapshots\nStandardized codes\nHeavy joins live here", fill=BLUE_SOFT, accent=BLUE, title_size=30, body_size=23)
    arrow(draw, (warehouse_rect[2], canvas.y(50)), (bi_rect[0] - 14, canvas.y(50)), color=TEAL, width=9, head=18)
    card(draw, bi_rect, "BI", "Fast KPI reads", fill=GREEN_SOFT, accent=GREEN, title_size=28, body_size=24)
    pill(draw, (canvas.x(33), canvas.y(10), canvas.x(69), canvas.y(17)), "Separate analytics protects production systems", GOLD_SOFT, NAVY)


def diagram_schema_strategy(canvas: Canvas):
    draw = canvas.draw
    for idx, (title, y, accent, soft, subtitle) in enumerate([
        ("Schema-on-write", 18, BLUE, BLUE_SOFT, "Validate before load"),
        ("Schema-on-read", 55, TEAL, TEAL_SOFT, "Interpret at query time"),
    ]):
        lane = (canvas.x(7), canvas.y(y), canvas.x(93), canvas.y(y + 22))
        box(draw, lane, fill=soft, outline=accent, radius=28, width=4)
        card(draw, (canvas.x(10), canvas.y(y + 3), canvas.x(24), canvas.y(y + 19)), title, subtitle, fill=WHITE, accent=accent, title_size=26, body_size=20)
        steps = ["Raw source", "Check rules", "Store", "Query"] if idx == 0 else ["Raw source", "Store raw", "Apply logic", "Query"]
        for i, step in enumerate(steps):
            sx = 29 + i * 15
            rect = (canvas.x(sx), canvas.y(y + 6), canvas.x(sx + 12), canvas.y(y + 16))
            card(draw, rect, step, fill=WHITE, accent=accent, title_size=22)
            if i < len(steps) - 1:
                arrow(draw, (canvas.x(sx + 12), canvas.y(y + 11)), (canvas.x(sx + 15) - 8, canvas.y(y + 11)), color=accent, width=6, head=14)
    pill(draw, (canvas.x(32), canvas.y(83), canvas.x(69), canvas.y(90)), "Hybrid reality: enforce core contracts, keep raw for change", GOLD_SOFT, NAVY)


def diagram_hybrid_architecture(canvas: Canvas):
    draw = canvas.draw
    sources = [
        ("Postgres\norders", 8, 28),
        ("Stripe\npayments", 8, 46),
        ("CRM\ncustomers", 8, 64),
        ("Web/app\nlogs", 8, 82),
    ]
    for label, x, y in sources:
        rect = (canvas.x(x), canvas.y(y - 8), canvas.x(x + 16), canvas.y(y + 4))
        card(draw, rect, label, fill=WHITE, accent=CORAL, title_size=22)
    bronze = (canvas.x(32), canvas.y(23), canvas.x(48), canvas.y(39))
    silver = (canvas.x(55), canvas.y(23), canvas.x(71), canvas.y(39))
    gold = (canvas.x(78), canvas.y(23), canvas.x(94), canvas.y(39))
    card(draw, bronze, "Bronze", "Raw landing", fill=BLUE_SOFT, accent=BLUE, title_size=28, body_size=20)
    card(draw, silver, "Silver", "Standardized", fill=TEAL_SOFT, accent=TEAL, title_size=28, body_size=20)
    card(draw, gold, "Gold", "Trusted marts", fill=GREEN_SOFT, accent=GREEN, title_size=28, body_size=20)
    for _, x, y in sources:
        arrow(draw, (canvas.x(x + 16), canvas.y(y - 2)), (bronze[0] - 10, canvas.y(31)), color=CORAL, width=6, head=14)
    arrow(draw, (bronze[2], canvas.y(31)), (silver[0] - 8, canvas.y(31)), color=BLUE, width=8, head=18)
    arrow(draw, (silver[2], canvas.y(31)), (gold[0] - 8, canvas.y(31)), color=TEAL, width=8, head=18)
    bi = (canvas.x(70), canvas.y(56), canvas.x(92), canvas.y(76))
    card(draw, bi, "Daily revenue +\ncustomer view", "Dashboards and KPI contracts", fill=WHITE, accent=GOLD, title_size=26, body_size=20)
    arrow(draw, (gold[0] + (gold[2] - gold[0]) // 2, gold[3]), (bi[0] + 26, bi[1] - 10), color=GOLD, width=8, head=18)
    pill(draw, (canvas.x(31), canvas.y(8), canvas.x(72), canvas.y(15)), "Hybrid = raw flexibility on the left, governed outputs on the right", GRAY_SOFT, NAVY)


def diagram_extraction_sources(canvas: Canvas):
    draw = canvas.draw
    columns = [10, 29, 48, 67]
    items = [
        ("Postgres", "Incremental query"),
        ("Stripe API", "Windowed pull"),
        ("CRM", "Master-data extract"),
        ("Log files", "Object-store pickup"),
    ]
    for x, (title, body) in zip(columns, items):
        rect = (canvas.x(x), canvas.y(18), canvas.x(x + 15), canvas.y(46))
        card(draw, rect, title, body, fill=WHITE, accent=CORAL, title_size=24, body_size=20)
        arrow(draw, (canvas.x(x + 7.5), canvas.y(46)), (canvas.x(x + 7.5), canvas.y(58)), color=CORAL, width=6, head=14)
    landing = (canvas.x(18), canvas.y(60), canvas.x(82), canvas.y(83))
    card(draw, landing, "Landing / staging zone", "Raw extracts arrive with source timestamp, batch id, and retry-safe metadata.", fill=BLUE_SOFT, accent=BLUE, title_size=30, body_size=22)
    pill(draw, (canvas.x(31), canvas.y(6), canvas.x(69), canvas.y(13)), "Extract incrementally so source systems stay responsive", GOLD_SOFT, NAVY)


def diagram_safe_writes(canvas: Canvas):
    draw = canvas.draw
    steps = [
        ("Batch D", BLUE, BLUE_SOFT),
        ("Staging +\ndedup", TEAL, TEAL_SOFT),
        ("MERGE\nfacts", GREEN, GREEN_SOFT),
        ("Update\ncheckpoint", GOLD, GOLD_SOFT),
        ("Trusted\ntarget", NAVY, GRAY_SOFT),
    ]
    xs = [8, 25, 42, 59, 77]
    for x, (title, accent, soft) in zip(xs, steps):
        rect = (canvas.x(x), canvas.y(34), canvas.x(x + 14), canvas.y(56))
        card(draw, rect, title, fill=soft, accent=accent, title_size=24)
    for i in range(len(xs) - 1):
        arrow(draw, (canvas.x(xs[i] + 14), canvas.y(45)), (canvas.x(xs[i + 1]) - 10, canvas.y(45)), color=BLUE, width=8, head=18)
    note = (canvas.x(22), canvas.y(68), canvas.x(81), canvas.y(84))
    card(draw, note, "Rerun safety", "Same batch can run again because dedup + merge protect the target from duplicates.", fill=WHITE, accent=CORAL, title_size=28, body_size=22)


def diagram_elt(canvas: Canvas):
    draw = canvas.draw
    raw = (canvas.x(12), canvas.y(25), canvas.x(28), canvas.y(44))
    load = (canvas.x(38), canvas.y(25), canvas.x(55), canvas.y(44))
    transform = (canvas.x(65), canvas.y(25), canvas.x(82), canvas.y(44))
    serve = (canvas.x(38), canvas.y(58), canvas.x(55), canvas.y(78))
    card(draw, raw, "Raw data", "APIs, DBs, files", fill=WHITE, accent=CORAL, title_size=28, body_size=22)
    card(draw, load, "Load first", "Land raw in the platform", fill=BLUE_SOFT, accent=BLUE, title_size=28, body_size=22)
    card(draw, transform, "Transform in\nwarehouse", "SQL / dbt / lakehouse compute", fill=TEAL_SOFT, accent=TEAL, title_size=26, body_size=20)
    card(draw, serve, "Trusted marts", "Business-ready outputs", fill=GREEN_SOFT, accent=GREEN, title_size=28, body_size=22)
    arrow(draw, (raw[2], canvas.y(34.5)), (load[0] - 10, canvas.y(34.5)), color=CORAL, width=8, head=18)
    arrow(draw, (load[2], canvas.y(34.5)), (transform[0] - 10, canvas.y(34.5)), color=BLUE, width=8, head=18)
    arrow(draw, (transform[0] + (transform[2] - transform[0]) // 2, transform[3]), (serve[0] + (serve[2] - serve[0]) // 2, serve[1] - 10), color=TEAL, width=8, head=18)
    pill(draw, (canvas.x(28), canvas.y(8), canvas.x(72), canvas.y(15)), "ELT shifts transformation into elastic analytics compute", GOLD_SOFT, NAVY)


def diagram_cdc_options(canvas: Canvas):
    draw = canvas.draw
    options = [
        ("Timestamp", "Easy to adopt", "Needs reliable updated_at", BLUE, BLUE_SOFT),
        ("Log based", "Best fidelity", "Requires WAL/binlog access", TEAL, TEAL_SOFT),
        ("Diff snapshot", "Works anywhere", "Expensive at scale", GOLD, GOLD_SOFT),
        ("Trigger", "Immediate signal", "Operational overhead", CORAL, CORAL_SOFT),
    ]
    for idx, (title, best, risk, accent, soft) in enumerate(options):
        x1 = canvas.x(5 + idx * 23)
        x2 = canvas.x(25 + idx * 23)
        rect = (x1, canvas.y(18), x2, canvas.y(78))
        card(draw, rect, title, fill=soft, accent=accent, title_size=28)
        pill(draw, (x1 + 20, canvas.y(36), x2 - 20, canvas.y(43)), "Best fit", WHITE, accent)
        draw_multiline(draw, (x1 + 20, canvas.y(44), x2 - 20, canvas.y(53)), best, size=22, color=INK)
        pill(draw, (x1 + 20, canvas.y(58), x2 - 20, canvas.y(65)), "Watch for", WHITE, accent)
        draw_multiline(draw, (x1 + 20, canvas.y(66), x2 - 20, canvas.y(75)), risk, size=20, color=MUTED)


def diagram_cdc_example(canvas: Canvas):
    draw = canvas.draw
    steps = [
        ("Orders table", CORAL, CORAL_SOFT),
        ("Read rows >\nwatermark", BLUE, BLUE_SOFT),
        ("Load batch D", TEAL, TEAL_SOFT),
        ("Commit new\nwatermark", GREEN, GREEN_SOFT),
    ]
    xs = [6, 28, 50, 72]
    for x, (title, accent, soft) in zip(xs, steps):
        rect = (canvas.x(x), canvas.y(34), canvas.x(x + 18), canvas.y(56))
        card(draw, rect, title, fill=soft, accent=accent, title_size=25)
    for i in range(len(xs) - 1):
        arrow(draw, (canvas.x(xs[i] + 18), canvas.y(45)), (canvas.x(xs[i + 1]) - 10, canvas.y(45)), color=BLUE, width=8, head=18)
    draw_multiline(draw, (canvas.x(22), canvas.y(12), canvas.x(78), canvas.y(22)), "If batch D fails, rerun from the last committed watermark", size=28, color=NAVY, bold=True)
    arrow(draw, (canvas.x(82), canvas.y(60)), (canvas.x(58), canvas.y(74)), color=CORAL, width=7, head=16)
    card(draw, (canvas.x(45), canvas.y(72), canvas.x(82), canvas.y(86)), "Rerun path", "Target stays correct because the checkpoint moves only after success.", fill=WHITE, accent=CORAL, title_size=27, body_size=20)


def diagram_sttm_flow(canvas: Canvas):
    draw = canvas.draw
    source = (canvas.x(7), canvas.y(28), canvas.x(28), canvas.y(70))
    rules = (canvas.x(38), canvas.y(22), canvas.x(62), canvas.y(76))
    target1 = (canvas.x(72), canvas.y(22), canvas.x(93), canvas.y(43))
    target2 = (canvas.x(72), canvas.y(55), canvas.x(93), canvas.y(76))
    card(draw, source, "Source fields", "crm_customer_id\norder_status\ncurrency\ntimezone", fill=WHITE, accent=CORAL, title_size=30, body_size=23)
    card(draw, rules, "STTM rules", "Map source keys\naudit defaults\nnormalize status\nstandardize nulls", fill=BLUE_SOFT, accent=BLUE, title_size=30, body_size=22)
    card(draw, target1, "dim_customer", "customer_key\nsegment\nregion", fill=TEAL_SOFT, accent=TEAL, title_size=26, body_size=21)
    card(draw, target2, "fact_order", "status\ncurrency\norder metrics", fill=GREEN_SOFT, accent=GREEN, title_size=26, body_size=21)
    arrow(draw, (source[2], canvas.y(49)), (rules[0] - 10, canvas.y(49)), color=CORAL, width=8, head=18)
    arrow(draw, (rules[2], canvas.y(40)), (target1[0] - 10, canvas.y(32)), color=BLUE, width=8, head=18)
    arrow(draw, (rules[2], canvas.y(58)), (target2[0] - 10, canvas.y(66)), color=BLUE, width=8, head=18)


def diagram_idempotency(canvas: Canvas):
    draw = canvas.draw
    first = (canvas.x(12), canvas.y(38), canvas.x(28), canvas.y(56))
    second = (canvas.x(34), canvas.y(38), canvas.x(50), canvas.y(56))
    rule = (canvas.x(56), canvas.y(30), canvas.x(74), canvas.y(64))
    target = (canvas.x(80), canvas.y(38), canvas.x(94), canvas.y(56))
    card(draw, first, "Batch D", fill=BLUE_SOFT, accent=BLUE, title_size=28)
    card(draw, second, "Rerun D", fill=BLUE_SOFT, accent=BLUE, title_size=28)
    card(draw, rule, "Dedup +\nMERGE rule", "Same input must\nlead to same state", fill=GOLD_SOFT, accent=GOLD, title_size=28, body_size=20)
    card(draw, target, "Target", "No extra rows", fill=GREEN_SOFT, accent=GREEN, title_size=28, body_size=22)
    arrow(draw, (first[2], canvas.y(47)), (rule[0] - 16, canvas.y(42)), color=BLUE, width=7, head=16)
    arrow(draw, (second[2], canvas.y(47)), (rule[0] - 16, canvas.y(52)), color=BLUE, width=7, head=16)
    arrow(draw, (rule[2], canvas.y(47)), (target[0] - 12, canvas.y(47)), color=GOLD, width=9, head=18)
    pill(draw, (canvas.x(28), canvas.y(12), canvas.x(72), canvas.y(19)), "Idempotent load = safe reruns after partial failure", CORAL_SOFT, NAVY)


def diagram_merge_vs_overwrite(canvas: Canvas):
    draw = canvas.draw
    left = (canvas.x(6), canvas.y(18), canvas.x(47), canvas.y(82))
    right = (canvas.x(53), canvas.y(18), canvas.x(94), canvas.y(82))
    card(draw, left, "MERGE", fill=BLUE_SOFT, accent=BLUE, title_size=34)
    draw_list(draw, left[0] + 24, left[1] + 102, left[2] - left[0] - 48, [
        "Best for row-level upserts",
        "Keeps unchanged rows in place",
        "Handles late updates cleanly",
    ], size=24)
    card(draw, right, "Partition overwrite", fill=TEAL_SOFT, accent=TEAL, title_size=34)
    draw_list(draw, right[0] + 24, right[1] + 102, right[2] - right[0] - 48, [
        "Best when entire partition is rebuilt",
        "Simple logic and fast replace",
        "Requires trusted partition boundary",
    ], size=24)
    pill(draw, (canvas.x(30), canvas.y(87), canvas.x(70), canvas.y(94)), "Choose by change granularity: rows vs full partitions", GOLD_SOFT, NAVY)


def diagram_multidimensional_model(canvas: Canvas):
    draw = canvas.draw
    fact = (canvas.x(39), canvas.y(34), canvas.x(62), canvas.y(56))
    dims = [
        ((canvas.x(13), canvas.y(18), canvas.x(31), canvas.y(34)), "Store\nDimension", CORAL, CORAL_SOFT),
        ((canvas.x(13), canvas.y(58), canvas.x(31), canvas.y(74)), "Product\nDimension", TEAL, TEAL_SOFT),
        ((canvas.x(69), canvas.y(18), canvas.x(87), canvas.y(34)), "Time\nDimension", GOLD, GOLD_SOFT),
        ((canvas.x(69), canvas.y(58), canvas.x(87), canvas.y(74)), "Measure\nrules", PURPLE, PURPLE_SOFT),
    ]
    card(draw, fact, "sales_fact", "store_key | product_key | date_key | sales | qty", fill=BLUE_SOFT, accent=BLUE, title_size=32, body_size=19)
    for rect, title, accent, soft in dims:
        card(draw, rect, title, fill=soft, accent=accent, title_size=28)
        arrow(draw, ((rect[2] if rect[0] < fact[0] else rect[0]), (rect[1] + rect[3]) // 2), ((fact[0] if rect[0] < fact[0] else fact[2]), canvas.y(45)), color=accent, width=7, head=16)
    pill(draw, (canvas.x(28), canvas.y(84), canvas.x(72), canvas.y(91)), "Fact grain in the center, dimensions add analysis context", GRAY_SOFT, NAVY)


def diagram_location_hierarchy(canvas: Canvas):
    draw = canvas.draw
    nodes = {
        "All": (50, 14),
        "Europe": (34, 30),
        "North America": (66, 30),
        "Germany": (24, 48),
        "Spain": (42, 48),
        "Canada": (58, 48),
        "Mexico": (76, 48),
        "Frankfurt": (24, 66),
        "Vancouver": (58, 66),
        "Toronto": (76, 66),
    }
    connections = [
        ("All", "Europe"), ("All", "North America"),
        ("Europe", "Germany"), ("Europe", "Spain"),
        ("North America", "Canada"), ("North America", "Mexico"),
        ("Germany", "Frankfurt"), ("Canada", "Vancouver"), ("Mexico", "Toronto"),
    ]
    for a, b in connections:
        arrow(draw, (canvas.x(nodes[a][0]), canvas.y(nodes[a][1] + 3)), (canvas.x(nodes[b][0]), canvas.y(nodes[b][1] - 3)), color=LINE, width=5, head=12)
    for title, (x, y) in nodes.items():
        rect = (canvas.x(x - 8), canvas.y(y - 4), canvas.x(x + 8), canvas.y(y + 4))
        accent = BLUE if y < 30 else TEAL if y < 48 else GOLD if y < 66 else CORAL
        soft = BLUE_SOFT if y < 30 else TEAL_SOFT if y < 48 else GOLD_SOFT if y < 66 else CORAL_SOFT
        card(draw, rect, title, fill=soft, accent=accent, title_size=22)
    pill(draw, (canvas.x(11), canvas.y(84), canvas.x(38), canvas.y(91)), "region", BLUE_SOFT, NAVY)
    pill(draw, (canvas.x(40), canvas.y(84), canvas.x(66), canvas.y(91)), "country", TEAL_SOFT, NAVY)
    pill(draw, (canvas.x(68), canvas.y(84), canvas.x(89), canvas.y(91)), "city", GOLD_SOFT, NAVY)


def diagram_scd_sequence(canvas: Canvas):
    draw = canvas.draw
    timeline = canvas.y(60)
    draw.line((canvas.x(12), timeline, canvas.x(90), timeline), fill=LINE, width=6)
    for pct, label in [(22, "2025-01"), (48, "Move"), (76, "2025-02")]:
        draw.line((canvas.x(pct), timeline - 14, canvas.x(pct), timeline + 14), fill=LINE, width=4)
        draw_multiline(draw, (canvas.x(pct - 6), timeline + 20, canvas.x(pct + 6), timeline + 50), label, size=20, color=MUTED, bold=True)
    old_row = (canvas.x(18), canvas.y(22), canvas.x(74), canvas.y(38))
    new_row = (canvas.x(44), canvas.y(40), canvas.x(90), canvas.y(56))
    card(draw, old_row, "customer_key 101 | region = North | valid_to = 2025-01-31", fill=BLUE_SOFT, accent=BLUE, title_size=23)
    card(draw, new_row, "customer_key 102 | region = Center | valid_from = 2025-02-01", fill=TEAL_SOFT, accent=TEAL, title_size=23)
    arrow(draw, (canvas.x(48), canvas.y(38)), (canvas.x(48), canvas.y(40)), color=CORAL, width=6, head=14)
    pill(draw, (canvas.x(18), canvas.y(8), canvas.x(84), canvas.y(15)), "Type 2 keeps old and new versions so past sales stay historically true", GOLD_SOFT, NAVY)


def diagram_star(canvas: Canvas, annotated: bool = False):
    draw = canvas.draw
    fact = (canvas.x(39), canvas.y(35), canvas.x(61), canvas.y(56))
    dims = {
        "dim_customer": (22, 22, BLUE, BLUE_SOFT),
        "dim_product": (22, 67, TEAL, TEAL_SOFT),
        "dim_date": (78, 22, GOLD, GOLD_SOFT),
        "dim_store": (78, 67, CORAL, CORAL_SOFT),
    }
    card(draw, fact, "sales_fact", "customer_key | product_key | date_key | amount | qty", fill=GRAY_SOFT, accent=NAVY, title_size=30, body_size=18)
    for title, (x, y, accent, soft) in dims.items():
        rect = (canvas.x(x - 10), canvas.y(y - 7), canvas.x(x + 10), canvas.y(y + 7))
        card(draw, rect, title, fill=soft, accent=accent, title_size=22)
        arrow(draw, (canvas.x(x), canvas.y(y + (7 if y < 50 else -7))), (canvas.x(50), canvas.y(45)), color=accent, width=6, head=14)
    if annotated:
        notes = [
            (63, 40, 90, 50, "Fewer joins"),
            (10, 10, 33, 18, "Denormalized dimensions"),
            (10, 82, 34, 90, "Easy for analysts"),
        ]
        for x1, y1, x2, y2, label in notes:
            pill(draw, (canvas.x(x1), canvas.y(y1), canvas.x(x2), canvas.y(y2)), label, WHITE, NAVY)


def diagram_star_vs_snowflake(canvas: Canvas):
    draw = canvas.draw
    left = (canvas.x(6), canvas.y(16), canvas.x(47), canvas.y(84))
    right = (canvas.x(53), canvas.y(16), canvas.x(94), canvas.y(84))
    card(draw, left, "Star schema", fill=BLUE_SOFT, accent=BLUE, title_size=34)
    card(draw, right, "Snowflake schema", fill=TEAL_SOFT, accent=TEAL, title_size=34)
    draw_multiline(draw, (left[0] + 20, left[1] + 68, left[2] - 20, left[1] + 120), "One fact joins denormalized dimensions", size=23, color=INK, bold=False)
    draw_multiline(draw, (right[0] + 20, right[1] + 68, right[2] - 20, right[1] + 120), "Dimensions split into hierarchy tables", size=23, color=INK, bold=False)
    star = (left[0] + 60, left[1] + 150, left[2] - 60, left[2] - 60)
    snow = (right[0] + 50, right[1] + 140, right[2] - 50, right[2] - 30)
    for center, accent, soft, label in [
        ((canvas.x(26), canvas.y(56)), BLUE, WHITE, "fact"),
        ((canvas.x(21), canvas.y(42)), BLUE, WHITE, "cust"),
        ((canvas.x(31), canvas.y(42)), BLUE, WHITE, "date"),
        ((canvas.x(21), canvas.y(70)), BLUE, WHITE, "prod"),
        ((canvas.x(31), canvas.y(70)), BLUE, WHITE, "store"),
    ]:
        rect = (center[0] - 52, center[1] - 26, center[0] + 52, center[1] + 26)
        card(draw, rect, label, fill=soft, accent=accent, title_size=20)
    for pt in [(canvas.x(21), canvas.y(42)), (canvas.x(31), canvas.y(42)), (canvas.x(21), canvas.y(70)), (canvas.x(31), canvas.y(70))]:
        arrow(draw, pt, (canvas.x(26), canvas.y(56)), color=BLUE, width=5, head=10)
    snow_nodes = {
        "fact": (73, 57, NAVY, WHITE),
        "product": (65, 42, TEAL, WHITE),
        "brand": (83, 42, TEAL, WHITE),
        "category": (89, 30, TEAL, WHITE),
        "date": (65, 72, GOLD, WHITE),
    }
    for label, (x, y, accent, soft) in snow_nodes.items():
        rect = (canvas.x(x) - 56, canvas.y(y) - 26, canvas.x(x) + 56, canvas.y(y) + 26)
        card(draw, rect, label, fill=soft, accent=accent, title_size=20)
    arrow(draw, (canvas.x(65), canvas.y(42)), (canvas.x(73), canvas.y(57)), color=TEAL, width=5, head=10)
    arrow(draw, (canvas.x(83), canvas.y(42)), (canvas.x(73), canvas.y(57)), color=TEAL, width=5, head=10)
    arrow(draw, (canvas.x(89), canvas.y(30)), (canvas.x(83), canvas.y(42)), color=TEAL, width=5, head=10)
    arrow(draw, (canvas.x(65), canvas.y(72)), (canvas.x(73), canvas.y(57)), color=GOLD, width=5, head=10)


def diagram_snowflake_hierarchy(canvas: Canvas):
    draw = canvas.draw
    nodes = [
        ("Electronics", 16, 48, BLUE, BLUE_SOFT),
        ("Apple", 38, 48, TEAL, TEAL_SOFT),
        ("iPhone", 60, 48, GOLD, GOLD_SOFT),
        ("iPhone 15 Pro", 82, 48, CORAL, CORAL_SOFT),
    ]
    for title, x, y, accent, soft in nodes:
        rect = (canvas.x(x - 9), canvas.y(y - 6), canvas.x(x + 9), canvas.y(y + 6))
        card(draw, rect, title, fill=soft, accent=accent, title_size=24)
    for i in range(len(nodes) - 1):
        arrow(draw, (canvas.x(nodes[i][1] + 9), canvas.y(48)), (canvas.x(nodes[i + 1][1] - 9), canvas.y(48)), color=LINE, width=6, head=12)
    pill(draw, (canvas.x(21), canvas.y(20), canvas.x(78), canvas.y(28)), "Category -> Brand -> Product Line -> SKU", GRAY_SOFT, NAVY)
    card(draw, (canvas.x(25), canvas.y(66), canvas.x(75), canvas.y(83)), "Snowflaking removes repeated hierarchy attributes from the base product dimension.", fill=WHITE, accent=TEAL, title_size=26, body_size=20)


def diagram_snowflaking(canvas: Canvas):
    draw = canvas.draw
    source = (canvas.x(8), canvas.y(34), canvas.x(33), canvas.y(60))
    brand = (canvas.x(48), canvas.y(24), canvas.x(67), canvas.y(40))
    category = (canvas.x(48), canvas.y(54), canvas.x(67), canvas.y(70))
    product = (canvas.x(75), canvas.y(39), canvas.x(93), canvas.y(57))
    card(draw, source, "Original product dimension", "product_key | product_name | brand_name | category_name", fill=BLUE_SOFT, accent=BLUE, title_size=28, body_size=19)
    card(draw, brand, "brand_dim", "brand_key | brand_name", fill=TEAL_SOFT, accent=TEAL, title_size=26, body_size=19)
    card(draw, category, "category_dim", "category_key | category_name", fill=GOLD_SOFT, accent=GOLD, title_size=24, body_size=19)
    card(draw, product, "product_dim", "product_key | brand_key | category_key", fill=WHITE, accent=NAVY, title_size=24, body_size=18)
    arrow(draw, (source[2], canvas.y(47)), (brand[0] - 8, canvas.y(32)), color=BLUE, width=6, head=14)
    arrow(draw, (source[2], canvas.y(47)), (category[0] - 8, canvas.y(62)), color=BLUE, width=6, head=14)
    arrow(draw, (brand[2], canvas.y(32)), (product[0] - 10, canvas.y(43)), color=TEAL, width=6, head=14)
    arrow(draw, (category[2], canvas.y(62)), (product[0] - 10, canvas.y(53)), color=GOLD, width=6, head=14)


def diagram_partition_pruning(canvas: Canvas):
    draw = canvas.draw
    fact = (canvas.x(8), canvas.y(28), canvas.x(36), canvas.y(68))
    card(draw, fact, "sales_fact = 1 TB", "365 daily partitions", fill=BLUE_SOFT, accent=BLUE, title_size=34, body_size=24)
    for i in range(10):
        fill = GOLD if i == 6 else "#C7D5E7"
        draw.rectangle((canvas.x(13 + i * 1.7), canvas.y(49), canvas.x(14.2 + i * 1.7), canvas.y(61)), fill=fill)
    focus = (canvas.x(48), canvas.y(30), canvas.x(70), canvas.y(60))
    card(draw, focus, "Filter = one day", "Reads ~1/365 of the table", fill=GOLD_SOFT, accent=GOLD, title_size=32, body_size=24)
    arrow(draw, (fact[2], canvas.y(48)), (focus[0] - 10, canvas.y(45)), color=GOLD, width=8, head=18)
    result = (canvas.x(76), canvas.y(28), canvas.x(93), canvas.y(60))
    card(draw, result, "~2.74 GB", "Instead of 1 TB", fill=GREEN_SOFT, accent=GREEN, title_size=34, body_size=24)
    note = (canvas.x(47), canvas.y(70), canvas.x(93), canvas.y(84))
    card(draw, note, "Cost and runtime fall because the engine skips untouched partitions.", fill=WHITE, accent=TEAL, title_size=25, body_size=20)


def diagram_join_cost(canvas: Canvas):
    draw = canvas.draw
    fact = (canvas.x(8), canvas.y(30), canvas.x(30), canvas.y(58))
    small = (canvas.x(38), canvas.y(22), canvas.x(56), canvas.y(40))
    large = (canvas.x(38), canvas.y(50), canvas.x(56), canvas.y(68))
    planner = (canvas.x(66), canvas.y(30), canvas.x(92), canvas.y(60))
    card(draw, fact, "Pruned fact", "~2.74 GB", fill=BLUE_SOFT, accent=BLUE, title_size=30, body_size=24)
    card(draw, small, "Small dims", "Broadcast join", fill=GREEN_SOFT, accent=GREEN, title_size=26, body_size=20)
    card(draw, large, "Large dims", "Shuffle risk", fill=CORAL_SOFT, accent=CORAL, title_size=26, body_size=20)
    card(draw, planner, "Engine outcome", "Keep dimensions compact so joins stay cheap.", fill=WHITE, accent=NAVY, title_size=30, body_size=22)
    arrow(draw, (fact[2], canvas.y(44)), (planner[0] - 16, canvas.y(44)), color=BLUE, width=8, head=18)
    arrow(draw, (small[2], canvas.y(31)), (planner[0] - 16, canvas.y(37)), color=GREEN, width=6, head=14)
    arrow(draw, (large[2], canvas.y(59)), (planner[0] - 16, canvas.y(51)), color=CORAL, width=6, head=14)


def diagram_query_flow(canvas: Canvas):
    draw = canvas.draw
    steps = [
        ("BI query", CORAL, CORAL_SOFT),
        ("Planner checks\npartition filter", BLUE, BLUE_SOFT),
        ("Prune fact\npartitions", GOLD, GOLD_SOFT),
        ("Join conformed\ndimensions", TEAL, TEAL_SOFT),
        ("Aggregate\nresult", GREEN, GREEN_SOFT),
    ]
    xs = [6, 24, 42, 60, 78]
    for x, (title, accent, soft) in zip(xs, steps):
        rect = (canvas.x(x), canvas.y(38), canvas.x(x + 14), canvas.y(58))
        card(draw, rect, title, fill=soft, accent=accent, title_size=22)
    for i in range(len(xs) - 1):
        arrow(draw, (canvas.x(xs[i] + 14), canvas.y(48)), (canvas.x(xs[i + 1]) - 8, canvas.y(48)), color=LINE, width=6, head=12)
    pill(draw, (canvas.x(29), canvas.y(18), canvas.x(71), canvas.y(26)), "Most cost is decided before the final aggregation", GRAY_SOFT, NAVY)


def diagram_query_antipattern(canvas: Canvas):
    draw = canvas.draw
    good = (canvas.x(8), canvas.y(22), canvas.x(42), canvas.y(78))
    bad = (canvas.x(58), canvas.y(22), canvas.x(92), canvas.y(78))
    card(draw, good, "Good query", fill=GREEN_SOFT, accent=GREEN, title_size=30)
    draw_list(draw, good[0] + 24, good[1] + 86, good[2] - good[0] - 48, [
        "Date predicate present",
        "Partitions pruned",
        "Predictable cost",
    ], size=24)
    card(draw, bad, "Anti-pattern", fill=CORAL_SOFT, accent=CORAL, title_size=30)
    draw_list(draw, bad[0] + 24, bad[1] + 86, bad[2] - bad[0] - 48, [
        "Missing time filter",
        "Full table scan",
        "Slots saturate under concurrency",
    ], size=24)
    arrow(draw, (good[2], canvas.y(50)), (bad[0] - 12, canvas.y(50)), color=LINE, width=5, head=12)
    pill(draw, (canvas.x(32), canvas.y(84), canvas.x(68), canvas.y(91)), "Guardrails should reject the anti-pattern automatically", GOLD_SOFT, NAVY)


def diagram_architecture_evolution(canvas: Canvas):
    draw = canvas.draw
    left = (canvas.x(8), canvas.y(20), canvas.x(42), canvas.y(82))
    right = (canvas.x(58), canvas.y(20), canvas.x(92), canvas.y(82))
    card(draw, left, "v1", fill=CORAL_SOFT, accent=CORAL, title_size=34)
    draw_list(draw, left[0] + 24, left[1] + 96, left[2] - left[0] - 48, [
        "One fact table",
        "Manual dimensions",
        "Few protections",
    ], size=24)
    card(draw, right, "v2", fill=TEAL_SOFT, accent=TEAL, title_size=34)
    draw_list(draw, right[0] + 24, right[1] + 96, right[2] - right[0] - 48, [
        "Conformed dimensions",
        "Partition-aware querying",
        "SQL guardrails and certified marts",
    ], size=24)
    arrow(draw, (left[2], canvas.y(50)), (right[0] - 14, canvas.y(50)), color=GOLD, width=9, head=18)
    pill(draw, (canvas.x(36), canvas.y(12), canvas.x(64), canvas.y(19)), "Planned migration beats broken dashboards", GRAY_SOFT, NAVY)


def diagram_governance_controls(canvas: Canvas):
    draw = canvas.draw
    layers = [
        ("Request", CORAL, CORAL_SOFT, 8),
        ("Metric layer", BLUE, BLUE_SOFT, 29),
        ("Guardrails", GOLD, GOLD_SOFT, 50),
        ("Policy", TEAL, TEAL_SOFT, 71),
    ]
    prev = None
    for title, accent, soft, x in layers:
        rect = (canvas.x(x), canvas.y(34), canvas.x(x + 17), canvas.y(56))
        card(draw, rect, title, fill=soft, accent=accent, title_size=28)
        if prev:
            arrow(draw, (prev[2], canvas.y(45)), (rect[0] - 10, canvas.y(45)), color=LINE, width=6, head=12)
        prev = rect
    note = (canvas.x(15), canvas.y(66), canvas.x(88), canvas.y(84))
    card(draw, note, "Require date filters | cap bytes scanned | route to approved metrics", fill=WHITE, accent=NAVY, title_size=26, body_size=20)


def build_diagram(key: str, size=(1600, 900)) -> Image.Image:
    canvas = Canvas(*size)
    mapping = {
        "dw_vs_lake": diagram_dw_vs_lake,
        "dw_choice": diagram_dw_choice,
        "why_dw": diagram_why_dw,
        "schema_strategy": diagram_schema_strategy,
        "hybrid_arch": diagram_hybrid_architecture,
        "extract_sources": diagram_extraction_sources,
        "safe_writes": diagram_safe_writes,
        "elt": diagram_elt,
        "cdc_options": diagram_cdc_options,
        "cdc_example": diagram_cdc_example,
        "sttm_flow": diagram_sttm_flow,
        "idempotency": diagram_idempotency,
        "merge_vs_overwrite": diagram_merge_vs_overwrite,
        "multidimensional": diagram_multidimensional_model,
        "location_hierarchy": diagram_location_hierarchy,
        "scd_sequence": diagram_scd_sequence,
        "star_default": lambda c: diagram_star(c, annotated=False),
        "star_annotated": lambda c: diagram_star(c, annotated=True),
        "star_vs_snowflake": diagram_star_vs_snowflake,
        "snowflake_hierarchy": diagram_snowflake_hierarchy,
        "snowflaking": diagram_snowflaking,
        "partition_pruning": diagram_partition_pruning,
        "join_cost": diagram_join_cost,
        "query_flow": diagram_query_flow,
        "query_antipattern": diagram_query_antipattern,
        "architecture_evolution": diagram_architecture_evolution,
        "governance_controls": diagram_governance_controls,
    }
    mapping[key](canvas)
    return canvas.image


def remove_shape(slide, shape):
    slide.shapes._spTree.remove(shape._element)


def replace_picture(slide, image_path: Path):
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pictures:
        raise RuntimeError("No picture found to replace.")
    picture = pictures[0]
    left, top, width, height = picture.left, picture.top, picture.width, picture.height
    remove_shape(slide, picture)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def replace_body_with_picture(slide, image_path: Path):
    title_shape = slide.shapes.title
    for shape in list(slide.shapes):
        if shape == title_shape:
            continue
        if getattr(shape, "is_placeholder", False) and ("Footer" in shape.name or "Slide Number" in shape.name):
            continue
        remove_shape(slide, shape)
    slide.shapes.add_picture(str(image_path), 731520, 1463040, 9723120, 4389120)


def render_assets(deck_name: str, updates: list[SlideUpdate]) -> dict[int, Path]:
    rendered: dict[int, Path] = {}
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)
    for update in updates:
        out = ASSETS_DIR / f"{Path(deck_name).stem}-slide-{update.slide_no}.png"
        image = build_diagram(update.key)
        image.save(out)
        rendered[update.slide_no] = out
    return rendered


def update_presentation(deck_name: str, updates: list[SlideUpdate]):
    rendered = render_assets(deck_name, updates)
    deck_path = SLIDES_DIR / deck_name
    presentation = Presentation(str(deck_path))
    for update in updates:
        slide = presentation.slides[update.slide_no - 1]
        if update.mode == "replace_body":
            replace_body_with_picture(slide, rendered[update.slide_no])
        else:
            replace_picture(slide, rendered[update.slide_no])
    presentation.save(str(deck_path))


DECKS = {
    "04-DWH-1-refined.pptx": [
        SlideUpdate(4, "dw_vs_lake"),
        SlideUpdate(5, "dw_choice"),
        SlideUpdate(7, "why_dw"),
        SlideUpdate(9, "schema_strategy"),
        SlideUpdate(10, "hybrid_arch"),
        SlideUpdate(12, "extract_sources"),
        SlideUpdate(14, "safe_writes"),
        SlideUpdate(15, "elt"),
        SlideUpdate(17, "cdc_options"),
        SlideUpdate(18, "cdc_example"),
        SlideUpdate(22, "sttm_flow"),
        SlideUpdate(24, "idempotency"),
        SlideUpdate(27, "merge_vs_overwrite"),
    ],
    "05-DWH-2-refined.pptx": [
        SlideUpdate(6, "multidimensional", mode="replace_body"),
        SlideUpdate(8, "location_hierarchy", mode="replace_body"),
        SlideUpdate(14, "scd_sequence"),
        SlideUpdate(15, "star_annotated"),
        SlideUpdate(16, "star_default"),
        SlideUpdate(17, "star_default"),
        SlideUpdate(19, "star_vs_snowflake"),
        SlideUpdate(20, "snowflake_hierarchy"),
        SlideUpdate(21, "snowflaking", mode="replace_body"),
        SlideUpdate(24, "partition_pruning"),
        SlideUpdate(25, "join_cost"),
        SlideUpdate(26, "query_flow"),
        SlideUpdate(29, "query_antipattern"),
        SlideUpdate(30, "architecture_evolution"),
        SlideUpdate(31, "governance_controls"),
    ],
}


def main():
    for deck_name, updates in DECKS.items():
        update_presentation(deck_name, updates)
        print(f"Updated {deck_name}")


if __name__ == "__main__":
    main()
