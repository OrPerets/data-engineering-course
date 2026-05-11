from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = "build/DWH-Schemas-enhanced.pptx"

W, H = Inches(13.333), Inches(7.5)
BG = "F6F7F2"
INK = "202428"
MUTED = "667085"
CHARCOAL = "253238"
TEAL = "0F766E"
TEAL_DARK = "115E59"
MINT = "D8F3EA"
AMBER = "D97706"
AMBER_SOFT = "FDECC8"
RED = "B42318"
RED_SOFT = "FDE4E1"
BLUE = "335C67"
BLUE_SOFT = "DDE8EA"
LINE = "D6DAD2"
WHITE = "FFFFFF"


def rgb(hex_color):
    hex_color = hex_color.replace("#", "")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_shape(slide, shape_type, x, y, w, h, fill, line=None, radius=True):
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.5, arrow=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    if arrow:
        conn.line.end_arrowhead = True
    return conn


def textbox(slide, text, x, y, w, h, size=18, bold=False, color=INK, align="left",
            font="Aptos", valign="top", margin=0.06):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "mid": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    for run in p.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def paragraph_box(slide, lines, x, y, w, h, size=15, color=INK, leading=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.space_after = Pt(4)
        p.line_spacing = leading
        for r in p.runs:
            r.font.name = "Aptos"
            r.font.size = Pt(size)
            r.font.color.rgb = rgb(color)
    return box


def title(slide, text, kicker=None, dark=False):
    color = WHITE if dark else INK
    if kicker:
        textbox(slide, kicker.upper(), 0.7, 0.34, 4.3, 0.25, 9, True, AMBER if dark else TEAL)
    textbox(slide, text, 0.66, 0.62, 8.9, 0.62, 26 if len(text) > 42 else 30, True, color, font="Aptos Display")


def footer(slide, n, dark=False):
    color = "C9D0CA" if dark else MUTED
    textbox(slide, "Or Peretz", 0.68, 7.02, 1.3, 0.22, 8, False, color)
    textbox(slide, str(n), 12.15, 7.02, 0.5, 0.22, 8, False, color, align="right")


def background(slide, dark=False):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 7.5, CHARCOAL if dark else BG)
    if not dark:
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, 13.333, 0.12, TEAL)


def section_chip(slide, text, x, y, fill=TEAL, color=WHITE):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 1.08, 0.32, fill)
    textbox(slide, text, x, y + 0.04, 1.08, 0.2, 8, True, color, align="center", margin=0.02)


def card(slide, x, y, w, h, head, body, accent=TEAL, fill=WHITE, head_size=14, body_size=12):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill, LINE)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, 0.08, h, accent)
    textbox(slide, head, x + 0.18, y + 0.14, w - 0.3, 0.28, head_size, True, INK)
    paragraph_box(slide, body if isinstance(body, list) else [body], x + 0.18, y + 0.52, w - 0.32, h - 0.62, body_size, MUTED)


def small_label(slide, text, x, y, w, fill=BLUE_SOFT, color=INK):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, 0.3, fill)
    textbox(slide, text, x, y + 0.055, w, 0.15, 8, True, color, align="center", margin=0.02)


def callout(slide, text, x, y, w, h=0.58, fill=MINT, color=TEAL_DARK, size=12):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill, "B8D8D0")
    textbox(slide, text, x + 0.12, y + 0.17, w - 0.24, h - 0.25, size, True, color, align="center", margin=0.01)


def db(slide, label, x, y, w=1.15, h=0.68, fill=BLUE_SOFT, line=BLUE, text_color=INK):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.CAN, x, y, w, h, fill, line)
    textbox(slide, label, x + 0.06, y + 0.22, w - 0.12, 0.2, 8.5, True, text_color, align="center", margin=0.01)


def arrow_box(slide, text, x, y, w, h, fill, color=WHITE):
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, x, y, w, h, fill)
    textbox(slide, text, x + 0.05, y + 0.16, w - 0.22, h - 0.3, 11, True, color, align="center", margin=0.02)


def add_table(slide, x, y, w, h, headers, rows, col_fills=None, font_size=11):
    cols = len(headers)
    col_w = w / cols
    row_h = h / (len(rows) + 1)
    for c, head in enumerate(headers):
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + c * col_w, y, col_w, row_h, CHARCOAL)
        textbox(slide, head, x + c * col_w + 0.04, y + 0.12, col_w - 0.08, row_h - 0.2, font_size, True, WHITE, align="center", margin=0.02)
    for r, row in enumerate(rows):
        for c, txt in enumerate(row):
            fill = col_fills[c] if col_fills else WHITE
            add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x + c * col_w, y + (r + 1) * row_h, col_w, row_h, fill, LINE)
            textbox(slide, txt, x + c * col_w + 0.05, y + (r + 1) * row_h + 0.08, col_w - 0.1, row_h - 0.1, font_size, c == 0, INK, align="center", margin=0.01)


def add_relation(slide, x1, y1, x2, y2, label=None):
    add_line(slide, x1, y1, x2, y2, TEAL, 2.0)
    if label:
        textbox(slide, label, (x1 + x2) / 2 - 0.3, (y1 + y2) / 2 - 0.12, 0.6, 0.2, 8, True, TEAL, align="center", margin=0)


prs = Presentation()
prs.slide_width = W
prs.slide_height = H
blank = prs.slide_layouts[6]


def new_slide(dark=False):
    s = prs.slides.add_slide(blank)
    background(s, dark)
    return s


# 1
s = new_slide(True)
title(s, "DWH Schema Design", "Data Engineering - Part 2", True)
textbox(s, "From raw events to reliable metrics, fast dashboards, and auditable history.", 0.7, 1.45, 7.6, 0.55, 19, False, "DDE8EA")
for i, lab in enumerate(["CRM", "Orders", "POS", "Web"]):
    db(s, lab, 1.0 + i * 1.35, 3.4, fill="334E52", line="7FB7AE", text_color=WHITE)
arrow_box(s, "STAGING", 6.45, 3.42, 1.35, 0.62, TEAL)
arrow_box(s, "CORE DWH", 8.05, 3.42, 1.55, 0.62, AMBER)
arrow_box(s, "MARTS", 9.88, 3.42, 1.22, 0.62, BLUE)
db(s, "BI / ML", 11.45, 3.37, 1.05, 0.72, fill="D8F3EA", line="7FB7AE")
textbox(s, "Professional goal: model decisions that survive scale, change, and audit.", 1.02, 5.65, 10.4, 0.46, 18, True, WHITE)
footer(s, 1, True)

# 2
s = new_slide()
title(s, "Learning Outcomes", "What students should be able to do")
items = [
    ("Declare grain", ["Define one row precisely", "Prevent double counting"]),
    ("Choose schema", ["Use star by default", "Snowflake only for clear tradeoffs"]),
    ("Engineer history", ["Select SCD type", "Load facts with correct keys"]),
    ("Operate the model", ["Partition, cluster, test", "Monitor scan and drift"]),
]
for i, (h, b) in enumerate(items):
    card(s, 0.8 + (i % 2) * 5.9, 1.62 + (i // 2) * 2.05, 5.25, 1.5, h, b, [TEAL, AMBER, BLUE, RED][i], WHITE, 17, 13)
footer(s, 2)

# 3
s = new_slide()
title(s, "The Engineering Problem", "Why schema design matters")
arrow_box(s, "Raw source records", 0.85, 2.2, 2.0, 0.75, BLUE)
arrow_box(s, "Modeled warehouse", 3.25, 2.2, 2.15, 0.75, TEAL)
arrow_box(s, "Certified metrics", 5.85, 2.2, 2.0, 0.75, AMBER)
arrow_box(s, "Business decisions", 8.25, 2.2, 2.15, 0.75, CHARCOAL)
for x in [2.85, 5.42, 7.88]:
    add_line(s, x, 2.58, x + 0.35, 2.58, MUTED, 2.0)
card(s, 0.92, 4.18, 2.85, 1.58, "Risk", ["Wrong grain turns one sale into many", "Natural-key drift breaks joins"], RED, WHITE, 17, 12)
card(s, 4.05, 4.18, 2.85, 1.58, "Cost", ["Missing predicates scan full facts", "Tiny files and skew waste slots"], AMBER, WHITE, 17, 12)
card(s, 7.18, 4.18, 2.85, 1.58, "Trust", ["Metric logic diverges in dashboards", "History is overwritten accidentally"], TEAL, WHITE, 17, 12)
card(s, 10.31, 4.18, 2.1, 1.58, "Goal", ["Stable, testable contracts"], BLUE, WHITE, 17, 12)
footer(s, 3)

# 4
s = new_slide()
title(s, "Modern DWH Layers", "A practical mental model")
layers = [
    ("Sources", "Operational truth; messy keys and formats", BLUE),
    ("Raw / staging", "Append, preserve, type lightly", TEAL),
    ("Core DWH", "Conformed dimensions and atomic facts", AMBER),
    ("Marts", "Use-case tables with certified semantics", RED),
    ("Consumers", "BI, reverse ETL, ML features, APIs", CHARCOAL),
]
for i, (h, b, c) in enumerate(layers):
    x = 0.65 + i * 2.45
    add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, 2.0, 2.05, 1.75, WHITE, LINE)
    add_shape(s, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, 2.0, 2.05, 0.16, c)
    textbox(s, h, x + 0.12, 2.35, 1.78, 0.28, 15, True, INK, align="center")
    paragraph_box(s, [b], x + 0.15, 2.78, 1.75, 0.62, 10.5, MUTED)
    if i < len(layers) - 1:
        add_line(s, x + 2.08, 2.88, x + 2.36, 2.88, c, 2)
callout(s, "Schema design lives mainly in Core DWH and marts, but it must respect source behavior and consumer workload.", 1.55, 5.28, 10.2, 0.56, MINT, TEAL_DARK, 11)
footer(s, 4)

# 5
s = new_slide()
title(s, "Fact and Dimension Contracts", "The two core table roles")
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.0, 2.05, 3.8, 2.4, WHITE, LINE)
textbox(s, "Fact table", 1.25, 2.28, 3.2, 0.35, 21, True, TEAL)
paragraph_box(s, ["- Atomic business event", "- Declared grain", "- Foreign keys to dimensions", "- Numeric measures"], 1.25, 2.85, 3.1, 1.15, 14)
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 8.55, 2.05, 3.8, 2.4, WHITE, LINE)
textbox(s, "Dimension table", 8.8, 2.28, 3.2, 0.35, 21, True, AMBER)
paragraph_box(s, ["- Descriptive context", "- Attributes and hierarchies", "- Surrogate key", "- Business-readable labels"], 8.8, 2.85, 3.1, 1.15, 14)
arrow_box(s, "JOIN BY SURROGATE KEY", 5.25, 2.85, 2.65, 0.7, CHARCOAL)
add_table(s, 2.0, 5.23, 9.4, 0.88, ["Question", "Fact answer", "Dimension answer"], [["What happened?", "Order line sold", "Product, customer, date"]], [BLUE_SOFT, WHITE, WHITE], 10)
footer(s, 5)

# 6
s = new_slide()
title(s, "Declare the Grain First", "The non-negotiable design move")
textbox(s, "Grain = what one row means.", 0.82, 1.42, 4.8, 0.35, 20, True, TEAL)
card(s, 0.82, 2.0, 3.55, 1.22, "Good grain", ["One row per order line", "Keys: order_id + line_id"], TEAL)
card(s, 0.82, 3.55, 3.55, 1.22, "Ambiguous grain", ["One row per order", "But product columns repeat"], RED)
card(s, 0.82, 5.1, 3.55, 1.22, "Engineering test", ["Can a unique key be enforced?", "Can every measure aggregate safely?"], AMBER)
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.3, 1.65, 6.8, 4.9, WHITE, LINE)
textbox(s, "Grain decision record", 5.62, 1.96, 4.5, 0.35, 19, True, INK)
add_table(s, 5.62, 2.62, 5.95, 2.75, ["Field", "Example"], [
    ["Fact name", "fact_order_line"],
    ["One row", "A sold product line in one order"],
    ["Unique key", "source_order_id + source_line_id"],
    ["Valid measures", "gross_amount, discount_amount, qty"],
    ["Invalid measure", "customer_lifetime_value"],
], [BLUE_SOFT, WHITE], 10)
footer(s, 6)

# 7
s = new_slide()
title(s, "Fact Table Types", "Different grains for different questions")
types = [
    ("Transaction fact", "One row per event", "sales, clicks, payments", TEAL),
    ("Periodic snapshot", "One row per period", "daily inventory, account balance", AMBER),
    ("Accumulating snapshot", "One row per lifecycle", "order fulfillment, ticket pipeline", BLUE),
]
for i, (h, b, ex, c) in enumerate(types):
    x = 0.85 + i * 4.05
    card(s, x, 1.65, 3.45, 3.8, h, [b, "Use for: " + ex], c, WHITE, 16, 12)
    y = 4.37
    for j in range(4):
        add_shape(s, MSO_AUTO_SHAPE_TYPE.OVAL, x + 0.55 + j * 0.58, y, 0.22, 0.22, c)
        if j < 3:
            add_line(s, x + 0.77 + j * 0.58, y + 0.11, x + 1.12 + j * 0.58, y + 0.11, c, 1.5)
callout(s, "Pick the fact type before choosing partitions, ETL strategy, or dashboard logic.", 2.12, 5.98, 9.1, 0.55, AMBER_SOFT, AMBER, 11)
footer(s, 7)

# 8
s = new_slide()
title(s, "Measure Semantics", "Aggregation rules are part of the schema")
add_table(s, 0.85, 1.62, 11.65, 3.8, ["Measure type", "Examples", "Safe aggregation", "Common failure"], [
    ["Additive", "revenue, qty", "SUM across all dimensions", "Currency conversion mixed late"],
    ["Semi-additive", "inventory_balance", "SUM by product, not by time", "Monthly total sums daily stock"],
    ["Non-additive", "conversion_rate", "Compute from numerator / denominator", "Averages of averages"],
    ["Derived metric", "gross_margin_pct", "Certified formula in semantic layer", "Different BI tools disagree"],
], [BLUE_SOFT, WHITE, WHITE, RED_SOFT], 10)
card(s, 2.0, 5.68, 9.2, 0.95, "Data engineer rule", ["Store atomic components when possible; calculate ratios close to the semantic layer."], TEAL, WHITE, 13, 11)
footer(s, 8)

# 9
s = new_slide()
title(s, "Dimension Strategy", "Attributes, hierarchies, and conformance")
db(s, "dim_product", 1.05, 2.45, 1.35, 0.82, fill=MINT, line=TEAL)
for i, lab in enumerate(["SKU", "Category", "Brand", "Department"]):
    lx = 3.05 + i * 1.42
    ly = 1.35 + i * 0.66
    add_line(s, 2.4, 2.86, lx, ly + 0.22, "9AB8B0", 1.0)
for i, lab in enumerate(["SKU", "Category", "Brand", "Department"]):
    lx = 3.05 + i * 1.42
    ly = 1.35 + i * 0.66
    add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, lx, ly, 1.18, 0.42, WHITE, LINE)
    textbox(s, lab, lx + 0.02, ly + 0.11, 1.14, 0.16, 8.5, True, INK, align="center")
card(s, 8.5, 1.55, 3.45, 1.05, "Conformed dimension", ["Same product_key and attributes reused by sales, returns, inventory."], TEAL)
card(s, 8.5, 2.95, 3.45, 1.05, "Role-playing dimension", ["Same dim_date used as order_date, ship_date, invoice_date."], AMBER)
card(s, 8.5, 4.35, 3.45, 1.05, "Degenerate dimension", ["Business ID stored in fact, e.g., order_number."], BLUE)
footer(s, 9)

# 10
s = new_slide()
title(s, "Keys: Natural, Surrogate, Durable", "Keep source volatility out of BI joins")
for i, (lab, desc, c) in enumerate([
    ("Natural key", "customer_id from CRM", BLUE),
    ("Durable key", "stable customer identity across systems", AMBER),
    ("Surrogate key", "warehouse integer used by facts", TEAL),
]):
    card(s, 0.92 + i * 4.15, 1.65, 3.4, 1.65, lab, [desc], c, WHITE, 16, 12)
    if i < 2:
        add_line(s, 4.35 + i * 4.15, 2.46, 4.86 + i * 4.15, 2.46, MUTED, 2)
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.25, 4.45, 10.75, 1.45, RED_SOFT, "EAA8A2")
textbox(s, "Anti-pattern", 1.55, 4.72, 1.7, 0.28, 15, True, RED)
paragraph_box(s, ["Joining BI facts directly on natural keys makes source merges, key reuse, and late-arriving dimensions visible to every dashboard."], 3.35, 4.68, 7.8, 0.52, 13, INK)
footer(s, 10)

# 11
s = new_slide()
title(s, "Slowly Changing Dimensions", "Pick history behavior intentionally")
add_table(s, 0.8, 1.62, 11.75, 3.95, ["Type", "Behavior", "Use when", "Risk"], [
    ["Type 0", "Never change", "Immutable attributes", "Bad source corrections ignored"],
    ["Type 1", "Overwrite value", "Fix errors, no history needed", "Past metrics reclassify"],
    ["Type 2", "Insert new version", "Historical truth matters", "ETL complexity and wider joins"],
    ["Type 3", "Store limited prior value", "Simple previous/current comparison", "Limited audit history"],
], [BLUE_SOFT, WHITE, WHITE, RED_SOFT], 10.5)
small_label(s, "Finance, compliance, and customer segmentation usually need Type 2 for material attributes.", 1.75, 5.95, 9.8, fill=MINT, color=TEAL_DARK)
footer(s, 11)

# 12
s = new_slide()
title(s, "SCD2 Customer Region Change", "Facts must point to the right dimension version")
add_line(s, 1.05, 3.1, 11.8, 3.1, CHARCOAL, 2)
for x, label, c in [(1.3, "2025-01-01\nRegion=North", TEAL), (6.0, "2025-06-10\nRegion=Center", AMBER), (10.35, "current", BLUE)]:
    add_shape(s, MSO_AUTO_SHAPE_TYPE.OVAL, x, 2.86, 0.38, 0.38, c)
    textbox(s, label, x - 0.5, 3.38, 1.45, 0.55, 10, True, INK, align="center")
card(s, 1.15, 1.35, 3.2, 0.9, "dim_customer v1", ["customer_key=101, valid_to=2025-06-09"], TEAL, WHITE, 13, 10)
card(s, 5.35, 1.35, 3.35, 0.9, "dim_customer v2", ["customer_key=219, valid_from=2025-06-10"], AMBER, WHITE, 13, 10)
card(s, 2.0, 5.0, 8.9, 0.85, "ETL join condition", ["Fact sale_date must be between dimension valid_from and valid_to."], BLUE, WHITE, 13, 11)
footer(s, 12)

# 13
s = new_slide()
title(s, "Star Schema: Default Design", "Simple joins, fast BI, clear ownership")
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 5.45, 2.55, 2.45, 1.08, CHARCOAL)
textbox(s, "fact_sales\norder_line grain", 5.55, 2.78, 2.25, 0.4, 12, True, WHITE, align="center")
dims = [
    ("dim_date", 2.0, 1.25, TEAL),
    ("dim_customer", 8.95, 1.25, TEAL),
    ("dim_product", 2.0, 4.65, AMBER),
    ("dim_store", 8.95, 4.65, AMBER),
]
for lab, x, y, c in dims:
    add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 2.05, 0.78, WHITE, c)
    textbox(s, lab, x + 0.1, y + 0.25, 1.85, 0.2, 12, True, INK, align="center")
    add_relation(s, x + 2.05 if x < 5 else x, y + 0.39, 6.68, 3.09)
small_label(s, "Default to star unless normalization has a measurable maintenance or correctness benefit.", 2.0, 6.15, 9.35, fill=AMBER_SOFT, color=AMBER)
footer(s, 13)

# 14
s = new_slide()
title(s, "Star Schema Design Checklist", "Review before building ETL")
checks = [
    ("1", "Grain declared", "One row meaning and unique key are enforceable"),
    ("2", "Measures classified", "Additive, semi-additive, non-additive rules documented"),
    ("3", "Dimensions conformed", "Shared keys and labels across facts"),
    ("4", "History modeled", "SCD type chosen per attribute group"),
    ("5", "Workload aligned", "Partition, clustering, and predicates match queries"),
]
for i, (num, h, b) in enumerate(checks):
    y = 1.38 + i * 0.93
    add_shape(s, MSO_AUTO_SHAPE_TYPE.OVAL, 0.95, y, 0.45, 0.45, TEAL if i < 3 else AMBER)
    textbox(s, num, 1.05, y + 0.105, 0.25, 0.1, 9, True, WHITE, align="center", margin=0)
    textbox(s, h, 1.65, y - 0.02, 3.2, 0.22, 14, True, INK)
    textbox(s, b, 1.65, y + 0.28, 8.8, 0.22, 11, False, MUTED)
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 9.75, 1.95, 2.45, 2.75, MINT, "7FB7AE")
textbox(s, "Model review gate", 10.0, 2.28, 1.95, 0.35, 17, True, TEAL_DARK, align="center")
paragraph_box(s, ["No production fact table should be created without this review."], 10.02, 3.05, 1.92, 0.8, 12, INK)
footer(s, 14)

# 15
s = new_slide()
title(s, "Snowflake Schema", "Normalize selected dimensions when it earns its cost")
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.9, 1.55, 5.4, 4.5, WHITE, LINE)
textbox(s, "Star product dimension", 1.2, 1.88, 4.0, 0.28, 15, True, TEAL)
db(s, "dim_product\nSKU+brand+category", 2.45, 3.25, 2.25, 0.9, fill=MINT, line=TEAL)
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 7.0, 1.55, 5.4, 4.5, WHITE, LINE)
textbox(s, "Snowflaked product hierarchy", 7.3, 1.88, 4.0, 0.28, 15, True, AMBER)
db(s, "dim_product", 7.58, 3.25, 1.5, 0.8, fill=AMBER_SOFT, line=AMBER)
db(s, "dim_brand", 9.35, 2.25, 1.38, 0.72, fill=BLUE_SOFT, line=BLUE)
db(s, "dim_category", 10.58, 3.85, 1.5, 0.72, fill=BLUE_SOFT, line=BLUE)
add_line(s, 9.05, 3.45, 9.42, 2.62, AMBER, 1.5)
add_line(s, 9.05, 3.7, 10.6, 4.18, AMBER, 1.5)
callout(s, "Snowflake reduces redundancy, but every extra join is a performance and usability tax.", 2.05, 6.27, 9.25, 0.55, RED_SOFT, RED, 11)
footer(s, 15)

# 16
s = new_slide()
title(s, "When Snowflaking Helps or Hurts", "A practical tradeoff")
add_line(s, 6.65, 1.45, 6.65, 6.3, LINE, 1.5)
add_line(s, 1.0, 3.86, 12.35, 3.86, LINE, 1.5)
textbox(s, "Lower query latency", 1.12, 1.55, 2.2, 0.25, 10, True, MUTED)
textbox(s, "Higher maintainability", 9.7, 1.55, 2.2, 0.25, 10, True, MUTED, align="right")
card(s, 1.2, 2.0, 4.55, 1.25, "Keep star", ["Small dimensions, BI-heavy workload, simple hierarchies"], TEAL)
card(s, 7.48, 2.0, 4.0, 1.25, "Snowflake selectively", ["Huge hierarchy reused by many dimensions"], AMBER)
card(s, 1.2, 4.45, 4.55, 1.25, "Avoid", ["More joins only to make ERD look normalized"], RED)
card(s, 7.48, 4.45, 4.0, 1.25, "Consider hybrid", ["Materialized mart hides normalized core"], BLUE)
footer(s, 16)

# 17
s = new_slide()
title(s, "Galaxy Schema and Bus Matrix", "Multiple facts, shared dimensions")
dims = [("dim_date", 1.0, 1.5), ("dim_product", 1.0, 3.0), ("dim_customer", 1.0, 4.5)]
for lab, x, y in dims:
    add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 1.75, 0.58, WHITE, TEAL)
    textbox(s, lab, x + 0.05, y + 0.2, 1.65, 0.15, 9, True, INK, align="center")
facts = [("fact_sales", 4.9, 2.0, TEAL), ("fact_returns", 4.9, 4.2, AMBER), ("fact_inventory", 8.0, 3.05, BLUE)]
for _, x, y in dims:
    for _, fx, fy, _ in facts:
        add_line(s, x + 1.75, y + 0.29, fx, fy + 0.36, LINE, 0.7)
for lab, x, y, c in facts:
    add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, 1.65, 0.72, c)
    textbox(s, lab, x + 0.05, y + 0.25, 1.55, 0.18, 9.5, True, WHITE, align="center")
add_table(s, 9.52, 1.65, 3.12, 3.45, ["Fact", "Date", "Prod.", "Cust."], [
    ["sales", "Y", "Y", "Y"],
    ["returns", "Y", "Y", "Y"],
    ["inv.", "Y", "Y", "N"],
], [BLUE_SOFT, WHITE, WHITE, WHITE], 9.5)
footer(s, 17)

# 18
s = new_slide()
title(s, "Bridge Tables for Many-to-Many", "Do not force dimensions into false one-to-many joins")
db(s, "dim_student", 1.2, 3.0, 1.65, 0.78, fill=MINT, line=TEAL)
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 4.25, 2.55, 2.1, 1.25, CHARCOAL)
textbox(s, "bridge_student_course\nstudent_key\ncourse_key\nweight", 4.35, 2.78, 1.9, 0.58, 8.5, True, WHITE, align="center")
db(s, "dim_course", 7.8, 3.0, 1.55, 0.78, fill=AMBER_SOFT, line=AMBER)
db(s, "fact_exam", 10.55, 3.0, 1.55, 0.78, fill=BLUE_SOFT, line=BLUE)
add_line(s, 2.85, 3.39, 4.25, 3.18, TEAL, 2)
add_line(s, 6.35, 3.18, 7.8, 3.39, AMBER, 2)
add_line(s, 9.35, 3.39, 10.55, 3.39, BLUE, 2)
card(s, 1.25, 5.1, 10.55, 0.86, "Use case", ["Customer belongs to multiple segments, product has multiple tags, student attends multiple courses."], TEAL, WHITE, 13, 11)
footer(s, 18)

# 19
s = new_slide()
title(s, "Incident: Full-Scan Dashboard Outage", "Operational symptoms expose modeling gaps")
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.95, 1.65, 11.45, 3.1, WHITE, LINE)
for i, (t, h, c) in enumerate([
    ("08:00", "Dashboard refresh every 5 minutes", TEAL),
    ("08:12", "No date predicate on 1 TB fact", AMBER),
    ("08:20", "Cluster slots saturated", RED),
    ("08:45", "Guardrail added and query fixed", BLUE),
]):
    x = 1.25 + i * 2.75
    add_shape(s, MSO_AUTO_SHAPE_TYPE.OVAL, x, 2.35, 0.45, 0.45, c)
    textbox(s, t, x - 0.25, 1.95, 0.95, 0.2, 10, True, c, align="center")
    textbox(s, h, x - 0.65, 2.98, 1.65, 0.55, 10.5, False, INK, align="center")
    if i < 3:
        add_line(s, x + 0.45, 2.58, x + 2.25, 2.58, LINE, 2)
card(s, 2.0, 5.35, 9.3, 0.82, "Root cause", ["The model and semantic layer did not enforce partition-filter usage for large facts."], RED, WHITE, 13, 11)
footer(s, 19)

# 20
s = new_slide()
title(s, "Partition and Cluster Strategy", "Physical design follows workload")
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.95, 1.7, 5.35, 4.7, WHITE, LINE)
textbox(s, "fact_sales partitioned by date_key", 1.2, 2.0, 3.9, 0.25, 14, True, TEAL)
for i, c in enumerate([TEAL, TEAL, MINT, MINT, AMBER_SOFT, AMBER_SOFT, BLUE_SOFT]):
    add_shape(s, MSO_AUTO_SHAPE_TYPE.RECTANGLE, 1.25 + i * 0.63, 3.0, 0.5, 1.7, c, WHITE)
textbox(s, "Query: last 7 days", 1.25, 5.05, 1.9, 0.22, 11, True, TEAL)
add_line(s, 1.25, 4.82, 2.38, 4.82, TEAL, 4)
add_table(s, 7.0, 1.7, 5.05, 3.4, ["Choice", "Good default", "Watch out"], [
    ["Partition", "Dominant time filter", "Too many tiny partitions"],
    ["Cluster / sort", "Secondary filters", "Skewed high-cardinality keys"],
    ["File size", "Large enough for scan efficiency", "Small-file explosion"],
], [BLUE_SOFT, WHITE, RED_SOFT], 9.5)
callout(s, "Data engineers own both logical correctness and physical scan behavior.", 3.0, 5.98, 7.45, 0.55, MINT, TEAL_DARK, 11)
footer(s, 20)

# 21
s = new_slide()
title(s, "Partition Pruning Cost Model", "A simple estimate students can apply")
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.0, 1.65, 5.35, 3.25, CHARCOAL)
textbox(s, "ScanCost = s * |F|", 1.45, 2.55, 4.4, 0.55, 30, True, WHITE, align="center")
paragraph_box(s, ["|F| = total fact size", "s = selected partition fraction"], 1.55, 3.55, 3.9, 0.58, 13, "DDE8EA")
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 7.1, 1.65, 4.9, 3.25, WHITE, LINE)
textbox(s, "Example", 7.42, 2.0, 1.7, 0.26, 16, True, TEAL)
paragraph_box(s, ["1 TB fact table partitioned by day", "One-day filter over 365 days", "s = 1 / 365", "Approx scan: 2.74 GB"], 7.42, 2.45, 3.8, 1.5, 14, INK)
callout(s, "Teaching point: partition pruning often reduces cost and latency roughly in proportion to bytes scanned.", 2.15, 5.58, 9.0, 0.68, WHITE, INK, 12)
footer(s, 21)

# 22
s = new_slide()
title(s, "Join Cost Intuition", "Why schema shape changes runtime")
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 0.95, 1.55, 5.45, 4.55, WHITE, LINE)
textbox(s, "Broadcast join", 1.25, 1.9, 2.2, 0.3, 17, True, TEAL)
db(s, "pruned fact", 1.45, 3.0, 1.55, 0.78, fill=MINT, line=TEAL)
db(s, "small dim", 4.1, 3.0, 1.35, 0.78, fill=AMBER_SOFT, line=AMBER)
add_line(s, 3.0, 3.39, 4.1, 3.39, TEAL, 2)
paragraph_box(s, ["Small dimension copied to workers.", "Usually fast after fact pruning."], 1.3, 4.45, 4.4, 0.65, 12, MUTED)
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 6.95, 1.55, 5.45, 4.55, WHITE, LINE)
textbox(s, "Shuffle join", 7.25, 1.9, 2.2, 0.3, 17, True, RED)
db(s, "large fact", 7.45, 3.0, 1.45, 0.78, fill=RED_SOFT, line=RED)
db(s, "large dim", 10.12, 3.0, 1.45, 0.78, fill=BLUE_SOFT, line=BLUE)
add_line(s, 8.9, 3.39, 10.12, 3.39, RED, 2)
paragraph_box(s, ["Both sides repartitioned by key.", "Skew or snowflaking can dominate runtime."], 7.3, 4.45, 4.4, 0.65, 12, MUTED)
footer(s, 22)

# 23
s = new_slide()
title(s, "Semantic Layer and Metric Contracts", "Prevent metric drift after the schema is built")
for i, (lab, c) in enumerate([("Physical tables", CHARCOAL), ("Models", TEAL), ("Certified metrics", AMBER), ("Dashboards / APIs", BLUE)]):
    add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 1.2 + i * 2.75, 2.15, 2.25, 0.82, c)
    textbox(s, lab, 1.3 + i * 2.75, 2.43, 2.05, 0.18, 10.5, True, WHITE, align="center")
    if i < 3:
        add_line(s, 3.45 + i * 2.75, 2.56, 3.88 + i * 2.75, 2.56, MUTED, 2)
card(s, 1.45, 4.35, 3.0, 1.05, "Metric contract", ["Name, formula, grain, filters, owner"], TEAL)
card(s, 5.15, 4.35, 3.0, 1.05, "Tests", ["Accepted ranges, null policy, reconciliation"], AMBER)
card(s, 8.85, 4.35, 3.0, 1.05, "Lineage", ["Tables, columns, dashboards, change history"], BLUE)
footer(s, 23)

# 24
s = new_slide()
title(s, "Governance Controls for Large Facts", "Controls should be executable, not aspirational")
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 4.8, 2.45, 3.6, 1.25, CHARCOAL)
textbox(s, "fact_sales\nlarge atomic table", 5.05, 2.82, 3.1, 0.4, 15, True, WHITE, align="center")
controls = [
    ("SQL guardrail", 1.0, 1.35, TEAL, "Block missing date filters"),
    ("Data contract", 8.95, 1.35, AMBER, "Required keys and grain checks"),
    ("Cost monitor", 1.0, 5.0, BLUE, "Bytes and partitions scanned"),
    ("Quality tests", 8.95, 5.0, RED, "Freshness, nulls, duplicates"),
]
for h, x, y, c, b in controls:
    card(s, x, y, 3.0, 0.9, h, [b], c, WHITE, 13, 10)
    add_line(s, x + (3.0 if x < 4 else 0), y + 0.45, 6.6, 3.08, c, 1.5)
footer(s, 24)

# 25
s = new_slide(True)
title(s, "Design Review Takeaways", "Before a model reaches production", True)
items = [
    "Declare grain and enforce uniqueness.",
    "Use star schema as the default BI contract.",
    "Model history explicitly with SCD choices.",
    "Conform dimensions before multiplying marts.",
    "Design partitioning and guardrails from real workload.",
]
for i, txt in enumerate(items):
    y = 1.55 + i * 0.78
    add_shape(s, MSO_AUTO_SHAPE_TYPE.OVAL, 1.05, y, 0.36, 0.36, AMBER if i % 2 else TEAL)
    textbox(s, str(i + 1), 1.145, y + 0.08, 0.17, 0.08, 8, True, WHITE, align="center", margin=0)
    textbox(s, txt, 1.7, y - 0.02, 8.9, 0.3, 19, True, WHITE)
add_shape(s, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, 8.25, 5.02, 4.35, 1.18, "334E52", "7FB7AE")
textbox(s, "Professional standard:\ncorrectness, performance, and operability are one design problem.", 8.52, 5.28, 3.8, 0.58, 12.5, True, "DDE8EA", align="center")
footer(s, 25, True)


prs.save(OUT)
print(OUT)
