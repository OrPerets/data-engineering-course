from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


TEMPLATE = Path("reference.pptx")
OUT = Path("build/08-TF-IDF-visual-reference.pptx")
USE_REFERENCE_TEMPLATE = True

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

COLORS = {
    "ink": "1F2933",
    "muted": "657383",
    "paper": "F6F7F9",
    "panel": "FFFFFF",
    "line": "D6DEE6",
    "teal": "087E8B",
    "mint": "D8F3ED",
    "coral": "F25F5C",
    "coral_light": "FFE4E1",
    "gold": "F5C542",
    "gold_light": "FFF4C7",
    "navy": "172A3A",
    "slate": "355070",
    "lav": "EAE7FF",
    "lav_dark": "5B5F97",
}

FONT_HEAD = "Trebuchet MS"
FONT_BODY = "Aptos"
FONT_CODE = "Consolas"


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    if transparency:
        shape.fill.transparency = transparency


def set_line(shape, color=None, width=1, transparency=0):
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(width)
    if transparency:
        shape.line.transparency = transparency


def add_bg(slide, color=COLORS["paper"]):
    if USE_REFERENCE_TEMPLATE:
        return
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, WIDE_W, WIDE_H)
    set_fill(shape, color)
    set_line(shape, None)


def add_text(slide, text, x, y, w, h, size=18, color=COLORS["ink"], bold=False,
             font=FONT_BODY, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    text = text.replace("\\n", "\n")
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    tf.vertical_anchor = valign
    return box


def add_rich_lines(slide, lines, x, y, w, h, size=15, color=COLORS["ink"],
                   leading=1.05, font=FONT_BODY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, item in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = leading
        p.space_after = Pt(3)
        parts = item if isinstance(item, list) else [(item, False, color)]
        for txt, is_bold, txt_color in parts:
            run = p.add_run()
            run.text = txt
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = is_bold
            run.font.color.rgb = rgb(txt_color)
    return box


def add_title(slide, title, subtitle=None, section=None, dark=False):
    color = "FFFFFF" if dark else COLORS["ink"]
    accent = COLORS["gold"] if dark else COLORS["teal"]
    if section:
        add_text(slide, section.upper(), 0.62, 0.34, 3.2, 0.3, 9, accent, True, FONT_BODY)
    add_text(slide, title, 0.62, 0.64, 8.8, 0.55, 25 if len(title) > 42 else 31,
             color, True, FONT_HEAD)
    if subtitle:
        add_text(slide, subtitle, 0.64, 1.17, 9.0, 0.36, 13, "D7DEE8" if dark else COLORS["muted"])


def footer(slide, n, total=35, dark=False):
    color = "9AA7B4"
    add_text(slide, "Or Peretz", 0.62, 7.08, 2.0, 0.22, 8, color)
    add_text(slide, str(n), 12.35, 7.08, 0.45, 0.22, 8, color, align=PP_ALIGN.RIGHT)
    add_text(slide, f"/ {total}", 12.77, 7.08, 0.38, 0.22, 8, color)


def card(slide, x, y, w, h, fill=COLORS["panel"], line=COLORS["line"], radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, line, 1)
    return shape


def label(slide, text, x, y, w, h=0.32, fill=COLORS["mint"], color=COLORS["teal"],
          size=10, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(shape, fill)
    set_line(shape, None)
    shape.text = text
    tf = shape.text_frame
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.name = FONT_BODY
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = rgb(color)
    return shape


def connector(slide, x1, y1, x2, y2, color=COLORS["teal"], width=2.0):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    set_line(line, color, width)
    return line


def add_step(slide, n, title, body, x, y, w, accent=COLORS["teal"]):
    card(slide, x, y, w, 1.12)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.22), Inches(0.48), Inches(0.48))
    set_fill(circ, accent)
    set_line(circ, None)
    add_text(slide, str(n), x + 0.18, y + 0.31, 0.48, 0.2, 12, "FFFFFF", True, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.82, y + 0.2, w - 1.02, 0.26, 13, COLORS["ink"], True)
    add_text(slide, body, x + 0.82, y + 0.52, w - 1.02, 0.42, 10.5, COLORS["muted"])


def add_token(slide, text, x, y, fill=COLORS["mint"], color=COLORS["ink"], w=None):
    width = w if w is not None else max(0.58, 0.18 + len(text) * 0.08)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(0.34))
    set_fill(shape, fill)
    set_line(shape, None)
    add_text(slide, text, x + 0.02, y + 0.08, width - 0.04, 0.14, 9.5, color, True, align=PP_ALIGN.CENTER)
    return width


def add_table(slide, rows, x, y, w, h, header_fill=COLORS["navy"], font_size=9.5):
    tbl_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = tbl_shape.table
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            set_fill(cell, header_fill if r_idx == 0 else ("FFFFFF" if r_idx % 2 else "EEF3F7"))
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if r_idx == 0 or c_idx > 0 else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(font_size)
                    run.font.bold = r_idx == 0
                    run.font.color.rgb = rgb("FFFFFF" if r_idx == 0 else COLORS["ink"])
    return tbl_shape


def add_code_box(slide, code, x, y, w, h, size=10.5):
    card(slide, x, y, w, h, fill="111827", line="111827", radius=False)
    add_text(slide, code, x + 0.18, y + 0.16, w - 0.36, h - 0.24, size, "E5E7EB", False, FONT_CODE)


def add_bar_chart(slide, items, x, y, w, h, max_value=None, accent=COLORS["teal"]):
    max_v = max_value or max(v for _, v, _ in items)
    row_h = h / len(items)
    for i, (name, value, color) in enumerate(items):
        yy = y + i * row_h
        add_text(slide, name, x, yy + 0.07, 1.25, 0.2, 9.5, COLORS["ink"], True)
        card(slide, x + 1.35, yy + 0.08, w - 1.78, 0.22, fill="E8EDF2", line="E8EDF2", radius=True)
        bw = (w - 1.78) * value / max_v
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 1.35), Inches(yy + 0.08), Inches(bw), Inches(0.22))
        set_fill(bar, color or accent)
        set_line(bar, None)
        add_text(slide, f"{value:.2f}", x + w - 0.36, yy + 0.05, 0.35, 0.2, 9, COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_matrix(slide, terms, docs, marks, x, y, cell=0.48):
    add_text(slide, "term / doc", x, y, 1.0, 0.24, 8.5, COLORS["muted"], True)
    for j, d in enumerate(docs):
        label(slide, d, x + 1.15 + j * cell, y - 0.02, cell - 0.04, 0.26, fill=COLORS["navy"], color="FFFFFF", size=8)
    for i, t in enumerate(terms):
        add_text(slide, t, x, y + 0.36 + i * cell, 1.05, 0.2, 9, COLORS["ink"], True)
        for j, d in enumerate(docs):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 1.15 + j * cell), Inches(y + 0.32 + i * cell), Inches(cell - 0.04), Inches(cell - 0.04))
            hit = (t, d) in marks
            set_fill(shape, COLORS["mint"] if hit else "F0F3F7")
            set_line(shape, "FFFFFF", 0.7)
            if hit:
                add_text(slide, "1", x + 1.15 + j * cell, y + 0.42 + i * cell, cell - 0.04, 0.12, 9, COLORS["teal"], True, align=PP_ALIGN.CENTER)


def section_slide(prs, n, title, subtitle, chips):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, COLORS["navy"])
    add_title(slide, title, subtitle, dark=True)
    x = 0.72
    for chip, color in chips:
        label(slide, chip, x, 3.7, 1.7, 0.42, fill=color, color=COLORS["ink"], size=11)
        x += 1.95
    footer(slide, n, dark=True)


def build():
    prs = Presentation(str(TEMPLATE)) if USE_REFERENCE_TEMPLATE and TEMPLATE.exists() else Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    if USE_REFERENCE_TEMPLATE:
        # Keep the template theme, masters, and layouts, but replace its sample slides.
        sld_id_lst = prs.slides._sldIdLst
        for sld_id in list(sld_id_lst):
            r_id = sld_id.rId
            prs.part.drop_rel(r_id)
            sld_id_lst.remove(sld_id)
        blank = prs.slide_layouts[8] if len(prs.slide_layouts) > 8 else prs.slide_layouts[6]
    else:
        blank = prs.slide_layouts[6]
    total = 35

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s, COLORS["navy"])
    add_text(s, "Text Analytics", 0.7, 0.62, 4.2, 0.45, 18, COLORS["teal"], True, FONT_HEAD)
    add_text(s, "Regex to TF-IDF", 0.7, 1.1, 7.0, 0.8, 42, COLORS["ink"], True, FONT_HEAD)
    add_text(s, "A visual, step-by-step path from messy operational text to ranked terms.", 0.72, 2.02, 6.2, 0.42, 15, COLORS["muted"])
    for i, (txt, fill) in enumerate([("RAW", COLORS["coral"]), ("REGEX", COLORS["gold"]), ("TOKENS", COLORS["mint"]), ("TF-IDF", COLORS["lav"])]):
        x = 0.82 + i * 1.55
        label(s, txt, x, 3.15, 1.15, 0.46, fill=fill, color=COLORS["ink"], size=11)
        if i < 3:
            connector(s, x + 1.18, 3.38, x + 1.45, 3.38, "FFFFFF", 1.5)
    add_code_box(s, '"Printer JAM!!! Ticket=TKT-1042"\\n-> clean tokens:\\n   printer, jam, ticket, tkt, 1042\\n-> score(jam) = 0.46', 7.15, 0.95, 5.25, 2.25, 13.5)
    footer(s, 1, total)

    # 2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Learning Path", "Each block builds the input needed by the next one.", "roadmap")
    steps = [
        ("1", "Messy text", "tickets, complaints, logs", COLORS["coral_light"]),
        ("2", "Regex", "extract and normalize fields", COLORS["gold_light"]),
        ("3", "Tokens", "stable units for counting", COLORS["mint"]),
        ("4", "TF-IDF", "local focus x global rarity", COLORS["lav"]),
        ("5", "Pipeline", "SQL stages and risks", "E8EDF2"),
    ]
    for i, (num, title, body, fill) in enumerate(steps):
        x = 0.72 + i * 2.45
        card(s, x, 2.25, 2.05, 2.25, fill=fill, line="FFFFFF")
        add_text(s, num, x + 0.16, 2.4, 0.42, 0.38, 18, COLORS["ink"], True, FONT_HEAD)
        add_text(s, title, x + 0.18, 3.05, 1.7, 0.26, 15, COLORS["ink"], True)
        add_text(s, body, x + 0.18, 3.48, 1.64, 0.5, 11, COLORS["muted"])
        if i < 4:
            connector(s, x + 2.08, 3.38, x + 2.32, 3.38)
    footer(s, 2, total)

    # 3
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Why Text Analytics Matters", "Most operational signals arrive as text before they become data.", "motivation")
    sources = [("Tickets", "delays, faults, queues"), ("Logs", "errors, timeouts, retries"), ("Complaints", "refunds, churn, severity"), ("Surveys", "free-text satisfaction"), ("Notes", "supplier and field context")]
    for i, (name, desc) in enumerate(sources):
        x = 0.75 + (i % 3) * 3.95
        y = 2.0 + (i // 3) * 1.55
        card(s, x, y, 3.35, 1.12)
        label(s, name, x + 0.2, y + 0.18, 1.12, 0.3, fill=COLORS["mint"], color=COLORS["teal"])
        add_text(s, desc, x + 0.2, y + 0.58, 2.8, 0.24, 11, COLORS["muted"])
    add_text(s, "The goal is not to read everything manually. The goal is to turn language into auditable features.", 7.0, 5.62, 4.95, 0.5, 16, COLORS["ink"], True)
    footer(s, 3, total)

    # 4
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "From Raw Text to Structured Features", "Regex is usually the first stabilizing step in the pipeline.", "regex")
    pipeline = [("raw", "Printer   JAM!!!\\nTicket=TKT-1042"), ("extract", "ticket_id = TKT-1042"), ("normalize", "printer jam ticket tkt 1042"), ("tokens", "printer | jam | ticket | tkt | 1042")]
    for i, (head, body) in enumerate(pipeline):
        x = 0.8 + i * 3.05
        card(s, x, 2.25, 2.45, 1.55, fill="FFFFFF")
        label(s, head.upper(), x + 0.2, 2.45, 0.95, 0.28, fill=COLORS["gold_light"], color=COLORS["ink"], size=8.5)
        add_text(s, body, x + 0.22, 2.95, 2.0, 0.46, 13, COLORS["ink"], True, FONT_CODE if i == 0 else FONT_BODY)
        if i < 3:
            connector(s, x + 2.5, 3.02, x + 2.86, 3.02, COLORS["teal"], 2)
    add_text(s, "Design rule: every transformation should be explainable, repeatable, and testable on examples.", 1.05, 5.08, 10.5, 0.42, 17, COLORS["teal"], True)
    footer(s, 4, total)

    # 5
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Regex Basics as Building Blocks", "A regex is a small grammar for matching text.", "regex")
    blocks = [("literal", "TKT", "match these characters"), ("class", "\\\\d", "match one digit"), ("quantifier", "{4,6}", "repeat 4 to 6 times"), ("boundary", "\\\\b", "do not leak into neighbors")]
    for i, (name, pat, desc) in enumerate(blocks):
        x = 0.75 + i * 3.05
        card(s, x, 2.0, 2.55, 2.3, fill=["FFFFFF", COLORS["mint"], COLORS["gold_light"], COLORS["coral_light"]][i])
        add_text(s, pat, x + 0.2, 2.28, 2.1, 0.42, 24, COLORS["ink"], True, FONT_CODE, PP_ALIGN.CENTER)
        add_text(s, name.upper(), x + 0.2, 3.15, 2.1, 0.24, 10, COLORS["teal"], True, align=PP_ALIGN.CENTER)
        add_text(s, desc, x + 0.28, 3.55, 1.95, 0.34, 11, COLORS["muted"], align=PP_ALIGN.CENTER)
    add_code_box(s, r"\bTKT-\d{4,6}\b", 3.6, 5.12, 6.0, 0.7, 23)
    footer(s, 5, total)

    # 6
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Character Classes: What Counts as a Match?", "Classes define the alphabet that a pattern can consume.", "regex lab")
    rows = [["Pattern", "Meaning", "Example match"], [r"\d", "one digit", "7 in delay 7 days"], [r"\d+", "one or more digits", "1200 in USD 1200"], [r"[A-Z]+", "uppercase letters", "ERR in ERR-404"], [r"\w+", "word characters", "printer"]]
    add_table(s, rows, 0.82, 1.85, 5.75, 2.55, font_size=10)
    card(s, 7.08, 1.92, 4.9, 2.35, fill="FFFFFF")
    add_text(s, "Scanner view", 7.38, 2.15, 2.0, 0.24, 14, COLORS["ink"], True)
    add_code_box(s, "delay 7 days\\n      ^\\n      \\\\d", 7.45, 2.62, 3.9, 0.95, 18)
    add_text(s, "Changing the class changes the evidence collected downstream.", 7.42, 3.82, 3.9, 0.3, 12, COLORS["muted"])
    footer(s, 6, total)

    # 7
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Quantifiers Encode Business Rules", "They specify how many times the previous unit may repeat.", "regex lab")
    add_code_box(s, r"TKT-\d{4,6}", 0.85, 1.85, 3.9, 0.82, 24)
    cases = [("TKT-1042", "ACCEPT", COLORS["mint"]), ("TKT-998877", "ACCEPT", COLORS["mint"]), ("TKT-12", "REJECT", COLORS["coral_light"]), ("TKT-1042A", "REJECT", COLORS["coral_light"])]
    for i, (txt, status, fill) in enumerate(cases):
        y = 3.05 + i * 0.72
        card(s, 0.95, y, 3.4, 0.48, fill=fill)
        add_text(s, txt, 1.15, y + 0.13, 1.55, 0.15, 11, COLORS["ink"], True, FONT_CODE)
        add_text(s, status, 3.15, y + 0.13, 0.82, 0.15, 10, COLORS["teal"] if status == "ACCEPT" else COLORS["coral"], True, align=PP_ALIGN.RIGHT)
    add_text(s, "{4,6}", 6.45, 2.05, 1.2, 0.4, 36, COLORS["coral"], True, FONT_HEAD, PP_ALIGN.CENTER)
    connector(s, 5.08, 2.3, 6.25, 2.3, COLORS["coral"], 2.5)
    add_text(s, "This is the rule: valid ticket IDs contain 4 to 6 digits.", 7.85, 2.05, 3.9, 0.46, 18, COLORS["ink"], True)
    add_text(s, "When this rule changes, historical data quality changes too. Document it.", 7.88, 3.02, 3.65, 0.5, 13, COLORS["muted"])
    footer(s, 7, total)

    # 8
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Capture Groups Turn Matches into Fields", "Parentheses mark the pieces we want to keep.", "regex lab")
    add_code_box(s, r"delay\s+(\d+)\s+days?", 0.85, 1.75, 5.2, 0.78, 22)
    add_code_box(s, "delay 3 days due to customs", 0.85, 3.0, 5.2, 0.62, 17)
    connector(s, 3.13, 2.55, 3.13, 2.95, COLORS["coral"], 2.5)
    label(s, "Group 1 = 3", 2.32, 4.02, 1.58, 0.42, fill=COLORS["coral_light"], color=COLORS["coral"], size=12)
    card(s, 7.0, 1.72, 4.85, 2.6)
    add_text(s, "Structured output", 7.35, 2.0, 2.4, 0.3, 17, COLORS["ink"], True)
    rows = [["field", "value"], ["delay_days", "3"], ["reason_text", "due to customs"]]
    add_table(s, rows, 7.35, 2.58, 3.7, 1.22, font_size=10)
    add_text(s, "The text remains messy, but the extracted feature is now queryable.", 7.35, 4.05, 3.6, 0.26, 11, COLORS["muted"])
    footer(s, 8, total)

    # 9
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Cleaning Example: One Record, Four Decisions", "Small choices in preprocessing change the vocabulary.", "regex lab")
    add_code_box(s, '"Printer   JAM!!! Ticket=TKT-1042"', 0.85, 1.65, 5.2, 0.62, 16)
    add_step(s, 1, "Extract ID", "TKT-1042 becomes a field", 0.85, 2.62, 3.1, COLORS["coral"])
    add_step(s, 2, "Lowercase", "Printer JAM -> printer jam", 4.08, 2.62, 3.1, COLORS["teal"])
    add_step(s, 3, "Punctuation", "replace with spaces", 7.31, 2.62, 3.1, COLORS["gold"])
    add_step(s, 4, "Tokenize", "split into stable terms", 10.54, 2.62, 2.05, COLORS["lav_dark"])
    x = 2.0
    for tok in ["printer", "jam", "ticket", "tkt", "1042"]:
        x += add_token(s, tok, x, 5.15, fill=COLORS["mint"]) + 0.12
    add_text(s, "Question for students: should TKT-1042 stay as one masked token?", 3.32, 5.92, 6.7, 0.28, 14, COLORS["teal"], True, align=PP_ALIGN.CENTER)
    footer(s, 9, total)

    # 10
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Regex Case: Amount Extraction", "Groups separate the currency from the numeric value.", "regex case")
    add_code_box(s, r"(USD|ILS|EUR)\s?(\d+(?:\.\d{1,2})?)", 0.78, 1.58, 6.4, 0.68, 17)
    add_code_box(s, "refund USD 120.50 requested", 0.78, 2.65, 6.4, 0.62, 17)
    connector(s, 3.0, 3.32, 3.0, 3.78, COLORS["gold"], 2.2)
    connector(s, 4.18, 3.32, 4.18, 3.78, COLORS["coral"], 2.2)
    label(s, "currency = USD", 2.18, 3.88, 1.68, 0.4, fill=COLORS["gold_light"], color=COLORS["ink"], size=11)
    label(s, "amount = 120.50", 3.98, 3.88, 1.88, 0.4, fill=COLORS["coral_light"], color=COLORS["coral"], size=11)
    card(s, 8.0, 1.62, 4.3, 3.15)
    add_text(s, "Non-capturing group", 8.32, 1.94, 2.8, 0.3, 17, COLORS["ink"], True)
    add_code_box(s, r"(?:\.\d{1,2})?", 8.32, 2.52, 2.9, 0.58, 16)
    add_text(s, "Used for structure, not output. It lets decimals be optional without creating another extracted field.", 8.32, 3.42, 3.35, 0.56, 12, COLORS["muted"])
    footer(s, 10, total)

    # 11
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Tokenization Rules Must Be Explicit", "Different rules create different feature spaces.", "tokens")
    examples = [("keep IDs", ["ticket", "tkt-1042", "jam"]), ("split IDs", ["ticket", "tkt", "1042", "jam"]), ("mask IDs", ["ticket", "<ticket_id>", "jam"])]
    for i, (name, toks) in enumerate(examples):
        x0 = 0.9 + i * 4.1
        card(s, x0, 2.0, 3.5, 2.25)
        label(s, name.upper(), x0 + 0.28, 2.28, 1.35, 0.32, fill=[COLORS["mint"], COLORS["gold_light"], COLORS["coral_light"]][i], color=COLORS["ink"])
        yy = 3.0
        xx = x0 + 0.32
        for tok in toks:
            xx += add_token(s, tok, xx, yy, fill="E8EDF2", w=max(0.8, len(tok) * 0.09 + 0.2)) + 0.1
            if xx > x0 + 2.7:
                yy += 0.48
                xx = x0 + 0.32
    add_text(s, "The rule is not cosmetic. It changes counts, vocabulary size, and model behavior.", 1.55, 5.42, 10.0, 0.38, 17, COLORS["teal"], True, align=PP_ALIGN.CENTER)
    footer(s, 11, total)

    # 12
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Operational Vocabulary", "TF-IDF uses three levels: term, document, corpus.", "tokens")
    card(s, 0.9, 1.85, 3.2, 3.5, fill=COLORS["mint"])
    add_text(s, "TERM", 1.15, 2.18, 1.4, 0.28, 15, COLORS["teal"], True)
    add_token(s, "jam", 1.25, 2.88, fill="FFFFFF")
    add_token(s, "vpn", 2.05, 2.88, fill="FFFFFF")
    add_token(s, "offline", 1.48, 3.42, fill="FFFFFF")
    add_text(s, "A cleaned word or phrase that we count.", 1.17, 4.35, 2.3, 0.32, 12, COLORS["ink"])
    card(s, 5.05, 1.85, 3.2, 3.5, fill=COLORS["gold_light"])
    add_text(s, "DOCUMENT", 5.3, 2.18, 1.9, 0.28, 15, COLORS["ink"], True)
    add_code_box(s, "T1: printer jam jam", 5.32, 2.82, 2.25, 0.48, 11)
    add_text(s, "One ticket, complaint, log line, or survey answer.", 5.32, 4.16, 2.35, 0.38, 12, COLORS["ink"])
    card(s, 9.2, 1.85, 3.2, 3.5, fill=COLORS["lav"])
    add_text(s, "CORPUS", 9.45, 2.18, 1.7, 0.28, 15, COLORS["lav_dark"], True)
    for j, t in enumerate(["T1", "T2", "T3", "T4", "T5"]):
        label(s, t, 9.55 + (j % 3) * 0.7, 2.85 + (j // 3) * 0.5, 0.5, 0.3, fill="FFFFFF", color=COLORS["ink"])
    add_text(s, "A versioned set of documents.", 9.52, 4.17, 2.2, 0.26, 12, COLORS["ink"])
    footer(s, 12, total)

    # 13
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "TF-IDF Intuition", "A high score needs both local focus and global rarity.", "tf-idf")
    card(s, 1.0, 2.0, 4.8, 2.75, fill=COLORS["mint"])
    add_text(s, "Term Frequency", 1.32, 2.35, 2.7, 0.32, 20, COLORS["teal"], True, FONT_HEAD)
    add_text(s, "How much does this document talk about the term?", 1.35, 3.12, 3.75, 0.42, 16, COLORS["ink"], True)
    add_bar_chart(s, [("jam", 0.67, COLORS["teal"]), ("printer", 0.33, COLORS["gold"])], 1.35, 3.92, 3.55, 0.75, max_value=1)
    card(s, 7.25, 2.0, 4.8, 2.75, fill=COLORS["lav"])
    add_text(s, "Inverse Document Frequency", 7.58, 2.35, 3.9, 0.32, 20, COLORS["lav_dark"], True, FONT_HEAD)
    add_text(s, "How rare is this term across the corpus?", 7.6, 3.12, 3.5, 0.42, 16, COLORS["ink"], True)
    add_bar_chart(s, [("outage", 1.10, COLORS["coral"]), ("jam", 0.69, COLORS["teal"])], 7.6, 3.92, 3.55, 0.75, max_value=1.1, accent=COLORS["coral"])
    add_text(s, "TF-IDF = local focus x global rarity", 3.42, 5.55, 6.7, 0.45, 22, COLORS["ink"], True, FONT_HEAD, PP_ALIGN.CENTER)
    footer(s, 13, total)

    # 14
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Notation Cheat Sheet", "Keep the symbols close to the worked example.", "tf-idf")
    rows = [["Symbol", "Meaning", "Example"], ["t", "term after preprocessing", "jam"], ["d", "one document", "T1"], ["N", "number of documents", "5"], ["count(t,d)", "count of t in d", "2"], ["|d|", "number of tokens in d", "3"], ["df(t)", "documents containing t", "2"]]
    add_table(s, rows, 0.95, 1.72, 6.05, 4.0, font_size=10.5)
    card(s, 7.7, 2.0, 4.2, 2.9, fill=COLORS["panel"])
    add_text(s, "Anchor example", 8.03, 2.28, 2.2, 0.26, 18, COLORS["ink"], True)
    add_code_box(s, "T1 = printer jam jam", 8.05, 2.9, 2.9, 0.55, 12.5)
    add_text(s, "For t = jam and d = T1:\\ncount(t,d) = 2\\n|d| = 3", 8.05, 3.7, 2.8, 0.72, 13, COLORS["ink"], True)
    footer(s, 14, total)

    # 15
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Worked Example: Ticket Corpus", "Five tiny tickets are enough to see the mechanics.", "worked example")
    rows = [["ticket", "text"], ["T1", "printer jam jam"], ["T2", "printer offline"], ["T3", "vpn outage"], ["T4", "vpn slow"], ["T5", "printer jam"]]
    add_table(s, rows, 0.85, 1.75, 5.2, 3.45, font_size=11)
    add_text(s, "Corpus size", 7.25, 2.05, 2.0, 0.26, 16, COLORS["muted"], True)
    add_text(s, "N = 5", 7.25, 2.42, 2.6, 0.75, 44, COLORS["coral"], True, FONT_HEAD)
    add_text(s, "Every calculation below uses this exact corpus definition.", 7.28, 3.48, 3.4, 0.4, 15, COLORS["ink"], True)
    footer(s, 15, total)

    # 16
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Step 1: Count Tokens in Each Document", "Counting only starts after cleaning and tokenization.", "worked example")
    docs = [("T1", ["printer", "jam", "jam"]), ("T2", ["printer", "offline"]), ("T3", ["vpn", "outage"]), ("T4", ["vpn", "slow"]), ("T5", ["printer", "jam"])]
    for i, (doc, toks) in enumerate(docs):
        x0 = 0.85 + (i % 3) * 4.0
        y0 = 1.8 + (i // 3) * 1.65
        card(s, x0, y0, 3.25, 1.2)
        label(s, doc, x0 + 0.18, y0 + 0.2, 0.55, 0.3, fill=COLORS["navy"], color="FFFFFF")
        xx = x0 + 0.88
        for tok in toks:
            xx += add_token(s, tok, xx, y0 + 0.23, fill=COLORS["mint"]) + 0.08
        add_text(s, f"length = {len(toks)}", x0 + 0.88, y0 + 0.76, 1.5, 0.2, 10, COLORS["muted"], True)
    footer(s, 16, total)

    # 17
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Step 2: Term Frequency Is Local", "TF measures a term inside one document.", "worked example")
    add_code_box(s, "T1 = printer jam jam", 0.9, 1.72, 4.7, 0.58, 16)
    add_text(s, "tf(jam, T1) = count(jam,T1) / |T1|", 0.95, 2.7, 5.6, 0.38, 20, COLORS["ink"], True, FONT_CODE)
    add_text(s, "= 2 / 3 = 0.67", 2.08, 3.35, 3.2, 0.52, 28, COLORS["coral"], True, FONT_HEAD, PP_ALIGN.CENTER)
    rows = [["term/doc", "calculation", "TF"], ["jam in T1", "2/3", "0.67"], ["offline in T2", "1/2", "0.50"], ["outage in T3", "1/2", "0.50"]]
    add_table(s, rows, 7.0, 1.8, 4.9, 2.25, font_size=10.5)
    add_text(s, "TF alone does not know whether a term is common or rare across tickets.", 7.1, 4.5, 4.55, 0.42, 15, COLORS["teal"], True)
    footer(s, 17, total)

    # 18
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Step 3: Document Frequency Counts Documents", "Repeated mentions inside the same document still count once.", "worked example")
    terms = ["jam", "offline", "outage", "printer", "vpn"]
    docs_names = ["T1", "T2", "T3", "T4", "T5"]
    marks = {("jam", "T1"), ("jam", "T5"), ("offline", "T2"), ("outage", "T3"), ("printer", "T1"), ("printer", "T2"), ("printer", "T5"), ("vpn", "T3"), ("vpn", "T4")}
    add_matrix(s, terms, docs_names, marks, 0.95, 1.85, cell=0.62)
    rows = [["term", "documents", "DF"], ["jam", "T1, T5", "2"], ["offline", "T2", "1"], ["outage", "T3", "1"]]
    add_table(s, rows, 7.25, 2.0, 4.35, 2.05, font_size=10.5)
    add_text(s, "DF is global. It requires aggregation across the corpus.", 7.3, 4.6, 4.1, 0.38, 16, COLORS["teal"], True)
    footer(s, 18, total)

    # 19
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Step 4: IDF Converts Commonness into Rarity Weight", "The +1 terms smooth the formula and avoid division by zero.", "worked example")
    add_text(s, "idf(t) = log((N + 1) / (df(t) + 1))", 0.95, 1.8, 6.4, 0.46, 22, COLORS["ink"], True, FONT_CODE)
    rows = [["term", "DF", "IDF"], ["jam", "2", "log(6/3)=0.69"], ["offline", "1", "log(6/2)=1.10"], ["outage", "1", "log(6/2)=1.10"]]
    add_table(s, rows, 0.95, 2.75, 5.35, 2.25, font_size=10.5)
    add_bar_chart(s, [("jam", 0.69, COLORS["teal"]), ("offline", 1.10, COLORS["coral"]), ("outage", 1.10, COLORS["coral"])], 7.15, 2.55, 4.35, 1.25, max_value=1.1)
    add_text(s, "Rare terms receive higher IDF under the same corpus definition.", 7.18, 4.32, 4.0, 0.35, 15, COLORS["ink"], True)
    footer(s, 19, total)

    # 20
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Step 5: Multiply TF by IDF", "Now local focus and global rarity meet.", "worked example")
    rows = [["term/doc", "TF", "IDF", "TF-IDF"], ["jam in T1", "0.67", "0.69", "0.46"], ["offline in T2", "0.50", "1.10", "0.55"], ["outage in T3", "0.50", "1.10", "0.55"]]
    add_table(s, rows, 0.85, 1.65, 6.2, 2.65, font_size=11)
    add_text(s, "Same TF, higher IDF", 7.65, 1.82, 3.2, 0.3, 18, COLORS["ink"], True)
    add_text(s, "offline and outage rank higher even with fewer raw mentions.", 7.68, 2.42, 3.6, 0.55, 15, COLORS["muted"])
    add_bar_chart(s, [("offline", 0.55, COLORS["coral"]), ("outage", 0.55, COLORS["coral"]), ("jam", 0.46, COLORS["teal"])], 7.65, 3.35, 3.95, 1.2, max_value=0.6, accent=COLORS["coral"])
    footer(s, 20, total)

    # 21
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Why Rare Operational Terms Can Win", "Raw counts and TF-IDF answer different questions.", "interpretation")
    add_text(s, "Raw frequency", 1.0, 1.75, 2.2, 0.3, 18, COLORS["ink"], True)
    add_bar_chart(s, [("printer", 3, COLORS["teal"]), ("jam", 3, COLORS["teal"]), ("offline", 1, COLORS["coral"]), ("outage", 1, COLORS["coral"])], 0.95, 2.25, 4.7, 1.7, max_value=3)
    add_text(s, "TF-IDF signal", 7.0, 1.75, 2.2, 0.3, 18, COLORS["ink"], True)
    add_bar_chart(s, [("offline", 0.55, COLORS["coral"]), ("outage", 0.55, COLORS["coral"]), ("jam", 0.46, COLORS["teal"]), ("printer", 0.18, COLORS["teal"])], 6.95, 2.25, 4.7, 1.7, max_value=0.6)
    add_text(s, "Interpretation: rare operational states may deserve attention even when volume is small.", 2.2, 5.35, 8.8, 0.4, 18, COLORS["teal"], True, align=PP_ALIGN.CENTER)
    footer(s, 21, total)

    # 22
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Sanity Check: Common Terms Fade", "If a term appears everywhere, its IDF approaches zero.", "interpretation")
    rows = [["term", "appears in", "effect"], ["please", "almost every ticket", "low IDF"], ["printer", "many printer tickets", "moderate signal"], ["outage", "few tickets", "higher signal"]]
    add_table(s, rows, 0.92, 1.85, 5.8, 2.45, font_size=11)
    card(s, 7.55, 1.82, 4.05, 2.45, fill=COLORS["coral_light"])
    add_text(s, "Important caution", 7.9, 2.12, 2.2, 0.28, 17, COLORS["coral"], True)
    add_text(s, "A high score is a signal for review, not an automatic business decision.", 7.9, 2.86, 3.1, 0.66, 16, COLORS["ink"], True)
    add_text(s, "Use business severity, recency, ownership, and quality flags alongside TF-IDF.", 7.9, 3.76, 3.2, 0.3, 11.5, COLORS["muted"])
    footer(s, 22, total)

    # 23
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Second Example: Customer Complaints", "The same principle works outside IT tickets.", "interpretation")
    rows = [["id", "text", "top signal"], ["C1", "billing error refund", "refund"], ["C2", "billing error", "common issue"], ["C3", "refund delayed", "delayed"], ["C4", "app crash", "crash"], ["C5", "billing error", "common issue"]]
    add_table(s, rows, 0.78, 1.65, 6.6, 3.35, font_size=10.5)
    card(s, 8.0, 1.85, 3.95, 2.9, fill=COLORS["lav"])
    add_text(s, "Contrast", 8.34, 2.2, 1.7, 0.28, 18, COLORS["lav_dark"], True)
    add_text(s, "billing is frequent; crash is rare and potentially severe.", 8.34, 2.85, 3.0, 0.6, 17, COLORS["ink"], True)
    add_text(s, "TF-IDF helps surface terms worth inspecting, then domain rules decide action.", 8.34, 3.78, 3.0, 0.42, 12, COLORS["muted"])
    footer(s, 23, total)

    # 24
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "SQL Lens: TF-IDF as a Pipeline", "The math becomes a set of materialized stages.", "data engineering")
    stages = [("tokens", "one row per token"), ("term_counts", "GROUP BY doc, term"), ("doc_lengths", "GROUP BY doc"), ("df", "GROUP BY term"), ("scores", "join + multiply")]
    for i, (name, desc) in enumerate(stages):
        x = 0.7 + i * 2.48
        card(s, x, 2.38, 2.0, 1.45, fill=["FFFFFF", COLORS["mint"], COLORS["gold_light"], COLORS["lav"], COLORS["coral_light"]][i])
        add_text(s, name, x + 0.16, 2.66, 1.6, 0.24, 14, COLORS["ink"], True)
        add_text(s, desc, x + 0.16, 3.08, 1.55, 0.32, 10.2, COLORS["muted"])
        if i < 4:
            connector(s, x + 2.02, 3.1, x + 2.28, 3.1, COLORS["teal"], 2)
    add_text(s, "Production pipelines usually materialize stages for debugging, replay, and cost control.", 1.55, 5.28, 9.8, 0.38, 16, COLORS["teal"], True, align=PP_ALIGN.CENTER)
    footer(s, 24, total)

    # 25
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "SQL Stage 1: Token Counts", "This produces sparse feature rows.", "sql lens")
    code = """SELECT
  doc_id,
  term,
  COUNT(*) AS term_count
FROM tokens
GROUP BY doc_id, term;"""
    add_code_box(s, code, 0.85, 1.6, 5.25, 2.58, 13)
    card(s, 7.0, 1.78, 4.55, 2.22)
    add_text(s, "Output shape", 7.3, 2.08, 1.8, 0.26, 17, COLORS["ink"], True)
    add_text(s, "one row per non-zero (doc_id, term) pair", 7.3, 2.72, 3.45, 0.34, 15, COLORS["teal"], True)
    add_text(s, "Bottleneck: token explosion and high-cardinality grouping.", 7.3, 3.36, 3.45, 0.3, 12, COLORS["muted"])
    footer(s, 25, total)

    # 26
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "SQL Stage 2: Document Lengths", "TF needs a denominator for every document.", "sql lens")
    code = """SELECT
  doc_id,
  COUNT(*) AS doc_len
FROM tokens
GROUP BY doc_id;"""
    add_code_box(s, code, 0.85, 1.6, 5.25, 2.25, 13)
    add_step(s, 1, "Count tokens", "document length is explicit", 6.95, 1.7, 4.5, COLORS["teal"])
    add_step(s, 2, "Cast numerics", "avoid integer division", 6.95, 3.05, 4.5, COLORS["gold"])
    add_step(s, 3, "Filter edge cases", "empty documents cannot divide", 6.95, 4.4, 4.5, COLORS["coral"])
    footer(s, 26, total)

    # 27
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "SQL Stage 3: Document Frequency", "DF is the expensive global aggregation.", "sql lens")
    code = """SELECT
  term,
  COUNT(DISTINCT doc_id) AS df
FROM tokens
GROUP BY term;"""
    add_code_box(s, code, 0.85, 1.6, 5.25, 2.25, 13)
    card(s, 7.0, 1.65, 4.6, 2.6, fill=COLORS["gold_light"])
    add_text(s, "Skew risk", 7.35, 1.98, 1.8, 0.3, 18, COLORS["ink"], True)
    add_text(s, "Common boilerplate terms create hot groups and uneven shuffle partitions.", 7.35, 2.75, 3.3, 0.6, 16, COLORS["ink"], True)
    add_text(s, "Watch max/median partition load, not only total runtime.", 7.35, 3.72, 3.2, 0.26, 11.5, COLORS["muted"])
    footer(s, 27, total)

    # 28
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "SQL Stage 4: Final TF-IDF Join", "Bring local and global statistics together.", "sql lens")
    code = """SELECT
  tc.doc_id,
  tc.term,
  (tc.term_count * 1.0 / dl.doc_len)
    * LOG((c.n_docs + 1.0) / (df.df + 1.0)) AS tfidf
FROM term_counts tc
JOIN doc_lengths dl USING (doc_id)
JOIN document_frequency df USING (term)
CROSS JOIN corpus_size c;"""
    add_code_box(s, code, 0.75, 1.45, 6.65, 4.25, 10.4)
    for i, (name, y) in enumerate([("term_counts", 1.78), ("doc_lengths", 2.7), ("document_frequency", 3.62), ("corpus_size", 4.54)]):
        label(s, name, 8.05, y, 2.05, 0.36, fill=["E8EDF2", COLORS["mint"], COLORS["lav"], COLORS["gold_light"]][i], color=COLORS["ink"], size=10)
        connector(s, 10.18, y + 0.18, 11.0, 3.16, COLORS["teal"], 1.4)
    card(s, 11.05, 2.7, 1.25, 0.9, fill=COLORS["coral_light"])
    add_text(s, "score", 11.22, 3.03, 0.82, 0.18, 13, COLORS["coral"], True, align=PP_ALIGN.CENTER)
    footer(s, 28, total)

    # 29
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "SQL Cost Consequences", "Each stage has a different driver and bottleneck.", "data engineering")
    rows = [["Operation", "Cardinality driver", "Bottleneck"], ["Tokenization", "total token count", "row explosion"], ["Term counts", "unique (doc, term)", "large GROUP BY"], ["DF", "unique terms + common terms", "hot-key aggregation"], ["Final join", "sparse feature rows", "memory and shuffle"]]
    add_table(s, rows, 0.78, 1.58, 7.2, 3.25, font_size=10.5)
    card(s, 8.55, 1.82, 3.4, 2.6, fill=COLORS["mint"])
    add_text(s, "Main lesson", 8.88, 2.16, 1.9, 0.28, 18, COLORS["teal"], True)
    add_text(s, "SQL can express TF-IDF. The engineering issue is predictable cost and stability at scale.", 8.88, 2.9, 2.55, 0.74, 16, COLORS["ink"], True)
    footer(s, 29, total)

    # 30
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Operational Risk: Stop-Word Skew", "Common terms have low value but high processing cost.", "risks")
    add_bar_chart(s, [("the", 100, COLORS["coral"]), ("please", 90, COLORS["coral"]), ("printer", 28, COLORS["teal"]), ("outage", 6, COLORS["gold"])], 1.0, 2.0, 5.2, 1.8, max_value=100)
    card(s, 7.15, 1.85, 4.6, 2.65, fill=COLORS["coral_light"])
    add_text(s, "Mitigation", 7.48, 2.18, 1.8, 0.28, 18, COLORS["coral"], True)
    add_rich_lines(s, [[("stop-word filtering", True, COLORS["ink"])], [("maximum DF threshold", True, COLORS["ink"])], [("partition load monitoring", True, COLORS["ink"])]], 7.52, 2.85, 3.4, 1.05, size=14)
    footer(s, 30, total)

    # 31
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Operational Risk: Vocabulary Explosion", "IDs, URLs, serials, and typos can become unique terms.", "risks")
    words = ["AB12X9", "url_884", "printer", "tkt_1042", "hash_9f3a", "vpn", "SN-778A", "pritner", "refund", "http_x"]
    x, y = 0.9, 2.0
    for i, word in enumerate(words):
        fill = COLORS["coral_light"] if any(c.isdigit() for c in word) or "_" in word or "-" in word else COLORS["mint"]
        x += add_token(s, word, x, y, fill=fill, w=max(0.86, len(word) * 0.09 + 0.25)) + 0.15
        if x > 6.3:
            y += 0.58
            x = 0.9
    card(s, 7.45, 1.85, 4.25, 2.85, fill=COLORS["gold_light"])
    add_text(s, "Mitigation", 7.78, 2.18, 1.8, 0.28, 18, COLORS["ink"], True)
    label(s, "mask identifiers", 7.78, 2.82, 1.65, 0.34, fill="FFFFFF", color=COLORS["ink"], size=10)
    label(s, "canonicalize typos", 9.6, 2.82, 1.78, 0.34, fill="FFFFFF", color=COLORS["ink"], size=10)
    label(s, "min DF threshold", 7.78, 3.36, 1.7, 0.34, fill="FFFFFF", color=COLORS["ink"], size=10)
    label(s, "vocab-size alerts", 9.65, 3.36, 1.68, 0.34, fill="FFFFFF", color=COLORS["ink"], size=10)
    add_text(s, "Do not let random identifiers become analytical vocabulary.", 7.78, 4.08, 3.05, 0.3, 12, COLORS["muted"])
    footer(s, 31, total)

    # 32
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Operational Risk: Edge Cases", "Quality handling is part of the feature pipeline.", "risks")
    rows = [["Edge case", "Risk", "Handling"], ["Empty text", "division by zero", "filter or mark invalid"], ["All stop words", "no useful features", "quality flag"], ["Very short text", "unstable ranking", "combine with rules"], ["Numeric-heavy text", "vocabulary explosion", "extract or mask patterns"]]
    add_table(s, rows, 0.85, 1.62, 7.7, 3.25, font_size=10.5)
    card(s, 9.05, 1.95, 3.0, 2.48, fill=COLORS["lav"])
    add_text(s, "Rule", 9.36, 2.3, 1.2, 0.28, 18, COLORS["lav_dark"], True)
    add_text(s, "Never hide quality problems inside a score.", 9.36, 3.0, 2.2, 0.58, 17, COLORS["ink"], True)
    footer(s, 32, total)

    # 33
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Make Quality Flags Visible", "Analysts and managers need to know when scores are based on weak text.", "risks")
    card(s, 0.85, 1.7, 11.6, 3.8, fill="FFFFFF")
    rows = [["doc", "top term", "tfidf", "quality"], ["T1", "jam", "0.46", "OK"], ["T2", "offline", "0.55", "OK"], ["T8", "-", "-", "EMPTY_TEXT"], ["T9", "sn_778a", "0.99", "NUMERIC_HEAVY"]]
    add_table(s, rows, 1.25, 2.08, 6.8, 2.45, font_size=10.5)
    label(s, "OK", 8.7, 2.42, 0.8, 0.34, fill=COLORS["mint"], color=COLORS["teal"])
    label(s, "EMPTY_TEXT", 8.7, 3.24, 1.45, 0.34, fill=COLORS["coral_light"], color=COLORS["coral"])
    label(s, "NUMERIC_HEAVY", 8.7, 4.06, 1.7, 0.34, fill=COLORS["gold_light"], color=COLORS["ink"])
    add_text(s, "A score without a data-quality state is hard to govern.", 2.2, 5.95, 8.6, 0.35, 17, COLORS["teal"], True, align=PP_ALIGN.CENTER)
    footer(s, 33, total)

    # 34
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_title(s, "Classroom Exercise", "Students can compute the ranking by hand before seeing SQL.", "practice")
    rows = [["doc", "text"], ["D1", "api timeout timeout"], ["D2", "api retry"], ["D3", "payment timeout"], ["D4", "payment failed"]]
    add_table(s, rows, 0.9, 1.55, 4.7, 2.45, font_size=11)
    add_step(s, 1, "Token counts", "count timeout in D1", 6.35, 1.55, 4.6, COLORS["teal"])
    add_step(s, 2, "TF", "divide by document length", 6.35, 2.9, 4.6, COLORS["gold"])
    add_step(s, 3, "DF and IDF", "compare timeout vs failed", 6.35, 4.25, 4.6, COLORS["coral"])
    add_text(s, "Prompt: which term should surface for D4, and why?", 2.05, 5.72, 8.8, 0.42, 19, COLORS["ink"], True, align=PP_ALIGN.CENTER)
    footer(s, 34, total)

    # 35
    s = prs.slides.add_slide(blank)
    add_bg(s, COLORS["navy"])
    add_title(s, "Takeaways", "Regex creates stable terms; TF-IDF ranks local focus against global rarity.")
    takeaways = [("Regex", "extract fields and clean text"), ("Tokens", "define the feature vocabulary"), ("TF-IDF", "rank rare focused signals"), ("SQL", "make cost and quality explicit")]
    for i, (head, body) in enumerate(takeaways):
        x = 0.95 + i * 3.05
        card(s, x, 2.35, 2.45, 1.55, fill=["FFFFFF", COLORS["mint"], COLORS["gold_light"], COLORS["coral_light"]][i], line="FFFFFF")
        add_text(s, head, x + 0.2, 2.68, 1.9, 0.26, 16, COLORS["ink"], True)
        add_text(s, body, x + 0.2, 3.1, 1.9, 0.34, 10.8, COLORS["muted"])
    add_text(s, "Practical rule: every text score needs preprocessing rules, example tests, and quality flags.", 1.35, 5.15, 10.5, 0.42, 18, COLORS["teal"], True, align=PP_ALIGN.CENTER)
    footer(s, 35, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
