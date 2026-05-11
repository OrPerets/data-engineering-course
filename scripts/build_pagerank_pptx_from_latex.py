from pathlib import Path
import math
import shutil
import tempfile
import zipfile

import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
TEX = ROOT / "lectures/07-mapreduce-advanced/pagerank_algorithm.tex"
REFERENCE = ROOT / "reference.pptx"
OUT = ROOT / "build/07-pagerank-algorithm.pptx"
ASSET_DIR = ROOT / "build/_pagerank_pptx_assets"
FLOW_DIAGRAM = ROOT / "diagrams/week7/week7_pagerank_iteration_flow.png"

EMU_PER_INCH = 914400

SLIDE_W = 13.333333333333
SLIDE_H = 7.5

RED = RGBColor(196, 111, 104)
BLACK = RGBColor(0, 0, 0)
MUTED = RGBColor(85, 85, 85)
BLUE = RGBColor(11, 79, 113)
LIGHT_BLUE = RGBColor(234, 245, 251)
GREEN = RGBColor(46, 125, 50)
LIGHT_GREEN = RGBColor(234, 246, 234)
ORANGE = RGBColor(184, 92, 0)
LIGHT_ORANGE = RGBColor(255, 243, 224)
GRAY = RGBColor(246, 248, 250)
MID_GRAY = RGBColor(144, 164, 174)
DARK_TEXT = RGBColor(38, 50, 56)
WHITE = RGBColor(255, 255, 255)


def clean_text(text):
    return text.replace("PageRank", "PageRank")


def add_blank(prs):
    for layout in prs.slide_layouts:
        if layout.name.lower() == "blank":
            return prs.slides.add_slide(layout)
    return prs.slides.add_slide(prs.slide_layouts[6])


def clear_template_slides(prs):
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def force_normal_open_view(path):
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        with zipfile.ZipFile(path) as zin:
            zin.extractall(tmp)

        view_props = tmp / "ppt" / "viewProps.xml"
        if view_props.exists():
            xml = view_props.read_text(encoding="utf-8")
            xml = xml.replace('lastView="sldMasterView"', 'lastView="sldView"')
            xml = xml.replace('lastView="sldLayoutView"', 'lastView="sldView"')
            view_props.write_text(xml, encoding="utf-8")

        out = path.with_suffix(".normal-view.pptx")
        with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in tmp.rglob("*"):
                if item.is_file():
                    zout.write(item, item.relative_to(tmp).as_posix())
        shutil.move(out, path)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def no_line(shape):
    try:
        shape.line.fill.background()
    except Exception:
        pass


def text_style(paragraph, size, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    paragraph.alignment = align
    paragraph.space_after = Pt(4)
    paragraph.line_spacing = 1.05
    for run in paragraph.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def set_text(shape, lines, sizes, bolds=None, aligns=None, color=BLACK):
    bolds = bolds or [False] * len(lines)
    aligns = aligns or [PP_ALIGN.LEFT] * len(lines)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = clean_text(line)
        text_style(p, sizes[i], bolds[i], color, aligns[i])


def add_textbox(slide, x, y, w, h, lines, sizes, bolds=None, aligns=None, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    no_line(box)
    set_text(box, lines, sizes, bolds, aligns, color)
    return box


def add_footer(slide, number):
    add_textbox(slide, 4.56, 6.95, 4.22, 0.28, ["Or Peretz"], [7.5], [False], [PP_ALIGN.CENTER], MUTED)
    add_textbox(slide, 9.56, 6.95, 3.11, 0.28, [str(number)], [10], [False], [PP_ALIGN.RIGHT], MUTED)


def add_title_frame(slide, title, number):
    add_textbox(
        slide,
        0.72,
        0.26,
        11.90,
        0.54,
        [title],
        [27 if len(title) > 38 else 30],
        [True],
        [PP_ALIGN.CENTER],
    )
    add_footer(slide, number)


def add_title_slide(prs):
    slide = add_blank(prs)
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(2.35), Inches(11.33), Inches(0.78)
    )
    set_fill(frame, WHITE)
    set_line(frame, RED, 1.0)
    add_textbox(slide, 1.1, 2.38, 11.1, 0.72, ["PageRank Algorithm"], [38], [True], [PP_ALIGN.CENTER])
    add_textbox(
        slide,
        2.0,
        3.28,
        9.33,
        0.62,
        ["From graphs to Markov chains to MapReduce"],
        [22],
        [False],
        [PP_ALIGN.CENTER],
        MUTED,
    )
    add_textbox(
        slide,
        2.1,
        4.30,
        9.1,
        0.52,
        ["Data Engineering Course"],
        [17],
        [False],
        [PP_ALIGN.CENTER],
        MUTED,
    )
    add_footer(slide, 1)


def add_section_slide(prs, title, subtitle, number):
    slide = add_blank(prs)
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(2.35), Inches(11.33), Inches(0.78)
    )
    set_fill(frame, WHITE)
    set_line(frame, RED, 1.0)
    add_textbox(slide, 1.1, 2.38, 11.1, 0.72, [title], [34], [True], [PP_ALIGN.CENTER])
    add_textbox(slide, 2.0, 3.25, 9.33, 0.65, [subtitle], [19], [False], [PP_ALIGN.CENTER], MUTED)
    add_footer(slide, number)


def add_bullets(slide, x, y, w, h, bullets, size=22, color=BLACK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    no_line(box)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"- {bullet}"
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
        p.line_spacing = 1.08
    return box


def add_callout(slide, x, y, w, h, title, body, line=BLUE, fill=LIGHT_BLUE):
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(box, fill)
    set_line(box, line, 1.4)
    lines = [title, body]
    set_text(box, lines, [16, 14], [True, False], [PP_ALIGN.LEFT, PP_ALIGN.LEFT], DARK_TEXT)
    return box


def add_node(slide, x, y, label, fill=LIGHT_BLUE, line=BLUE):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.72), Inches(0.72))
    set_fill(shp, fill)
    set_line(shp, line, 2.0)
    set_text(shp, [label], [17], [True], [PP_ALIGN.CENTER], DARK_TEXT)
    return shp


def add_arrow(slide, x1, y1, x2, y2, color=BLUE, width=2.0):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    try:
        conn.line.end_arrowhead = True
    except Exception:
        pass
    return conn


def add_graph(slide, nodes, edges):
    shapes = {}
    for name, (x, y, fill, line) in nodes.items():
        shapes[name] = add_node(slide, x, y, name, fill, line)
    for src, dst, label in edges:
        s = nodes[src]
        d = nodes[dst]
        x1, y1 = s[0] + 0.36, s[1] + 0.36
        x2, y2 = d[0] + 0.36, d[1] + 0.36
        add_arrow(slide, x1, y1, x2, y2)
        if label:
            lx, ly = (x1 + x2) / 2, (y1 + y2) / 2
            add_textbox(slide, lx - 0.35, ly - 0.20, 0.95, 0.28, [label], [10], [False], [PP_ALIGN.CENTER], MUTED)


def render_math(name, latex, fontsize=28, dpi=240):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / f"{name}.png"
    plt.rcParams["mathtext.fontset"] = "dejavusans"
    fig = plt.figure(figsize=(0.01, 0.01), dpi=dpi)
    fig.patch.set_alpha(0)
    text = fig.text(0, 0, latex, fontsize=fontsize, color="black")
    fig.canvas.draw()
    bbox = text.get_window_extent(renderer=fig.canvas.get_renderer()).expanded(1.08, 1.35)
    w, h = max(bbox.width / dpi, 0.2), max(bbox.height / dpi, 0.2)
    plt.close(fig)

    fig = plt.figure(figsize=(w, h), dpi=dpi)
    fig.patch.set_alpha(0)
    fig.text(0.02, 0.38, latex, fontsize=fontsize, color="black")
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)

    im = Image.open(path).convert("RGBA")
    bbox = im.getbbox()
    if bbox:
        im.crop(bbox).save(path)
    return path


def add_math(slide, latex, x, y, w, name, fontsize=28):
    path = render_math(name, latex, fontsize)
    im = Image.open(path)
    aspect = im.height / im.width
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(w * aspect))


def add_table(slide, x, y, w, h, data, font_size=13):
    rows = len(data)
    cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c])
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            for run in para.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(font_size if r else font_size + 1)
                run.font.bold = r == 0
                run.font.color.rgb = WHITE if r == 0 else BLACK
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if r % 2 else GRAY
    return table_shape


def add_formula_slide(
    prs,
    title,
    number,
    formula,
    bullets,
    formula_name,
    formula_width=8.0,
    formula_size=24,
    formula_y=1.78,
):
    slide = add_blank(prs)
    add_title_frame(slide, title, number)
    formula_x = (SLIDE_W - formula_width) / 2
    add_math(slide, formula, formula_x, formula_y, formula_width, formula_name, formula_size)
    add_bullets(slide, 1.25, 3.72, 10.8, 2.35, bullets, 19)
    return slide


def make_deck():
    if not TEX.exists():
        raise FileNotFoundError(TEX)
    if not REFERENCE.exists():
        raise FileNotFoundError(REFERENCE)

    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(REFERENCE)
    clear_template_slides(prs)
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_title_slide(prs)

    slide = add_blank(prs)
    add_title_frame(slide, "Learning Path", 2)
    add_bullets(
        slide,
        1.2,
        1.45,
        5.4,
        4.5,
        [
            "Start with graphs: nodes, edges, direction.",
            "Define the terms PageRank needs.",
            "Run the first algorithm by hand.",
            "Connect the algorithm to Markov chains.",
            "Improve it with damping.",
            "Scale one iteration with MapReduce.",
        ],
        21,
    )
    add_callout(slide, 7.05, 1.70, 4.65, 2.05, "Core idea", "A page is important if important pages link to it.")
    add_callout(slide, 7.05, 4.05, 4.65, 1.55, "Engineering idea", "Each iteration is a large distributed aggregation over graph edges.", GREEN, LIGHT_GREEN)

    add_section_slide(prs, "Part 1: Graphs", "The language PageRank is built on", 3)

    slide = add_blank(prs)
    add_title_frame(slide, "What Is a Graph?", 4)
    add_bullets(
        slide,
        1.0,
        1.45,
        5.6,
        3.2,
        [
            "A graph represents objects and relationships.",
            "Vertices / nodes are the objects.",
            "Edges are the relationships.",
            "Formally: G = (V, E).",
        ],
        22,
    )
    add_graph(
        slide,
        {
            "A": (7.0, 1.55, LIGHT_BLUE, BLUE),
            "B": (9.55, 1.00, LIGHT_BLUE, BLUE),
            "C": (9.55, 2.35, LIGHT_BLUE, BLUE),
            "D": (11.35, 1.70, LIGHT_BLUE, BLUE),
        },
        [("A", "B", ""), ("A", "C", ""), ("B", "D", ""), ("C", "D", "")],
    )
    add_callout(slide, 7.25, 4.35, 4.65, 0.95, "Web interpretation", "Nodes are pages; directed edges are hyperlinks.", GREEN, LIGHT_GREEN)

    slide = add_blank(prs)
    add_title_frame(slide, "Directed Web Graph", 5)
    add_bullets(
        slide,
        1.0,
        1.35,
        5.7,
        4.5,
        [
            "An undirected edge means mutual connection.",
            "A directed edge has a source and destination.",
            "A web link is directed: A -> B does not imply B -> A.",
            "PageRank operates on a directed graph.",
        ],
        21,
    )
    add_graph(
        slide,
        {
            "A": (7.2, 1.30, LIGHT_BLUE, BLUE),
            "B": (10.3, 1.30, LIGHT_BLUE, BLUE),
            "C": (8.75, 3.25, LIGHT_BLUE, BLUE),
        },
        [("A", "B", ""), ("A", "C", ""), ("B", "C", ""), ("C", "A", "")],
    )
    add_textbox(slide, 7.2, 5.05, 4.3, 0.45, ["Edges: A->B, A->C, B->C, C->A"], [14], [False], [PP_ALIGN.CENTER], MUTED)

    slide = add_blank(prs)
    add_title_frame(slide, "Terms Needed for PageRank", 6)
    add_table(
        slide,
        0.95,
        1.30,
        11.45,
        3.65,
        [
            ["Term", "Meaning"],
            ["Outlinks", "Pages that p links to"],
            ["Inlinks", "Pages that link to p"],
            ["Out-degree", "Number of outgoing links"],
            ["Gamma(p)", "The out-neighbor set of p"],
            ["Rank PR(p)", "Structural importance score of p"],
        ],
        15,
    )
    add_callout(slide, 1.1, 5.25, 11.1, 0.75, "Intuition", "A link from a strong page contributes more than a link from a weak page.")

    add_section_slide(prs, "Part 2: First Algorithm", "Rank flows through outgoing links", 7)

    add_formula_slide(
        prs,
        "Simple PageRank Update",
        8,
        r"$PR_{t+1}(p)=\sum_{q\rightarrow p}\frac{PR_t(q)}{|\Gamma(q)|}$",
        [
            "Every page distributes its current rank equally among outlinks.",
            "Every page sums the contributions it receives.",
            "Initialize all pages equally: PR0(p) = 1 / N.",
        ],
        "simple_update",
        6.9,
        22,
    )

    slide = add_blank(prs)
    add_title_frame(slide, "Simple Algorithm", 9)
    code = [
        "initialize PR[p] = 1 / N",
        "repeat until convergence:",
        "  newPR[p] = 0 for every page",
        "  for each page q:",
        "    contribution = PR[q] / outdegree(q)",
        "    for each outlink p of q:",
        "      newPR[p] += contribution",
        "  PR = newPR",
    ]
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.05), Inches(1.30), Inches(6.05), Inches(4.65))
    set_fill(box, GRAY)
    set_line(box, MID_GRAY, 1.0)
    set_text(box, code, [16] * len(code), [False] * len(code), [PP_ALIGN.LEFT] * len(code), BLACK)
    add_callout(slide, 7.55, 1.55, 4.2, 1.20, "Distribute", "Each source page sends rank to its out-neighbors.", BLUE, LIGHT_BLUE)
    add_callout(slide, 7.55, 3.10, 4.2, 1.20, "Collect", "Each destination page sums incoming rank.", GREEN, LIGHT_GREEN)
    add_callout(slide, 7.55, 4.65, 4.2, 1.20, "Repeat", "Continue until rank changes become small.", ORANGE, LIGHT_ORANGE)

    slide = add_blank(prs)
    add_title_frame(slide, "Numerical Example: Graph", 10)
    add_graph(
        slide,
        {
            "A": (2.05, 1.75, LIGHT_BLUE, BLUE),
            "B": (5.25, 1.75, LIGHT_BLUE, BLUE),
            "C": (3.65, 4.05, LIGHT_BLUE, BLUE),
        },
        [("A", "B", ""), ("A", "C", ""), ("B", "C", ""), ("C", "A", "")],
    )
    add_table(
        slide,
        7.35,
        1.48,
        4.5,
        2.2,
        [
            ["Page", "Outlinks", "Out-degree"],
            ["A", "[B, C]", "2"],
            ["B", "[C]", "1"],
            ["C", "[A]", "1"],
        ],
        14,
    )
    add_callout(slide, 7.35, 4.15, 4.5, 0.90, "Initialization", "N = 3, so every page starts with PR0 = 1/3.", GREEN, LIGHT_GREEN)

    slide = add_blank(prs)
    add_title_frame(slide, "Iteration 1: Distribute Rank", 11)
    add_table(
        slide,
        1.0,
        1.35,
        11.2,
        3.05,
        [
            ["Source", "Current rank", "Out-degree", "Contribution(s)"],
            ["A", "1/3", "2", "A->B = 1/6, A->C = 1/6"],
            ["B", "1/3", "1", "B->C = 1/3"],
            ["C", "1/3", "1", "C->A = 1/3"],
        ],
        14,
    )
    add_callout(slide, 1.35, 4.85, 10.3, 0.8, "Principle", "Outgoing rank is split equally across outgoing links.")

    slide = add_blank(prs)
    add_title_frame(slide, "Iteration 1: Collect Contributions", 12)
    add_table(
        slide,
        0.85,
        1.30,
        11.6,
        3.35,
        [
            ["Page", "Incoming contributions", "PR1"],
            ["A", "from C: 1/3", "1/3 = 0.3333"],
            ["B", "from A: 1/6", "1/6 = 0.1667"],
            ["C", "from A: 1/6, from B: 1/3", "1/2 = 0.5000"],
        ],
        15,
    )
    add_callout(slide, 1.2, 5.05, 10.9, 0.75, "Reading the result", "C becomes strong because both A and B point to it.", GREEN, LIGHT_GREEN)

    slide = add_blank(prs)
    add_title_frame(slide, "Iteration 2: Use PR1 as Input", 13)
    add_table(
        slide,
        0.85,
        1.25,
        11.6,
        3.9,
        [
            ["Page", "PR1", "Incoming contributions in iteration 2", "PR2"],
            ["A", "1/3", "from C: 1/2", "1/2 = 0.5000"],
            ["B", "1/6", "from A: 1/6", "1/6 = 0.1667"],
            ["C", "1/2", "from A: 1/6, from B: 1/6", "1/3 = 0.3333"],
        ],
        13,
    )
    add_callout(slide, 1.2, 5.45, 10.9, 0.65, "The rank keeps moving through the graph until it stabilizes.", "", ORANGE, LIGHT_ORANGE)

    add_section_slide(prs, "Part 3: Markov Chains", "The probability view of PageRank", 14)

    slide = add_blank(prs)
    add_title_frame(slide, "Random Surfer Interpretation", 15)
    add_bullets(
        slide,
        1.05,
        1.35,
        5.7,
        4.2,
        [
            "A user is currently on a page.",
            "The user chooses one outgoing link uniformly.",
            "The user moves to that linked page.",
            "After many steps, visit frequency becomes PageRank.",
        ],
        21,
    )
    add_graph(
        slide,
        {
            "A": (7.2, 1.35, LIGHT_BLUE, BLUE),
            "B": (10.1, 1.35, LIGHT_BLUE, BLUE),
            "C": (8.65, 3.50, LIGHT_BLUE, BLUE),
        },
        [("A", "B", "1/2"), ("A", "C", "1/2"), ("B", "C", "1"), ("C", "A", "1")],
    )

    slide = add_blank(prs)
    add_title_frame(slide, "Transition Matrix", 16)
    add_textbox(
        slide,
        0.95,
        1.35,
        5.9,
        1.45,
        ["M_ij = 1 / outdeg(j) if j -> i", "M_ij = 0 otherwise"],
        [22, 22],
        [False, False],
        [PP_ALIGN.CENTER, PP_ALIGN.CENTER],
    )
    add_textbox(
        slide,
        7.35,
        1.10,
        3.9,
        2.15,
        ["M =", "[ 0    0    1 ]", "[ 1/2  0    0 ]", "[ 1/2  1    0 ]"],
        [24, 23, 23, 23],
        [True, False, False, False],
        [PP_ALIGN.CENTER] * 4,
    )
    add_math(slide, r"$\mathbf{r}_{t+1}=M\mathbf{r}_t$", 4.50, 4.45, 4.25, "matrix_update", 32)
    add_callout(slide, 1.1, 5.45, 11.0, 0.65, "Meaning", "Rows are destinations, columns are sources. Each column distributes one page's rank.")

    slide = add_blank(prs)
    add_title_frame(slide, "Stationary Distribution", 17)
    add_math(slide, r"$\mathbf{r}=M\mathbf{r}$", 1.2, 1.45, 3.7, "stationary", 38)
    add_bullets(
        slide,
        5.2,
        1.35,
        6.7,
        3.6,
        [
            "At convergence, the rank vector stops changing.",
            "The stable vector is a stationary distribution.",
            "It is also an eigenvector of M with eigenvalue 1.",
        ],
        21,
    )
    add_callout(slide, 1.4, 4.85, 10.7, 0.95, "Two equivalent views", "PageRank is both iterative graph message passing and repeated sparse matrix-vector multiplication.", ORANGE, LIGHT_ORANGE)

    add_section_slide(prs, "Part 4: Damping", "Making PageRank robust on real web graphs", 18)

    slide = add_blank(prs)
    add_title_frame(slide, "Problems in the Simple Algorithm", 19)
    add_callout(slide, 0.95, 1.35, 3.55, 1.55, "Dangling node", "A page with no outlinks. Its rank has nowhere to go.", ORANGE, LIGHT_ORANGE)
    add_callout(slide, 4.90, 1.35, 3.55, 1.55, "Spider trap", "A group of pages can absorb rank forever.", ORANGE, LIGHT_ORANGE)
    add_callout(slide, 8.85, 1.35, 3.55, 1.55, "Disconnected graph", "The web is messy; not every page reaches every other page.", ORANGE, LIGHT_ORANGE)
    add_graph(
        slide,
        {
            "X": (5.35, 3.65, LIGHT_ORANGE, ORANGE),
            "Y": (7.35, 3.65, LIGHT_ORANGE, ORANGE),
            "Z": (6.35, 5.10, LIGHT_ORANGE, ORANGE),
        },
        [("X", "Y", ""), ("Y", "X", ""), ("Y", "Z", ""), ("Z", "Y", "")],
    )

    slide = add_blank(prs)
    add_title_frame(slide, "What Is Damping?", 20)
    add_bullets(
        slide,
        1.05,
        1.35,
        5.55,
        3.6,
        [
            "With probability d, follow an outgoing link.",
            "With probability 1-d, teleport to a random page.",
            "Common choice: d = 0.85.",
            "Teleportation prevents permanent trapping.",
        ],
        21,
    )
    add_callout(slide, 7.0, 1.55, 4.55, 1.20, "85%", "Follow a hyperlink", BLUE, LIGHT_BLUE)
    add_callout(slide, 7.0, 3.15, 4.55, 1.20, "15%", "Jump to any page uniformly", GREEN, LIGHT_GREEN)
    add_callout(slide, 7.0, 4.75, 4.55, 0.95, "Effect", "Every page remains reachable.", ORANGE, LIGHT_ORANGE)

    add_formula_slide(
        prs,
        "Damped PageRank Formula",
        21,
        r"$PR_{t+1}(p)=\frac{1-d}{N}+d\left(\sum_{q\rightarrow p}\frac{PR_t(q)}{|\Gamma(q)|}+\frac{M_t}{N}\right)$",
        [
            "(1-d)/N is teleportation mass.",
            "The sum is normal link-based contribution.",
            "M_t/N redistributes dangling-node mass.",
        ],
        "damped_formula",
        8.0,
        22,
    )

    slide = add_blank(prs)
    add_title_frame(slide, "Damped Numerical Example", 22)
    add_graph(
        slide,
        {
            "A": (1.4, 1.50, LIGHT_BLUE, BLUE),
            "B": (4.05, 1.50, LIGHT_BLUE, BLUE),
            "C": (2.72, 3.65, LIGHT_BLUE, BLUE),
            "D": (6.30, 2.55, LIGHT_ORANGE, ORANGE),
        },
        [("A", "B", ""), ("A", "C", ""), ("B", "C", ""), ("C", "A", "")],
    )
    add_textbox(slide, 6.05, 3.45, 1.35, 0.30, ["dangling"], [12], [False], [PP_ALIGN.CENTER], MUTED)
    add_bullets(
        slide,
        7.75,
        1.35,
        4.4,
        2.5,
        [
            "N = 4",
            "PR0 for each page = 0.25",
            "d = 0.85",
            "M0 = PR0(D) = 0.25",
        ],
        19,
    )
    add_math(slide, r"$\frac{M_0}{N}=0.0625,\quad \frac{1-d}{N}=0.0375$", 7.55, 4.25, 4.65, "damped_terms", 25)

    slide = add_blank(prs)
    add_title_frame(slide, "Damped Example: Output", 23)
    add_table(
        slide,
        0.8,
        1.25,
        11.8,
        3.8,
        [
            ["Page", "Link contribution", "Dangling share", "Formula", "PR1"],
            ["A", "0.2500", "0.0625", "0.0375 + 0.85(0.3125)", "0.3031"],
            ["B", "0.1250", "0.0625", "0.0375 + 0.85(0.1875)", "0.1969"],
            ["C", "0.3750", "0.0625", "0.0375 + 0.85(0.4375)", "0.4094"],
            ["D", "0.0000", "0.0625", "0.0375 + 0.85(0.0625)", "0.0906"],
        ],
        12,
    )
    add_callout(slide, 1.05, 5.45, 11.1, 0.65, "Check", "The ranks still sum to 1.0000, so no rank mass leaked.", GREEN, LIGHT_GREEN)

    add_section_slide(prs, "Part 5: MapReduce", "Scaling one PageRank iteration", 24)

    slide = add_blank(prs)
    add_title_frame(slide, "Why PageRank Needs Distributed Computing", 25)
    add_bullets(
        slide,
        1.05,
        1.35,
        5.9,
        4.4,
        [
            "The web graph can contain billions of pages.",
            "The edge list can contain many billions of links.",
            "Each iteration scans the graph.",
            "Each iteration emits one contribution per edge.",
            "This is naturally a distributed aggregation problem.",
        ],
        21,
    )
    add_callout(slide, 7.35, 1.55, 4.65, 1.05, "Record format", "(page, current rank, adjacency list)", BLUE, LIGHT_BLUE)
    add_table(
        slide,
        7.25,
        3.0,
        4.85,
        2.35,
        [
            ["Page", "Rank", "Outlinks"],
            ["A", "0.25", "[B, C]"],
            ["B", "0.25", "[C]"],
            ["C", "0.25", "[A]"],
            ["D", "0.25", "[]"],
        ],
        12,
    )

    slide = add_blank(prs)
    add_title_frame(slide, "MapReduce Plan", 26)
    add_bullets(
        slide,
        1.05,
        1.30,
        10.8,
        4.8,
        [
            "Map: emit rank contributions to destination pages.",
            "Map: also emit the adjacency list to preserve graph structure.",
            "Shuffle: group all values by destination page.",
            "Reduce: sum contributions and apply the damped formula.",
            "Repeat: feed the output into the next iteration.",
        ],
        22,
    )

    slide = add_blank(prs)
    add_title_frame(slide, "PageRank MapReduce Flow", 27)
    if FLOW_DIAGRAM.exists():
        slide.shapes.add_picture(str(FLOW_DIAGRAM), Inches(0.85), Inches(1.18), width=Inches(11.65))
    else:
        add_callout(slide, 1.1, 2.5, 11.0, 1.0, "Missing figure", str(FLOW_DIAGRAM), ORANGE, LIGHT_ORANGE)

    slide = add_blank(prs)
    add_title_frame(slide, "Mapper Logic", 28)
    code = [
        "map(page p, rank r, adjacency list L):",
        "  emit(p, STRUCTURE(L))",
        "  if L is empty:",
        "    add r to DANGLING_MASS",
        "  else:",
        "    c = r / length(L)",
        "    for each destination q in L:",
        "      emit(q, CONTRIBUTION(c))",
    ]
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), Inches(1.25), Inches(6.25), Inches(4.85))
    set_fill(box, GRAY)
    set_line(box, MID_GRAY, 1.0)
    set_text(box, code, [15] * len(code), [False] * len(code), [PP_ALIGN.LEFT] * len(code), BLACK)
    add_callout(slide, 7.65, 1.55, 4.35, 1.25, "Contribution", "Rank mass is sent to outlinks.", BLUE, LIGHT_BLUE)
    add_callout(slide, 7.65, 3.20, 4.35, 1.25, "Structure", "Adjacency list must survive for the next iteration.", GREEN, LIGHT_GREEN)
    add_callout(slide, 7.65, 4.85, 4.35, 1.00, "Dangling", "No outlinks means no emitted contribution.", ORANGE, LIGHT_ORANGE)

    slide = add_blank(prs)
    add_title_frame(slide, "Shuffle and Reducer Logic", 29)
    add_table(
        slide,
        0.85,
        1.25,
        5.65,
        2.75,
        [
            ["Key", "Grouped values"],
            ["A", "0.250 + structure [B,C]"],
            ["B", "0.125 + structure [C]"],
            ["C", "0.125, 0.250 + structure [A]"],
            ["D", "structure []"],
        ],
        11,
    )
    code = [
        "reduce(page p, values):",
        "  L = adjacency list",
        "  s = sum contribution values",
        "  newRank = (1-d)/N + d*(s + M/N)",
        "  emit(p, newRank, L)",
    ]
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.0), Inches(1.35), Inches(5.1), Inches(3.0))
    set_fill(box, GRAY)
    set_line(box, MID_GRAY, 1.0)
    set_text(box, code, [14] * len(code), [False] * len(code), [PP_ALIGN.LEFT] * len(code), BLACK)
    add_callout(slide, 1.05, 4.75, 11.05, 0.75, "Reducer responsibility", "Separate numeric contributions from graph structure; then emit updated rank plus adjacency list.", GREEN, LIGHT_GREEN)

    slide = add_blank(prs)
    add_title_frame(slide, "Step-by-Step Map Output", 30)
    add_table(
        slide,
        0.8,
        1.16,
        11.8,
        4.85,
        [
            ["Source", "Emitted key", "Emitted value"],
            ["A", "B", "contribution 0.125"],
            ["A", "C", "contribution 0.125"],
            ["A", "A", "structure [B, C]"],
            ["B", "C", "contribution 0.250"],
            ["B", "B", "structure [C]"],
            ["C", "A", "contribution 0.250"],
            ["C", "C", "structure [A]"],
            ["D", "D", "structure []"],
            ["D", "global", "dangling mass 0.250"],
        ],
        11,
    )

    slide = add_blank(prs)
    add_title_frame(slide, "Reducer Output", 31)
    add_table(
        slide,
        1.05,
        1.35,
        10.9,
        2.65,
        [
            ["Page", "New rank", "Preserved structure"],
            ["A", "0.3031", "[B, C]"],
            ["B", "0.1969", "[C]"],
            ["C", "0.4094", "[A]"],
            ["D", "0.0906", "[]"],
        ],
        15,
    )
    add_callout(slide, 1.15, 4.55, 10.7, 0.90, "One iteration complete", "The next job uses these ranks as input and repeats the same pattern.", BLUE, LIGHT_BLUE)

    slide = add_blank(prs)
    add_title_frame(slide, "Engineering Notes", 32)
    add_callout(slide, 0.95, 1.25, 3.55, 1.35, "Combiner", "Safe for summing numeric contributions; do not drop structure.", BLUE, LIGHT_BLUE)
    add_callout(slide, 4.90, 1.25, 3.55, 1.35, "Dangling mass", "Track with a counter, special key, or auxiliary aggregation.", ORANGE, LIGHT_ORANGE)
    add_callout(slide, 8.85, 1.25, 3.55, 1.35, "Convergence", "Aggregate |PR(t+1)-PR(t)| across pages.", GREEN, LIGHT_GREEN)
    add_callout(slide, 0.95, 3.35, 3.55, 1.35, "Cost", "Each iteration shuffles O(|E|) contributions.", BLUE, LIGHT_BLUE)
    add_callout(slide, 4.90, 3.35, 3.55, 1.35, "Skew", "Popular pages can overload a reducer.", ORANGE, LIGHT_ORANGE)
    add_callout(slide, 8.85, 3.35, 3.55, 1.35, "Mitigation", "Use partial aggregation for hot keys.", GREEN, LIGHT_GREEN)

    slide = add_blank(prs)
    add_title_frame(slide, "Summary", 33)
    add_bullets(
        slide,
        1.05,
        1.35,
        10.9,
        4.8,
        [
            "A web graph is a directed graph of pages and hyperlinks.",
            "Basic PageRank sends rank through outgoing links.",
            "The Markov-chain view treats ranks as long-run visit probabilities.",
            "Damping adds teleportation and handles traps/dangling nodes.",
            "MapReduce implements PageRank as distributed contribution aggregation.",
        ],
        22,
    )

    slide = add_blank(prs)
    add_title_frame(slide, "Practice Questions", 34)
    add_bullets(
        slide,
        0.95,
        1.25,
        11.6,
        5.3,
        [
            "Compute one simple PageRank iteration for a four-page graph.",
            "Write the transition matrix for the graph.",
            "Add damping with d = 0.85 and compare the result.",
            "Explain why the adjacency list must be preserved in MapReduce.",
            "Identify the bottleneck when one page receives many links.",
        ],
        20,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    force_normal_open_view(OUT)
    return OUT


if __name__ == "__main__":
    out = make_deck()
    print(out)
