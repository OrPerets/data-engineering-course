from copy import deepcopy

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


SRC = "build/DWH-Schemas.pptx"
OUT = "build/DWH-Schemas-refined.pptx"


def set_text(shape, text):
    """Replace text while keeping the existing placeholder/shape styling."""
    shape.text = text


def textbox(slide, text, x, y, w, h, size=18, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for idx, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0, 0, 0)
    return box


def footer_textbox(slide, text, x, y, w, h, align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.name = "Calibri"
        run.font.size = Pt(7)
        run.font.color.rgb = RGBColor(128, 128, 128)
    return box


def set_font_size(shape, size):
    if not hasattr(shape, "text_frame"):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)


def insert_title_content_slide(prs, index, title, body):
    """Create a course-template title/content slide and insert it at index."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    for shape in slide.shapes:
        if shape.has_text_frame and shape != slide.shapes.title and shape.text == "":
            # The first large empty placeholder is the content body.
            if shape.width > Inches(5) and shape.height > Inches(2):
                shape.text = body
                break

    sld_id_lst = prs.slides._sldIdLst
    new_sld_id = sld_id_lst[-1]
    sld_id_lst.remove(new_sld_id)
    sld_id_lst.insert(index, new_sld_id)
    footer_textbox(slide, "Or Peretz", 4.56, 6.95, 4.22, 0.4)
    footer_textbox(slide, "0", 9.56, 6.95, 3.11, 0.4, PP_ALIGN.RIGHT)
    return slide


def renumber_slides(prs):
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip().isdigit():
                shape.text = str(idx)


prs = Presentation(SRC)


# 1
set_text(prs.slides[0].shapes[0], "DWH Schema Design (Part 2)")

# 2
set_text(prs.slides[1].shapes[0], "Learning Objectives")
set_text(
    prs.slides[1].shapes[1],
    "Declare fact grain, keys, measures, and aggregation rules\n"
    "Choose star, snowflake, or galaxy schemas by workload\n"
    "Model history, table behavior, and conformed dimensions\n"
    "Design partitioning and governance controls for production use",
)

# 3
set_text(prs.slides[2].shapes[0], "Why Schema Design Matters")
set_text(
    prs.slides[2].shapes[1],
    "Schema design is where business meaning becomes executable\n"
    "Wrong grain or keys silently corrupt KPIs and downstream models\n"
    "Physical choices such as partitioning determine cost and latency\n"
    "Professional DWH design balances correctness, performance, and operability",
)

# 4
set_text(prs.slides[3].shapes[0], "Fact and Dimension Contracts")
set_text(
    prs.slides[3].shapes[1],
    "Fact table: measurable business events at one declared grain\n"
    "Dimension table: descriptive context used to filter, group, and explain facts\n"
    "Measure: numeric value with explicit aggregation semantics\n"
    "Hierarchy: governed drill path such as day -> month -> quarter -> year",
)

# 5
set_text(prs.slides[4].shapes[0], "The Multi-Dimensional Model")

# 6
set_text(prs.slides[5].shapes[0], "Dimensional Modeling")
set_text(
    prs.slides[5].shapes[1],
    "Dimensions organize attributes into stable analysis paths\n"
    "E.g., Time: day -> week -> month -> quarter -> year\n"
    "E.g., Product: SKU -> product line -> brand -> category\n"
    "Conformed dimensions reuse the same keys and labels across facts\n"
    "Role-playing dimensions reuse one table in different roles, e.g., order_date and ship_date",
)

# 7
set_text(prs.slides[6].shapes[0], "Concept Hierarchy: Location Dimension")

# 8
set_text(prs.slides[7].shapes[0], "Declaring Fact Grain")
set_text(
    prs.slides[7].shapes[1],
    "Grain defines exactly what one row in the fact table represents\n"
    "Example: one row per sold order line, not one row per order\n"
    "The grain determines valid joins, uniqueness checks, and aggregations\n"
    "If the grain is ambiguous, double counting is likely\n"
    "Document grain and natural uniqueness before writing ETL",
)

insert_title_content_slide(
    prs,
    8,
    "Grain: What It Is",
    "The grain is the business event level captured by the fact table\n"
    "It answers the question: what does one row mean?\n"
    "Grain is not the table name, source file, or dashboard metric\n"
    "Once declared, every key and measure must match that row meaning\n"
    "Changing grain later usually breaks joins, aggregates, and history",
)

insert_title_content_slide(
    prs,
    9,
    "Common Grain Options",
    "Transaction grain: one row per event, e.g., order line, payment, click\n"
    "Periodic snapshot grain: one row per entity per period, e.g., daily inventory\n"
    "Accumulating snapshot grain: one row per process lifecycle, e.g., fulfillment case\n"
    "Aggregate grain: pre-summarized rows, e.g., product-day sales totals\n"
    "Factless grain: event existence only, e.g., attendance or promotion eligibility",
)

insert_title_content_slide(
    prs,
    10,
    "Choosing the Right Grain",
    "Start with the lowest atomic grain needed for audit and flexible analysis\n"
    "Use aggregate grains only for performance marts, not as the only source of truth\n"
    "Avoid mixing grains in one fact table, such as order header and order line\n"
    "Check that the proposed grain has a stable unique key\n"
    "Write the grain statement before designing dimensions or ETL logic",
)

insert_title_content_slide(
    prs,
    11,
    "Table Behavior Types",
    "Table behavior describes how rows are written and changed over time\n"
    "Append-only tables only insert new rows and preserve full event history\n"
    "Mutable tables allow updates and deletes when the latest state matters\n"
    "Late-arriving data reaches the warehouse after the business event happened\n"
    "These choices affect ETL, audits, reconciliation, and downstream trust",
)

insert_title_content_slide(
    prs,
    12,
    "Append-Only vs. Mutable Tables",
    "Append-only: best for logs, transactions, CDC history, and replayability\n"
    "Mutable: best for current-state dimensions, operational mirrors, and corrections\n"
    "Append-only simplifies audit and backfill, but consumers may need latest-state views\n"
    "Mutable tables are easier for current-state queries, but they can hide history\n"
    "A common pattern is append-only raw layer plus curated mutable serving tables",
)

insert_title_content_slide(
    prs,
    13,
    "Late-Arriving Data",
    "Late-arriving fact: event_time is earlier than load_time\n"
    "Late-arriving dimension: the fact arrives before the matching dimension row\n"
    "Do not drop these records; stage, track, and reconcile them\n"
    "Keep event_time separate from load_time\n"
    "Define unknown-key, backfill, and restatement policy",
)

# 9
set_text(prs.slides[14].shapes[0], "Measure Types and Aggregation Rules")
set_text(
    prs.slides[14].shapes[1],
    "Additive: revenue, quantity - safe to sum across all dimensions\n"
    "Semi-additive: inventory_balance - sum across product, not across time\n"
    "Non-additive: conversion_rate - calculate from numerator and denominator\n"
    "Derived metrics should be certified in the semantic layer\n"
    "Store atomic components when possible; avoid averages of averages",
)

# 10
set_text(prs.slides[15].shapes[0], "Natural, Durable, and Surrogate Keys")
set_text(
    prs.slides[15].shapes[1],
    "Natural key: source identifier such as customer_id\n"
    "Durable key: stable identity across systems and merges\n"
    "Surrogate key: warehouse-managed key used in fact joins\n"
    "Facts should join dimensions by surrogate key\n"
    "Late-arriving dimensions need an unknown-key policy",
)

# 11
set_text(prs.slides[16].shapes[0], "Slowly Changing Dimensions (SCD)")
set_text(
    prs.slides[16].shapes[1],
    "Type 0: preserve original value, no changes after insert\n"
    "Type 1: overwrite old value when history is not required\n"
    "Type 2: keep historical versions with validity window and current flag\n"
    "Type 3: keep limited previous value in extra columns\n"
    "Choose SCD behavior by audit, finance, and analytical truth requirements",
)

# 12
set_text(prs.slides[17].shapes[0], "SCD Type 2 Example: Customer Region Change")
set_text(
    prs.slides[17].shapes[1],
    "Customer moves from North to Center\n"
    "Type 1: all historical sales are reclassified as Center\n"
    "Type 2: past sales remain North, new sales use Center\n"
    "Fact rows join to the dimension version valid on the sale date\n"
    "Finance and compliance usually require Type 2 for material attributes",
)
set_text(prs.slides[17].shapes[3], "SCD Type 2: versioned dimension rows")

# 13
set_text(prs.slides[18].shapes[0], "Star Schema: Default BI Design")

# 14
set_text(prs.slides[19].shapes[2], "The Classic Star Schema")
set_text(
    prs.slides[19].shapes[3],
    "A fact table sits at the center and connects to denormalized dimensions.\n"
    "Each dimension key appears once in the fact table primary/unique key.\n"
    "Dimension tables contain business labels, attributes, and hierarchies.\n"
    "Star schemas are usually easiest for BI tools, analysts, and query optimizers.\n"
    "Default to star unless normalization has a measurable benefit.",
)

# 15
set_text(prs.slides[20].shapes[0], "Snowflake Schema")
set_text(
    prs.slides[20].shapes[1],
    "Snowflaking normalizes selected dimension attributes into sub-dimensions\n"
    "It can reduce redundancy in large, shared, hierarchical dimensions\n"
    "It adds joins, ETL complexity, and harder self-service BI\n"
    "Use selectively when maintenance or consistency benefits outweigh query cost\n"
    "Common compromise: normalized core plus star marts",
)

# 16
set_text(prs.slides[21].shapes[0], "Snowflake Schema: Normalized Dimensions")

# 17
set_text(prs.slides[22].shapes[0], "Snowflake Example: Product and Geography")

# 18
set_text(prs.slides[23].shapes[0], "Snowflaking a Product Dimension")
set_text(
    prs.slides[23].shapes[1],
    "Split shared hierarchy attributes into reference tables\n"
    "Product stores brand_key and category_key\n"
    "Benefit: less repetition and stronger governance\n"
    "Cost: extra joins unless a mart denormalizes back",
)
set_text(prs.slides[23].shapes[2], "Product key\nProduct name\nProduct code\nBrand key")
set_text(prs.slides[23].shapes[3], "Brand key\nBrand name\nCategory key")
set_text(prs.slides[23].shapes[4], "Category key\nProduct category")
set_font_size(prs.slides[23].shapes[2], 14)
set_font_size(prs.slides[23].shapes[3], 14)
set_font_size(prs.slides[23].shapes[4], 14)

# 19
set_text(prs.slides[24].shapes[0], "Incident: Full-Scan Dashboard Outage")
set_text(
    prs.slides[24].shapes[1],
    "Dashboard refreshed every 5 minutes on a 1 TB atomic fact table\n"
    "A missing date predicate forced full scans on every refresh\n"
    "Concurrent scans saturated warehouse slots and delayed other pipelines\n"
    "Root cause: no model-level or semantic-layer guardrail for partition filters\n"
    "Fix: enforce required time predicates for large facts and monitor partitions scanned",
)

# 20
set_text(prs.slides[25].shapes[0], "Partitioning and Clustering Strategy")
set_text(
    prs.slides[25].shapes[1],
    "Partition large facts by the dominant time filter, usually date_key\n"
    "Align partition granularity with query patterns and retention policy\n"
    "Cluster or sort on common filters such as customer_key or product_key\n"
    "Avoid tiny partitions and small files; they add metadata overhead\n"
    "Validate design with bytes scanned, partitions read, and p95 latency",
)

# 21: keep original formula/content objects to preserve the lecture style.
set_text(prs.slides[26].shapes[0], "Partition Pruning Cost Model")

# 22
set_text(prs.slides[27].shapes[0], "Join Cost Intuition")

# 23
set_text(prs.slides[28].shapes[0], "Architecture Evolution: From Model to Contract")
set_text(
    prs.slides[28].shapes[1],
    "v1: one large table and manual metric definitions\n"
    "v2: star schema, conformed dimensions, and semantic contracts\n"
    "Add SCD2, bridges, and snowflaking only when needed\n"
    "Migrate with compatibility views to avoid breaking dashboards\n"
    "Version, test, document, and monitor schema changes",
)
set_text(prs.slides[28].shapes[3], "Architecture evolution")

# 24
set_text(prs.slides[29].shapes[0], "Governance Controls for Large Facts")
set_text(
    prs.slides[29].shapes[1],
    "SQL guardrail: block large-fact queries without required time filters\n"
    "Semantic layer: certified metric formulas, grain, filters, and owners\n"
    "Data contracts: required keys, null policies, and uniqueness checks\n"
    "Quality tests: freshness, duplicates, and reconciliation\n"
    "Cost monitoring: bytes scanned, slots consumed, and query p95",
)

# 25
set_text(prs.slides[30].shapes[0], "Operational Metrics and Review Checklist")
set_text(
    prs.slides[30].shapes[1],
    "Bytes scanned and partitions read per dashboard refresh\n"
    "Query latency p50/p95, queue time, and spill indicators\n"
    "Duplicate fact rows, unmatched dimension keys, and unknown-key rate\n"
    "Freshness lag, late-arriving volume, and reconciliation differences\n"
    "Compaction backlog, small-file count, and growth of the largest facts",
)

renumber_slides(prs)
prs.save(OUT)
print(OUT)
