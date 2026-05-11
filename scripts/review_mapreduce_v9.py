from pathlib import Path

import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "build/06-map-reduce.applied.v8.pptx"
OUT = ROOT / "build/06-map-reduce.applied.v9-reviewed.pptx"
ASSET_DIR = ROOT / "build/_pptx_check/v9_assets"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BLUE = RGBColor(49, 116, 181)
DARK_BLUE = RGBColor(35, 76, 119)
MUTED = RGBColor(85, 85, 85)
BLACK = RGBColor(0, 0, 0)


def remove_shape(shape):
    shape._element.getparent().remove(shape._element)


def no_outline(shape):
    try:
        shape.line.fill.background()
    except Exception:
        pass


def set_run_style(run, size=None, bold=None, color=BLACK, name="Calibri"):
    run.font.name = name
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    run.font.color.rgb = color


def style_text_frame(shape, title_size=32, body_size=22, sentence_size=None, subtitle_align=PP_ALIGN.LEFT):
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.word_wrap = True
    for pi, p in enumerate(tf.paragraphs):
        if not p.text.strip():
            continue
        if pi == 0:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                set_run_style(run, title_size, True)
        else:
            p.alignment = subtitle_align
            for run in p.runs:
                set_run_style(run, sentence_size or body_size, False)


def set_box_text(shape, paragraphs, sizes, bolds=None, alignments=None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    bolds = bolds or [False] * len(paragraphs)
    alignments = alignments or [PP_ALIGN.LEFT] * len(paragraphs)
    for i, text in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = alignments[i]
        p.space_after = Pt(5 if i else 7)
        for run in p.runs:
            set_run_style(run, sizes[i], bolds[i])


def add_note(slide, text, x, y, w, h, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for r in p.runs:
        set_run_style(r, size, False, MUTED)
    no_outline(box)
    return box


def render_math(name, latex, fontsize=30, dpi=240):
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / f"{name}.png"
    plt.rcParams["mathtext.fontset"] = "dejavusans"
    fig = plt.figure(figsize=(0.01, 0.01), dpi=dpi)
    fig.patch.set_alpha(0)
    text = fig.text(0, 0, latex, fontsize=fontsize, color="black")
    fig.canvas.draw()
    bbox = text.get_window_extent(renderer=fig.canvas.get_renderer()).expanded(1.10, 1.45)
    w, h = bbox.width / dpi, bbox.height / dpi
    plt.close(fig)

    fig = plt.figure(figsize=(w, h), dpi=dpi)
    fig.patch.set_alpha(0)
    fig.text(0.02, 0.42, latex, fontsize=fontsize, color="black")
    fig.savefig(path, transparent=True, bbox_inches="tight", pad_inches=0.04)
    plt.close(fig)

    # Trim transparent padding so PowerPoint sizing is predictable.
    im = Image.open(path).convert("RGBA")
    bbox = im.getbbox()
    if bbox:
        im.crop(bbox).save(path)
    return path


def add_math_line(slide, path, x, y, w):
    im = Image.open(path)
    aspect = im.height / im.width
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(w * aspect))


def remove_off_canvas(slide):
    for shape in list(slide.shapes):
        if shape.left >= SLIDE_W or shape.top >= SLIDE_H:
            remove_shape(shape)


def remove_pictures(slide, predicate):
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and predicate(shape):
            remove_shape(shape)


def normalize_borders_and_text(prs):
    for si, slide in enumerate(prs.slides, 1):
        remove_off_canvas(slide)
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.shape_type in {
                MSO_SHAPE_TYPE.PLACEHOLDER,
                MSO_SHAPE_TYPE.TEXT_BOX,
            }:
                no_outline(shape)

        # Normalize main title text without disturbing diagram labels.
        first = slide.shapes[0] if slide.shapes else None
        if first is not None and getattr(first, "has_text_frame", False):
            if si == 1:
                style_text_frame(first, title_size=38, body_size=20)
            elif si in {10, 11, 12, 13, 15, 16, 17, 18, 19}:
                style_text_frame(first, title_size=32, body_size=18, sentence_size=18, subtitle_align=PP_ALIGN.CENTER)
            else:
                style_text_frame(first, title_size=32, body_size=22)

        # Full-text slides had very large body text; keep it readable but consistent.
        if si in {2, 3, 7, 14, 21, 23, 24, 26} and len(slide.shapes) > 1:
            body = slide.shapes[1]
            if getattr(body, "has_text_frame", False):
                for p in body.text_frame.paragraphs:
                    for run in p.runs:
                        set_run_style(run, 24, False)


def fix_formal_model(slide):
    remove_pictures(slide, lambda s: True)
    title_box = slide.shapes[0]
    title_box.left, title_box.top, title_box.width, title_box.height = Inches(0.72), Inches(0.48), Inches(4.05), Inches(4.10)
    set_box_text(
        title_box,
        [
            "Formal Model",
            "D is the input multiset.",
            "Map emits intermediate pairs.",
            "Reduce receives one complete key group.",
            "K₂ is the grouping key.",
        ],
        [32, 18.5, 18.5, 18.5, 18.5],
        [True, False, False, False, False],
        [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )

    equations = [
        ("model_input", r"$D \in \mathcal{M}(K_1 \times V_1)$", 32, 5.15, 1.05, 4.15),
        ("model_map", r"$m:K_1 \times V_1 \rightarrow \mathcal{M}(K_2 \times V_2)$", 32, 5.15, 2.20, 5.35),
        ("model_reduce", r"$\rho:K_2 \times \mathcal{M}(V_2) \rightarrow \mathcal{M}(K_3 \times V_3)$", 32, 5.15, 3.35, 5.80),
    ]
    for name, latex, size, x, y, w in equations:
        add_math_line(slide, render_math(name, latex, size), x, y, w)
    add_note(slide, "M(X) denotes a multiset over domain X.", 5.17, 4.45, 5.5, 0.35, 15)


def fix_formal_algorithm(slide):
    remove_pictures(slide, lambda s: True)
    left = slide.shapes[0]
    left.left, left.top, left.width, left.height = Inches(0.76), Inches(0.48), Inches(4.05), Inches(5.05)
    set_box_text(
        left,
        [
            "Formal Algorithm",
            "1. Map all records.",
            "2. Partition by key.",
            "3. Group values by key.",
            "4. Reduce each group.",
            "Invariant: same key -> complete group.",
        ],
        [32, 18.2, 18.2, 18.2, 18.2, 16.5],
        [True, False, False, False, False, False],
        [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    equations = [
        ("alg_map", r"$I=\mathrm{merge}\{m(k_1,v_1):(k_1,v_1)\in D\}$", 29, 5.10, 1.05, 5.95),
        ("alg_partition", r"$p(k_2)=h(k_2)\ \mathrm{mod}\ R$", 29, 5.00, 2.05, 3.35),
        ("alg_group", r"$G(k_2)=\{v_2:(k_2,v_2)\in I\}$", 29, 5.00, 3.05, 4.95),
        ("alg_reduce", r"$O=\mathrm{merge}\{\rho(k_2,G(k_2)):k_2\in\mathrm{keys}(I)\}$", 29, 5.10, 4.05, 6.15),
    ]
    for name, latex, size, x, y, w in equations:
        add_math_line(slide, render_math(name, latex, size), x, y, w)
    add_note(slide, "merge = multiset union/concatenation of emitted pairs.", 5.03, 5.05, 6.6, 0.35, 14.5)


def fix_shuffle_cost(slide):
    remove_pictures(slide, lambda s: s.left < Inches(6.5))
    # Enlarge and center the existing phase diagram while keeping it clear of the right arc.
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape.left, shape.top, shape.width = Inches(7.00), Inches(1.10), Inches(3.05)
            shape.height = Inches(4.85)

    title = slide.shapes[0]
    title.left, title.top, title.width, title.height = Inches(0.75), Inches(0.35), Inches(11.0), Inches(0.58)
    set_box_text(title, ["Shuffle Cost and Runtime"], [32], [True], [PP_ALIGN.CENTER])

    body = slide.shapes[1]
    body.left, body.top, body.width, body.height = Inches(1.05), Inches(2.55), Inches(5.10), Inches(2.65)
    set_box_text(
        body,
        [
            "E = emitted intermediate pairs",
            "s = serialized pair size",
            "Use combiners to reduce E.",
            "Use compact encoding to reduce s.",
            "Slowest reducer bounds completion.",
        ],
        [19, 19, 19, 19, 19],
        [False] * 5,
        [PP_ALIGN.LEFT] * 5,
    )

    equations = [
        ("shuffle_cost", r"$C_{\mathrm{shuffle}}=E\cdot s$", 28, 1.05, 1.28, 3.15),
        ("runtime_total", r"$T_{\mathrm{total}}=T_{\mathrm{map}}+T_{\mathrm{shuffle}}+T_{\mathrm{reduce}}$", 28, 1.05, 1.85, 4.85),
    ]
    for name, latex, size, x, y, w in equations:
        add_math_line(slide, render_math(name, latex, size), x, y, w)


def fix_join_slide(slide):
    title = slide.shapes[0]
    title.left, title.top, title.width, title.height = Inches(0.75), Inches(0.35), Inches(11.0), Inches(0.58)
    set_box_text(title, ["Joins in MapReduce"], [32], [True], [PP_ALIGN.CENTER])

    body = slide.shapes[1]
    body.left, body.top, body.width, body.height = Inches(0.78), Inches(1.33), Inches(3.65), Inches(3.5)
    set_box_text(
        body,
        [
            "Reduce-side: shuffle both tables by join key.",
            "Broadcast: copy the small table to each mapper.",
            "Decision: move both tables or replicate the small side.",
        ],
        [20, 20, 20],
        [False, False, False],
        [PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape.left, shape.top, shape.width, shape.height = Inches(4.65), Inches(1.18), Inches(6.85), Inches(4.05)


def fix_skew_slide(slide):
    title = slide.shapes[0]
    title.left, title.top, title.width, title.height = Inches(0.75), Inches(0.35), Inches(11.0), Inches(0.58)
    set_box_text(title, ["Failure Mode: Data Skew"], [32], [True], [PP_ALIGN.CENTER])

    body = slide.shapes[1]
    body.left, body.top, body.width, body.height = Inches(0.78), Inches(1.35), Inches(3.65), Inches(3.55)
    set_box_text(
        body,
        [
            "One hot key dominates one reducer.",
            "That reducer spills or fails.",
            "The whole job waits for it.",
            "Retries do not fix structural skew.",
        ],
        [20.5, 20.5, 20.5, 20.5],
        [False] * 4,
        [PP_ALIGN.LEFT] * 4,
    )
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape.left, shape.top, shape.width, shape.height = Inches(4.35), Inches(1.17), Inches(6.65), Inches(3.95)


def tune_word_count_slides(prs):
    # Keep phase labels clear of the right decorative arc.
    for si in range(15, 20):
        slide = prs.slides[si - 1]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip() == "Reduce":
                shape.left = Inches(11.18)
                shape.width = Inches(1.0)
        if si == 19:
            for shape in slide.shapes:
                # Move final reduce output and its incoming arrows into the safe visual area.
                if shape.left >= Inches(10.20):
                    shape.left -= Inches(0.55)


def main():
    prs = Presentation(SRC)
    normalize_borders_and_text(prs)

    fix_formal_model(prs.slides[4])
    fix_formal_algorithm(prs.slides[5])
    fix_shuffle_cost(prs.slides[19])
    fix_join_slide(prs.slides[21])
    fix_skew_slide(prs.slides[24])
    tune_word_count_slides(prs)

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
